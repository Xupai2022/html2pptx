"""
样式映射器
将CSS样式映射为python-pptx样式
"""

import re
from typing import Optional, Dict
from pptx.util import Pt
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

from src.utils.color_parser import ColorParser
from src.utils.logger import setup_logger

logger = setup_logger(__name__)


class StyleMapper:
    """样式映射器"""

    # 默认样式
    DEFAULT_FONT_SIZE = 20
    DEFAULT_FONT_NAME = 'Microsoft YaHei'

    @staticmethod
    def parse_font_size(font_size_str: str) -> Optional[int]:
        """
        解析字体大小

        Args:
            font_size_str: 字体大小字符串(如'48px', '2em')

        Returns:
            点值(pt)
        """
        if not font_size_str:
            return None

        # 提取数字
        match = re.match(r'([\d.]+)(px|pt|em)?', font_size_str.strip())
        if not match:
            return None

        value = float(match.group(1))
        unit = match.group(2) or 'px'

        if unit == 'px':
            # 假设96 DPI, 1px ≈ 0.75pt
            return int(value * 0.75)
        elif unit == 'pt':
            return int(value)
        elif unit == 'em':
            # 1em = 16px (默认)
            return int(value * 16 * 0.75)

        return int(value)

    @staticmethod
    def parse_font_weight(font_weight_str: str) -> bool:
        """
        解析字体粗细

        Args:
            font_weight_str: 字体粗细(如'700', 'bold')

        Returns:
            是否加粗
        """
        if not font_weight_str:
            return False

        font_weight_str = font_weight_str.strip().lower()

        if font_weight_str in ('bold', 'bolder'):
            return True

        try:
            weight = int(font_weight_str)
            return weight >= 600
        except ValueError:
            return False

    @staticmethod
    def apply_text_style(text_frame, css_style: Dict[str, str], inline_style: Dict[str, str] = None):
        """
        应用文本样式

        Args:
            text_frame: python-pptx文本框对象
            css_style: CSS样式字典
            inline_style: 内联样式字典
        """
        # 合并样式(内联样式优先)
        merged_style = {**css_style}
        if inline_style:
            merged_style.update(inline_style)

        # 应用到段落
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                # 字体大小
                font_size_str = merged_style.get('font-size')
                if font_size_str:
                    font_size = StyleMapper.parse_font_size(font_size_str)
                    if font_size:
                        run.font.size = Pt(font_size)

                # 字体颜色
                color_str = merged_style.get('color')
                if color_str:
                    color = ColorParser.parse_color(color_str)
                    if color:
                        run.font.color.rgb = color

                # 字体粗细
                font_weight = merged_style.get('font-weight')
                if font_weight:
                    run.font.bold = StyleMapper.parse_font_weight(font_weight)

                # 字体名称
                font_family = merged_style.get('font-family')
                if font_family:
                    # 移除引号和多余空格
                    font_family = font_family.strip('\'"').split(',')[0].strip()
                    run.font.name = font_family
                else:
                    run.font.name = StyleMapper.DEFAULT_FONT_NAME

    @staticmethod
    def get_text_alignment(align_str: str):
        """
        获取文本对齐方式

        Args:
            align_str: 对齐字符串('left', 'center', 'right')

        Returns:
            PP_ALIGN枚举值
        """
        align_map = {
            'left': PP_ALIGN.LEFT,
            'center': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT,
            'justify': PP_ALIGN.JUSTIFY,
        }
        return align_map.get(align_str.lower(), PP_ALIGN.LEFT)

    @staticmethod
    def get_vertical_alignment(align_str: str):
        """
        获取垂直对齐方式

        Args:
            align_str: 对齐字符串('top', 'middle', 'bottom')

        Returns:
            MSO_ANCHOR枚举值
        """
        align_map = {
            'top': MSO_ANCHOR.TOP,
            'middle': MSO_ANCHOR.MIDDLE,
            'bottom': MSO_ANCHOR.BOTTOM,
        }
        return align_map.get(align_str.lower(), MSO_ANCHOR.TOP)

    @staticmethod
    def parse_padding(padding_str: str) -> tuple:
        """
        解析padding值

        Args:
            padding_str: padding字符串(如'20px', '10px 20px')

        Returns:
            (top, right, bottom, left) 像素值
        """
        if not padding_str:
            return (0, 0, 0, 0)

        parts = padding_str.split()
        values = [int(re.search(r'\d+', p).group()) for p in parts if re.search(r'\d+', p)]

        if len(values) == 1:
            return (values[0], values[0], values[0], values[0])
        elif len(values) == 2:
            return (values[0], values[1], values[0], values[1])
        elif len(values) == 3:
            return (values[0], values[1], values[2], values[1])
        elif len(values) == 4:
            return tuple(values)

        return (0, 0, 0, 0)

    @staticmethod
    def parse_border_width(border_str: str) -> int:
        """
        解析边框宽度

        Args:
            border_str: 边框字符串(如'4px solid red')

        Returns:
            宽度像素值
        """
        if not border_str:
            return 0

        match = re.search(r'(\d+)px', border_str)
        return int(match.group(1)) if match else 0
