"""
HTML转PPTX主程序
"""

import sys
import re
from pathlib import Path

from src.parser.html_parser import HTMLParser
from src.parser.css_parser import CSSParser
from src.renderer.pptx_builder import PPTXBuilder
from src.converters.text_converter import TextConverter
from src.converters.table_converter import TableConverter
from src.converters.shape_converter import ShapeConverter
from src.converters.chart_converter import ChartConverter
from src.converters.timeline_converter import TimelineConverter
from src.converters.svg_converter import SvgConverter
from src.utils.logger import setup_logger
from src.utils.unit_converter import UnitConverter
from src.utils.color_parser import ColorParser
from src.utils.chart_capture import ChartCapture
from src.utils.font_manager import get_font_manager
from src.utils.style_computer import get_style_computer
from src.utils.content_height_calculator import ContentHeightCalculator
from pptx.util import Pt
from pptx.enum.text import MSO_ANCHOR, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor

logger = setup_logger(__name__)


class HTML2PPTX:
    """HTML转PPTX转换器"""

    def __init__(self, html_path: str, existing_presentation=None, template_file: str = None, use_shared_presentation=None, use_stable_chart_capture=False):
        """
        初始化转换器

        Args:
            html_path: HTML文件路径
            existing_presentation: 可选的现有Presentation对象，用于批量转换
            template_file: 模板文件路径（兼容旧版本）
            use_shared_presentation: 共享的Presentation对象（用于批量转换）
            use_stable_chart_capture: 是否使用稳定的截图版本
        """
        self.html_path = html_path
        self.html_parser = HTMLParser(html_path)
        # 使用完整的HTML soup来初始化CSS解析器，以便解析head中的style标签
        self.css_parser = CSSParser(self.html_parser.full_soup)

        # 优先使用use_shared_presentation，其次是existing_presentation
        if use_shared_presentation:
            self.pptx_builder = PPTXBuilder(use_shared_presentation)
        elif existing_presentation:
            self.pptx_builder = PPTXBuilder(existing_presentation)
        else:
            self.pptx_builder = PPTXBuilder()

        # 保存稳定版本标志
        self.use_stable_chart_capture = use_stable_chart_capture

        # 初始化全局字体管理器和样式计算器
        self.font_manager = get_font_manager(self.css_parser)
        self.style_computer = get_style_computer(self.css_parser)
        
        # 初始化内容高度计算器
        self.height_calculator = ContentHeightCalculator(self.css_parser, self.style_computer)

        # 设置HTML文件ID，避免缓存冲突
        if hasattr(self.style_computer, 'set_html_file_id'):
            self.style_computer.set_html_file_id(html_path)
        if hasattr(self.style_computer.font_size_extractor, 'set_html_file_id'):
            self.style_computer.font_size_extractor.set_html_file_id(html_path)

        # 记录所有SVG转换器实例，用于清理临时文件
        self.svg_converters = []

    def convert(self, output_path: str):
        """
        执行转换

        Args:
            output_path: 输出PPTX路径
        """
        logger.info("=" * 50)
        logger.info("开始HTML转PPTX转换")
        logger.info("=" * 50)

        # 获取所有幻灯片
        slides = self.html_parser.get_slides()

        for slide_html in slides:
            logger.info(f"\n处理幻灯片...")

            # 创建空白幻灯片
            pptx_slide = self.pptx_builder.add_blank_slide()

            # 初始化转换器
            text_converter = TextConverter(pptx_slide, self.css_parser)
            table_converter = TableConverter(pptx_slide, self.css_parser)
            shape_converter = ShapeConverter(pptx_slide, self.css_parser)

            # 1. 添加顶部装饰条
            shape_converter.add_top_bar()

            # 2. 添加标题和副标题
            title_info = self.html_parser.get_title_info(slide_html)
            if title_info:
                # content-section的padding-top是20px
                title_end_y = text_converter.convert_title(
                    title_info['text'],
                    title_info['subtitle'],
                    x=80,
                    y=20,
                    is_cover=title_info.get('is_cover', False),
                    title_classes=title_info.get('classes', []),
                    h1_element=title_info.get('h1_element')
                )
                # space-y-10的第一个子元素紧接标题区域（无上间距）
                y_offset = title_end_y
            else:
                # 没有标题时使用默认位置（content-section padding-top）
                y_offset = 20

            # 3. 处理内容区域

            # 统一处理所有容器：优先查找space-y-10容器，如果没有则处理content-section的直接子元素
            space_y_container = slide_html.find('div', class_='space-y-10')
            content_section = slide_html.find('div', class_='content-section')

            if space_y_container:
                logger.info("找到space-y-10容器，开始处理直接子元素")
                # 按顺序遍历直接子元素
                is_first_container = True
                container_count = 0
                for container in space_y_container.find_all(recursive=False):
                    if not container.name or container.name in ['nav', 'script', 'style']:
                        continue

                    container_count += 1
                    container_classes = container.get('class', [])
                    logger.info(f"处理容器 #{container_count}: tag={container.name}, class={container_classes}")

                    # space-y-10: 第一个元素无上间距，后续元素有40px间距
                    if not is_first_container:
                        y_offset += 40  # space-y-10间距
                        logger.info(f"添加space-y-10间距40px，当前y_offset={y_offset}")
                    is_first_container = False

                    # 根据class路由到对应的处理方法
                    try:
                        old_y = y_offset
                        y_offset = self._process_container(container, pptx_slide, y_offset, shape_converter)
                        logger.info(f"容器处理完成，y_offset从{old_y}变为{y_offset}")
                        if y_offset == old_y:
                            logger.warning(f"警告：容器{container_classes}的y_offset没有变化，可能内容未正确处理")
                    except Exception as e:
                        logger.error(f"处理容器时出错: {e}, container={container_classes}")
                        import traceback
                        logger.error(f"错误堆栈: {traceback.format_exc()}")
                        # 继续处理下一个容器
                        continue
            else:
                # 新的HTML结构：直接处理content-section的直接子元素（跳过标题区域）
                logger.info("未找到space-y-10容器，处理content-section的直接子元素")

                # 跳过标题区域（第一个包含h1/h2的mb-6容器）
                containers = []
                skip_first_mb = True  # 默认跳过第一个mb容器
                for child in content_section.children:
                    if hasattr(child, 'get') and child.get('class'):
                        classes = child.get('class', [])

                        # 检查是否是标题区域
                        has_title = child.find('h1') or child.find('h2')

                        # 如果是第一个mb容器且有标题，则跳过
                        # 但要确保它不包含其他重要内容（如grid、stat-card等）
                        is_title_container = False
                        if skip_first_mb and any(cls in ['mb-6', 'mb-4', 'mb-8'] for cls in classes) and has_title:
                            # 检查是否真的是纯标题容器（不包含grid、card等内容）
                            has_content = any(cls in classes for cls in ['grid', 'stat-card', 'data-card', 'risk-card', 'flex'])
                            # 检查是否只包含标题元素和简单装饰
                            has_only_titles = True
                            for elem in child.find_all(['div', 'span', 'i']):
                                elem_classes = elem.get('class', [])
                                # 如果找到非装饰性的类，说明不是纯标题容器
                                if any(cls in elem_classes for cls in ['w-', 'h-', 'primary-bg']) and not any(cls in elem_classes for cls in ['fas', 'fa-']):
                                    # 这是装饰条，允许存在
                                    continue
                                elif any(cls in elem_classes for cls in ['fas', 'fa-']):
                                    # 图标也允许
                                    continue
                                elif elem.get('style') and any(prop in elem.get('style', '') for prop in ['width', 'height', 'background']):
                                    # 内联样式的装饰元素
                                    continue

                            if not has_content:
                                is_title_container = True
                                skip_first_mb = False  # 跳过后设置为false

                        if is_title_container:
                            continue  # 跳过纯标题容器

                        # 其他容器都保留
                        if child.name:
                            containers.append(child)

                # 处理所有容器
                for container in containers:
                    if container.name:
                        # 添加间距（模拟mb-6等间距）
                        if containers.index(container) > 0:
                            y_offset += 40  # 间距
                        y_offset = self._process_container(container, pptx_slide, y_offset, shape_converter)

            # 4. 添加页码
            page_num = self.html_parser.get_page_number(slide_html)
            if page_num:
                shape_converter.add_page_number(page_num)

        # 保存PPTX
        self.pptx_builder.save(output_path)

        # 清理所有临时PNG文件
        self._cleanup_temp_files()

        logger.info("=" * 50)
        logger.info(f"转换完成! 输出: {output_path}")
        logger.info("=" * 50)

    def _cleanup_temp_files(self):
        """
        清理所有临时文件
        """
        # 清理所有SVG转换器生成的临时PNG文件
        for svg_converter in self.svg_converters:
            svg_converter.cleanup_temp_files()
        self.svg_converters.clear()

        # 清理当前目录下可能残留的临时PNG文件
        import os
        import glob
        pattern = "svg_screenshot_*.png"
        for png_file in glob.glob(pattern):
            try:
                os.remove(png_file)
                logger.info(f"已删除残留临时文件: {png_file}")
            except Exception as e:
                logger.warning(f"删除残留临时文件失败 {png_file}: {e}")

    def _process_container(self, container, pptx_slide, y_offset, shape_converter):
        """
        处理单个容器，根据其类型路由到相应的处理方法

        Args:
            container: 容器元素
            pptx_slide: PPTX幻灯片
            y_offset: 当前Y坐标
            shape_converter: 形状转换器

        Returns:
            下一个元素的Y坐标
        """
        container_classes = container.get('class', [])

        # 检测封面页容器（优先级最高）
        if 'cover-content' in container_classes or 'cover-info' in container_classes:
            logger.info(f"识别为封面页容器: {container_classes}，不添加背景")
            # 封面页容器不添加背景，直接处理内容
            return self._convert_cover_container(container, pptx_slide, y_offset)

        # 根据class路由到对应的处理方法
        # 优先检测grid布局（包含grid类）
        if 'grid' in container_classes:
            # 网格容器（新的Tailwind结构）
            logger.info(f"识别为grid容器: {container_classes}")
            return self._convert_grid_container(container, pptx_slide, y_offset, shape_converter)
        elif 'stats-container' in container_classes:
            # 顶层stats-container（不在stat-card内）
            logger.info(f"识别为stats-container: {container_classes}")
            return self._convert_stats_container(container, pptx_slide, y_offset)
        elif 'stat-card' in container_classes:
            logger.info(f"识别为stat-card: {container_classes}")
            return self._convert_stat_card(container, pptx_slide, y_offset)
        elif 'data-card' in container_classes:
            logger.info(f"识别为data-card: {container_classes}")
            return self._convert_data_card(container, pptx_slide, shape_converter, y_offset)
        elif 'strategy-card' in container_classes:
            logger.info(f"识别为strategy-card: {container_classes}")
            return self._convert_strategy_card(container, pptx_slide, y_offset)
        elif 'risk-card' in container_classes:
            return self._convert_risk_card(container, pptx_slide, shape_converter, y_offset)
        elif 'flex' in container_classes and 'gap-6' in container_classes:
            # 检查是否包含SVG图表的flex容器
            svgs_in_container = container.find_all('svg')
            if svgs_in_container:
                logger.info(f"检测到包含 {len(svgs_in_container)} 个SVG的flex容器")
                return self._convert_flex_charts_container(container, pptx_slide, y_offset, shape_converter)
            else:
                # 底部信息容器（包含bullet-point的flex布局）
                return self._convert_bottom_info(container, pptx_slide, y_offset)
        elif 'flex' in container_classes and 'justify-between' in container_classes:
            # 底部信息容器（包含bullet-point的flex布局）
            return self._convert_bottom_info(container, pptx_slide, y_offset)
        elif 'flex-1' in container_classes and 'overflow-hidden' in container_classes:
            # 先检查是否是居中容器
            has_justify_center = 'justify-center' in container_classes
            has_flex_col = 'flex-col' in container_classes
            has_items_center = 'items-center' in container_classes

            # 如果同时有居中相关的类，优先作为居中容器处理
            if has_justify_center and (has_flex_col or has_items_center):
                logger.info(f"检测到居中容器（flex-1 overflow-hidden variant）: {container_classes}")
                return self._convert_centered_container(container, pptx_slide, y_offset, shape_converter)
            else:
                # 内容容器（包含多个子容器）
                return self._convert_content_container(container, pptx_slide, y_offset, shape_converter)
        else:
            # 首先检查是否包含SVG元素
            svgs_in_container = container.find_all('svg')
            if svgs_in_container:
                logger.info(f"检测到容器包含 {len(svgs_in_container)} 个SVG元素")
                # 初始化SVG转换器
                svg_converter = SvgConverter(pptx_slide, self.css_parser, self.html_path, self.use_stable_chart_capture)
                self.svg_converters.append(svg_converter)  # 记录实例

                # 如果是单个SVG，直接转换
                if len(svgs_in_container) == 1:
                    svg_elem = svgs_in_container[0]

                    # 检查是否有标题
                    h3_elem = container.find('h3')
                    if h3_elem:
                        # 处理标题
                        title_text = h3_elem.get_text(strip=True)
                        if title_text:
                            text_box = pptx_slide.shapes.add_textbox(
                                UnitConverter.px_to_emu(80),
                                UnitConverter.px_to_emu(y_offset),
                                UnitConverter.px_to_emu(1760),
                                UnitConverter.px_to_emu(30)
                            )
                            text_frame = text_box.text_frame
                            text_frame.text = title_text
                            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

                            for paragraph in text_frame.paragraphs:
                                paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                                for run in paragraph.runs:
                                    run.font.size = Pt(20)
                                    run.font.bold = True
                                    run.font.name = self.font_manager.get_font('body')
                                    run.font.color.rgb = ColorParser.parse_color('rgb(10, 66, 117)')

                            y_offset += 40

                    # 转换SVG
                    # 使用SVG的原始尺寸，不进行缩放
                    svg_width, svg_height = svg_converter._get_svg_dimensions(svg_elem)
                    logger.info(f"SVG原始尺寸: {svg_width}x{svg_height}px")

                    # 检查父容器是否有flex居中布局
                    parent = svg_elem.parent
                    is_centered = False
                    if parent and 'class' in parent.attrs:
                        classes = parent.get('class', [])
                        if any('justify-center' in str(c) for c in classes):
                            is_centered = True
                            logger.info(f"检测到SVG居中布局: {classes}")

                    # 使用SVG原始尺寸
                    chart_width = svg_width

                    # 如果是居中布局，计算居中位置
                    left = 80  # 默认左边距
                    if is_centered:
                        # 计算居中位置：(幻灯片宽度 - SVG宽度) / 2
                        left = (1920 - chart_width) / 2
                        logger.info(f"SVG居中显示，左边距: {left}px")

                    chart_height = svg_converter.convert_svg(
                        svg_elem,
                        container,
                        left,
                        y_offset,
                        chart_width,
                        0
                    )

                    # 更新y_offset，继续处理容器中的其他元素
                    y_offset += chart_height + 20

                    # 移除已处理的SVG元素，继续处理其他元素
                    svg_elem.decompose()

                    # 移除已处理的标题（如果有）
                    if h3_elem:
                        h3_elem.decompose()

                    # 继续处理容器中的其他元素
                    return self._convert_content_container(container, pptx_slide, y_offset, shape_converter)
                else:
                    # 多个SVG，使用水平布局
                    chart_height = svg_converter.convert_multiple_svgs(
                        container,
                        80,
                        y_offset,
                        1760,
                        gap=24
                    )

                    # 更新y_offset，继续处理容器中的其他元素
                    y_offset += chart_height + 20

                    # 继续处理容器中的其他元素
                    return self._convert_content_container(container, pptx_slide, y_offset, shape_converter)

            # 先检查是否包含已知子容器（优先级高于数字列表检测）
            # 这样可以避免将包含数字的stat-card误判为数字列表
            has_data_cards = container.find_all('div', class_='data-card')
            has_stat_cards = container.find_all('div', class_='stat-card')
            has_grid = container.find('div', class_='grid')
            
            if has_data_cards or has_stat_cards or has_grid:
                logger.info(f"容器 {container_classes} 包含已知子容器，递归处理")
                return self._convert_content_container(container, pptx_slide, y_offset, shape_converter)
            
            # 检测是否包含多个数字列表项（如多个toc-item）
            toc_items = container.find_all('div', class_='toc-item')
            if len(toc_items) > 1:
                return self._convert_numbered_list_group(container, pptx_slide, y_offset)
            elif 'toc-item' in container_classes or self._has_numbered_list_pattern(container):
                # 单个数字列表项
                return self._convert_numbered_list_container(container, pptx_slide, y_offset)

            # 检测flex容器（放在最后，避免误判）
            if 'flex-1' in container_classes or 'flex' in container_classes:
                # 检查flex容器内是否包含网格布局
                grid_child = container.find('div', class_='grid')
                if grid_child:
                    # 如果flex容器内只有一个grid子容器，直接处理grid
                    direct_children = [child for child in container.children if hasattr(child, 'name') and child.name]
                    if len(direct_children) == 1 and direct_children[0] == grid_child:
                        logger.info("flex容器内只包含一个网格容器，直接处理网格布局")
                        return self._convert_grid_container(grid_child, pptx_slide, y_offset, shape_converter)

                # flex容器 - 增强检测，处理居中布局（优先检测）
                # 检查是否是居中容器
                has_justify_center = 'justify-center' in container_classes
                has_items_center = 'items-center' in container_classes
                has_flex_col = 'flex-col' in container_classes
                has_overflow_hidden = 'overflow-hidden' in container_classes
                has_flex_1 = 'flex-1' in container_classes or 'flex' in container_classes

                # 如果是居中布局的flex容器（增强检测逻辑）
                # 条件1：有justify-center和items-center（水平垂直居中）
                # 条件2：有justify-center和flex-col（垂直居中）
                # 条件3：有justify-center且是flex容器（更宽松的检测）
                # 条件4：有flex-col和justify-center（特别处理垂直居中）
                if (has_justify_center and has_items_center) or \
                   (has_flex_col and has_justify_center) or \
                   (has_justify_center and has_flex_1):
                    logger.info(f"检测到居中容器: {container_classes}")
                    return self._convert_centered_container(container, pptx_slide, y_offset, shape_converter)

                # 普通flex容器
                return self._convert_flex_container(container, pptx_slide, y_offset, shape_converter)

            # 未知容器类型，先检查是否包含已知子容器
            # 检查是否包含data-card、stat-card、grid等已知容器
            has_data_cards = container.find_all('div', class_='data-card')
            has_stat_cards = container.find_all('div', class_='stat-card')
            has_grid = container.find('div', class_='grid')
            
            if has_data_cards or has_stat_cards or has_grid:
                logger.info(f"未知容器类型 {container_classes} 包含已知子容器，递归处理")
                return self._convert_content_container(container, pptx_slide, y_offset, shape_converter)
            
            # 真正的未知容器类型，记录警告但尝试处理
            logger.warning(f"遇到未知容器类型: {container_classes}，尝试降级处理")
            return self._convert_generic_card(container, pptx_slide, y_offset, card_type='unknown')

    def _convert_grid_container(self, container, pptx_slide, y_start, shape_converter):
        """
        转换网格容器（如grid grid-cols-2 gap-6）

        Args:
            container: 网格容器元素
            pptx_slide: PPTX幻灯片
            y_start: 起始Y坐标
            shape_converter: 形状转换器

        Returns:
            下一个元素的Y坐标
        """
        logger.info("处理网格容器")
        classes = container.get('class', [])

        # 获取网格列数
        num_columns = 2  # 默认2列
        for cls in classes:
            if cls.startswith('grid-cols-') and hasattr(self.css_parser, 'tailwind_grid_columns'):
                columns = self.css_parser.tailwind_grid_columns.get(cls)
                if columns:
                    num_columns = columns
                    logger.info(f"检测到网格列数: {num_columns}")
                    break

        # 获取间距 - 优先从CSS，然后检查inline style，最后使用默认值
        gap = self.css_parser.get_gap_size('.grid')  # 从CSS获取间距
        
        # 检查是否有gap-*类覆盖
        for cls in classes:
            if cls.startswith('gap-'):
                gap_from_class = self.css_parser.get_gap_size(f'.{cls}')
                if gap_from_class != 20:  # 如果不是默认值，使用它
                    gap = gap_from_class
                    logger.info(f"从Tailwind类 {cls} 检测到网格间距: {gap}px")
                    break
        
        # 检查inline style中的gap
        inline_style = container.get('style', '')
        if 'gap:' in inline_style:
            import re
            gap_match = re.search(r'gap:\s*(\d+)px', inline_style)
            if gap_match:
                gap = int(gap_match.group(1))
                logger.info(f"从inline style检测到网格间距: {gap}px")
        
        if gap == 20:
            logger.debug(f"使用默认网格间距: {gap}px")

        # 获取所有子元素
        children = []
        for child in container.children:
            if hasattr(child, 'name') and child.name:
                children.append(child)

        # 计算布局
        total_width = 1760  # 可用宽度
        item_width = (total_width - (num_columns - 1) * gap) // num_columns

        current_y = y_start
        max_y_in_row = y_start

        for idx, child in enumerate(children):
            col = idx % num_columns
            row = idx // num_columns

            if col == 0 and idx > 0:
                # 新行开始
                current_y = max_y_in_row + gap

            x = 80 + col * (item_width + gap)
            y = current_y

            # 处理子元素
            child_classes = child.get('class', [])

            # 检查是否包含chart-container（用于SVG图表）
            has_chart_container = child.find('div', class_='chart-container') is not None
            has_svg = child.find('svg') is not None

            # 检查是否需要添加左边框（data-card有左边框特性）
            needs_left_border = False
            if 'data-card' in child_classes or 'stat-card' in child_classes:
                # 检查CSS中是否有border-left样式
                css_style = self.css_parser.get_style(f".{child_classes[0]}") if child_classes else {}
                if 'border-left' in css_style or 'data-card' in child_classes:
                    needs_left_border = True

            if 'stat-card' in child_classes:
                # 专门处理网格中的stat-card
                child_y = self._convert_grid_stat_card(child, pptx_slide, shape_converter, x, y, item_width)
            elif 'data-card' in child_classes:
                # data-card需要左边框，但我们不能直接调用_convert_data_card
                # 因为它会自己添加边框。我们需要特殊处理
                child_y = self._convert_grid_data_card(child, pptx_slide, shape_converter, x, y, item_width)
            elif 'risk-card' in child_classes:
                # 专门处理网格中的risk-card
                child_y = self._convert_grid_risk_card(child, pptx_slide, shape_converter, x, y, item_width)
            elif has_chart_container or has_svg:
                # 处理包含SVG图表的子元素（如slide_003.html）
                logger.info(f"网格子元素包含SVG图表，使用SVG处理逻辑")
                child_y = self._convert_grid_svg_chart(child, pptx_slide, shape_converter, x, y, item_width)
            else:
                # 降级处理
                child_y = self._convert_generic_card(child, pptx_slide, y, card_type='grid-item')

                # 如果需要左边框，在这里添加
                if needs_left_border:
                    # 计算实际高度
                    actual_height = child_y - y
                    if actual_height > 0:
                        shape_converter.add_border_left(x, y, actual_height, 4)

            max_y_in_row = max(max_y_in_row, child_y)

        return max_y_in_row + 20  # 返回下一行的起始位置

    def _convert_grid_data_card(self, card, pptx_slide, shape_converter, x, y, width):
        """
        转换网格中的data-card（带左侧竖线）

        Args:
            card: data-card元素
            pptx_slide: PPTX幻灯片
            shape_converter: 形状转换器
            x: X坐标
            y: Y坐标
            width: 宽度

        Returns:
            下一个元素的Y坐标
        """
        logger.info("处理网格中的data-card")

        # 初始化变量
        risk_color = None
        bg_color = None

        # 添加data-card背景色
        bg_color_str = 'rgba(10, 66, 117, 0.03)'  # 从CSS获取的背景色
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.dml.color import RGBColor

        # 使用ContentHeightCalculator动态计算data-card高度
        estimated_height = self.height_calculator.calculate_data_card_height(card, width)

        bg_shape = pptx_slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            UnitConverter.px_to_emu(x),
            UnitConverter.px_to_emu(y),
            UnitConverter.px_to_emu(width),
            UnitConverter.px_to_emu(estimated_height)
        )
        bg_shape.fill.solid()
        bg_rgb, alpha = ColorParser.parse_rgba(bg_color_str)
        if bg_rgb:
            if alpha < 1.0:
                bg_rgb = ColorParser.blend_with_white(bg_rgb, alpha)
            bg_shape.fill.fore_color.rgb = bg_rgb
        bg_shape.line.fill.background()
        logger.info(f"添加data-card背景色，高度={estimated_height}px")

        current_y = y + 15  # 顶部padding，与CSS中的15px保持一致

        # 1. 首先处理h3标题
        h3_elem = card.find('h3')
        if h3_elem:
            h3_text = h3_elem.get_text(strip=True)
            if h3_text:
                text_left = UnitConverter.px_to_emu(x + 20)
                text_top = UnitConverter.px_to_emu(current_y)
                text_box = pptx_slide.shapes.add_textbox(
                    text_left, text_top,
                    UnitConverter.px_to_emu(width - 40), UnitConverter.px_to_emu(30)
                )
                text_frame = text_box.text_frame
                text_frame.text = h3_text
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        font_size_pt = self.style_computer.get_font_size_pt(h3_elem)
                        run.font.size = Pt(font_size_pt)
                        run.font.color.rgb = ColorParser.get_primary_color()
                        run.font.bold = True
                        run.font.name = self.font_manager.get_font('h3')

                current_y += 40  # 28px字体 + 12px margin-bottom
                logger.info(f"渲染h3标题: {h3_text}")

        # 2. 处理risk-item或bullet-point
        risk_items = card.find_all('div', class_='risk-item')
        bullet_points = card.find_all('div', class_='bullet-point')
        # 同时检查space-y-3容器内的flex items-start结构（slide_011.html的实际结构）
        if not bullet_points:
            space_y_containers = card.find_all('div', class_='space-y-3')
            for container in space_y_containers:
                flex_items = container.find_all('div', class_='flex')
                for flex_item in flex_items:
                    if flex_item.find('i') and flex_item.find('p'):
                        bullet_points.append(flex_item)

        if risk_items:
            logger.info(f"找到 {len(risk_items)} 个risk-item")
            self._process_risk_items(risk_items, card, pptx_slide, x, y, width, current_y)
        elif bullet_points:
            logger.info(f"找到 {len(bullet_points)} 个bullet-point")
            self._process_bullet_points(bullet_points, card, pptx_slide, x, y, width, current_y)

        # 3. 如果没有risk-item和bullet-point，使用原来的逻辑处理其他内容
        if not risk_items and not bullet_points:
            # 提取文本内容
            text_elements = []
            for elem in card.descendants:
                if hasattr(elem, 'name') and elem.name in ['p', 'h1', 'h2', 'h3', 'h4', 'div', 'span']:
                    # 只提取没有子块级元素的文本节点
                    if not elem.find_all(['div', 'p', 'h1', 'h2', 'h3']):
                        text = elem.get_text(strip=True)
                        if text and len(text) > 2:
                            text_elements.append(elem)

            # 渲染文本
            for elem in text_elements[:5]:  # 最多5个元素
                text = elem.get_text(strip=True)
                if text:
                    # 获取元素的class以判断样式
                    elem_classes = elem.get('class', [])
                    
                    text_left = UnitConverter.px_to_emu(x + 20)
                    text_top = UnitConverter.px_to_emu(current_y)
                    text_box = pptx_slide.shapes.add_textbox(
                        text_left, text_top,
                        UnitConverter.px_to_emu(width - 40), UnitConverter.px_to_emu(30)
                    )
                    text_frame = text_box.text_frame
                    text_frame.text = text
                    text_frame.word_wrap = True

                    # 判断是否需要加粗（stat-value通常加粗）
                    is_bold = 'stat-value' in elem_classes or 'font-bold' in elem_classes or 'font-semibold' in elem_classes
                    # 判断是否为主色（stat-value和stat-label通常有颜色）
                    is_primary = 'stat-value' in elem_classes or 'primary-color' in elem_classes
                    is_gray = 'stat-label' in elem_classes or 'text-gray-600' in elem_classes

                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            # 使用样式计算器获取字体大小
                            font_size_px = self.style_computer.get_font_size_pt(elem)
                            run.font.size = Pt(font_size_px) if font_size_px else Pt(16)
                            run.font.name = self.font_manager.get_font('body')
                            
                            # 应用加粗
                            if is_bold:
                                run.font.bold = True
                            
                            # 应用颜色
                            if is_primary:
                                run.font.color.rgb = ColorParser.get_primary_color()
                            elif is_gray:
                                run.font.color.rgb = RGBColor(102, 102, 102)  # #666

                    # 根据元素类型计算合适的间距
                    if 'stat-value' in elem_classes:
                        # stat-value后的间距较小（CSS: margin-bottom: 8px）
                        # 考虑字体大小42px * 1.2行高 + 8px margin ≈ 58px
                        current_y += 58
                    elif 'stat-label' in elem_classes:
                        # stat-label后的间距（20px字体 + 正常间距）
                        # 考虑字体大小20px * 1.2行高 + 12px margin ≈ 36px
                        current_y += 36
                    elif 'mt-3' in elem_classes:
                        # p标签带mt-3，需要额外的顶部间距
                        # mt-3 = 12px，加上段落本身高度
                        current_y += 35
                    else:
                        # 默认间距
                        current_y += 35

        # 添加左边框 - 使用与背景相同的高度，确保竖线不会过长
        border_height = estimated_height  # 使用背景矩形的高度
        shape_converter.add_border_left(x, y, border_height, 4)

        # 确保返回正确的位置，保留底部padding
        # 返回背景矩形的底部位置 + 额外间距
        return y + estimated_height + 5

    def _calculate_text_width(self, text: str, font_size: Pt) -> int:
        """
        计算文本的像素宽度

        Args:
            text: 要计算的文本
            font_size: 字体大小（Pt）

        Returns:
            int: 文本的宽度（像素）
        """
        # 将Pt转换为Px
        font_size_px = int(font_size.pt * 0.75)

        # 计算字符宽度
        total_width = 0
        for char in text:
            # 中文字符宽度约为字体大小的1倍
            if '\u4e00' <= char <= '\u9fff':
                char_width = font_size_px
            # 英文字母和数字宽度约为字体大小的0.6倍
            elif char.isalnum() or char in '.,;:!?\'"()[]{}-+/\\=_@#%&*':
                char_width = int(font_size_px * 0.6)
            # 空格宽度约为字体大小的0.3倍
            elif char == ' ':
                char_width = int(font_size_px * 0.3)
            # 其他符号
            else:
                char_width = font_size_px

            total_width += char_width

        return total_width

    def _process_bullet_points(self, bullet_points, card, pptx_slide, x, y, width, current_y):
        """
        处理bullet-point列表

        Args:
            bullet_points: bullet-point元素列表
            card: 父data-card元素
            pptx_slide: PPTX幻灯片
            x: 起始X坐标
            y: 起始Y坐标（相对于幻灯片的绝对位置）
            width: 容器宽度
            current_y: 当前Y坐标偏移（相对于容器的位置）
        """
        from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
        from pptx.dml.color import RGBColor

        # 使用current_y而不是y作为起始位置，因为current_y已经考虑了标题的偏移
        actual_y = current_y if current_y > y else y

        for bp in bullet_points:
            # 获取图标
            icon_elem = bp.find('i')
            icon_char = None
            icon_color = ColorParser.get_primary_color()

            if icon_elem:
                icon_classes = icon_elem.get('class', [])
                # 使用_get_icon_char函数获取图标字符
                icon_char = self._get_icon_char(icon_classes)

                # 根据图标类确定颜色
                if 'text-red-600' in icon_classes or 'risk-high' in icon_classes:
                    icon_color = RGBColor(220, 38, 38)  # 红色
                elif 'text-orange-600' in icon_classes or 'risk-medium' in icon_classes:
                    icon_color = RGBColor(234, 88, 12)  # 橙色
                elif 'text-yellow-600' in icon_classes or 'text-green-600' in icon_classes or 'risk-low' in icon_classes:
                    icon_color = RGBColor(34, 197, 94)  # 绿色
                elif 'text-blue-600' in icon_classes:
                    icon_color = RGBColor(59, 130, 246)  # 蓝色
                elif 'primary-color' in icon_classes:
                    icon_color = ColorParser.get_primary_color()

            # 获取段落元素
            p_elem = bp.find('p')
            if p_elem:
                # 检查是否包含priority-tag
                priority_tag = p_elem.find('span', class_='priority-tag')
                main_text = ""
                tag_text = ""
                tag_color = None

                if priority_tag:
                    # 提取标签文本
                    tag_text = priority_tag.get_text(strip=True)
                    tag_classes = priority_tag.get('class', [])

                    # 确定标签颜色
                    if 'priority-high' in tag_classes:
                        tag_color = RGBColor(239, 68, 68)  # 红色
                        tag_bg_color = RGBColor(254, 226, 226)  # 浅红色背景
                    elif 'priority-medium' in tag_classes:
                        tag_color = RGBColor(251, 146, 60)  # 橙色
                        tag_bg_color = RGBColor(255, 237, 213)  # 浅橙色背景
                    elif 'priority-low' in tag_classes:
                        tag_color = RGBColor(209, 177, 0)  # 黄色
                        tag_bg_color = RGBColor(250, 204, 21)  # 浅黄色背景

                    # 移除标签后获取主文本
                    priority_tag.extract()
                    main_text = p_elem.get_text(strip=True)
                else:
                    main_text = p_elem.get_text(strip=True)

                # 创建文本框，使用actual_y确保不会与标题重叠
                text_box = pptx_slide.shapes.add_textbox(
                    UnitConverter.px_to_emu(x + 20),
                    UnitConverter.px_to_emu(actual_y),
                    UnitConverter.px_to_emu(width - 40),
                    UnitConverter.px_to_emu(35)
                )
                text_frame = text_box.text_frame
                text_frame.clear()

                # 添加段落
                p = text_frame.paragraphs[0]

                # 添加图标
                if icon_char:
                    icon_run = p.add_run()
                    icon_run.text = icon_char + " "
                    # 图标字体大小与文本同步
                    icon_font_size_pt = self._get_font_size_pt(p_elem, default_px=25)
                    icon_run.font.size = Pt(icon_font_size_pt)
                    icon_run.font.color.rgb = icon_color
                    icon_run.font.name = self.font_manager.get_font('body')

                # 添加主文本
                if main_text:
                    text_run = p.add_run()
                    text_run.text = main_text

                    # 获取字体大小（从CSS解析）
                    font_size_pt = self._get_font_size_pt(p_elem, default_px=25)
                    text_run.font.size = Pt(font_size_pt)

                    text_run.font.name = self.font_manager.get_font('body')
                    text_run.font.color.rgb = RGBColor(51, 51, 51)  # 深灰色

                # 添加标签文本（如果存在）
                if tag_text and tag_color:
                    tag_run = p.add_run()
                    tag_run.text = " " + tag_text
                    # 标签字体稍小
                    tag_font_size_pt = max(12, self._get_font_size_pt(p_elem, default_px=25) - 4)
                    tag_run.font.size = Pt(tag_font_size_pt)
                    tag_run.font.color.rgb = tag_color
                    tag_run.font.bold = True
                    tag_run.font.name = self.font_manager.get_font('body')

                # 设置段落格式
                p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                p.space_before = Pt(0)
                p.space_after = Pt(4)

                actual_y += 35  # 每个bullet-point占35px

        logger.info(f"处理了 {len(bullet_points)} 个bullet-point")

    def _process_risk_items(self, risk_items, card, pptx_slide, x, y, width, current_y):
        """
        处理risk-item列表

        Args:
            risk_items: risk-item元素列表
            card: 父data-card元素
            pptx_slide: PPTX幻灯片
            x: 起始X坐标
            y: 起始Y坐标
            width: 容器宽度
            current_y: 当前Y坐标偏移
        """
        from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
        from pptx.dml.color import RGBColor
        from src.utils.unit_converter import UnitConverter
        from pptx.util import Pt

        logger.info(f"开始处理{len(risk_items)}个risk-item")

        # 使用current_y作为起始位置，与bullet-point保持一致
        actual_y = current_y if current_y > y else y

        for risk_item in risk_items:
            # 获取图标
            icon_elem = risk_item.find('i')
            icon_char = None
            icon_color = RGBColor(220, 38, 38)  # 默认红色（risk-icon）

            if icon_elem:
                icon_classes = icon_elem.get('class', [])
                icon_char = self._get_icon_char(icon_classes)

                # 根据图标类调整颜色
                if 'bullet-icon' in icon_classes:
                    icon_color = ColorParser.get_primary_color()

            # 处理文本内容
            text_container = risk_item.find('div')
            if text_container:
                # 处理主标题和风险等级
                first_p = text_container.find('p')
                if first_p:
                    # 提取strong文本
                    strong_elem = first_p.find('strong')
                    main_text = strong_elem.get_text(strip=True) if strong_elem else first_p.get_text(strip=True)

                    # 提取risk-level
                    risk_level_elem = first_p.find('span', class_='risk-level')
                    risk_text = ""
                    risk_color = RGBColor(220, 38, 38)  # 默认红色

                    if risk_level_elem:
                        risk_text = risk_level_elem.get_text(strip=True)
                        risk_classes = risk_level_elem.get('class', [])

                        # 根据风险等级设置颜色
                        if 'risk-high' in risk_classes:
                            risk_color = RGBColor(220, 38, 38)  # 红色
                        elif 'risk-medium' in risk_classes:
                            risk_color = RGBColor(245, 158, 11)  # 橙色
                        elif 'risk-low' in risk_classes:
                            risk_color = RGBColor(59, 130, 246)  # 蓝色

                    # 创建主文本文本框（与bullet-point保持一致的位置）
                    if main_text:
                        text_box = pptx_slide.shapes.add_textbox(
                            UnitConverter.px_to_emu(x + 20),  # 与bullet-point一致的左边距
                            UnitConverter.px_to_emu(actual_y),
                            UnitConverter.px_to_emu(width - 40),  # 与bullet-point一致的宽度
                            UnitConverter.px_to_emu(25)
                        )
                        text_frame = text_box.text_frame
                        text_frame.clear()

                        # 添加段落
                        p = text_frame.paragraphs[0]

                        # 添加图标（如果有）
                        if icon_char:
                            icon_run = p.add_run()
                            icon_run.text = icon_char + " "
                            icon_run.font.size = Pt(20)
                            icon_run.font.color.rgb = icon_color
                            icon_run.font.name = self.font_manager.get_font('icon')

                        # 添加主文本
                        text_run = p.add_run()
                        text_run.text = main_text
                        text_run.font.size = Pt(22)
                        text_run.font.bold = True
                        text_run.font.color.rgb = RGBColor(51, 51, 51)  # 深灰色
                        text_run.font.name = self.font_manager.get_font('p')

                        # 添加风险等级标签（如果有）
                        if risk_text:
                            risk_run = p.add_run()
                            risk_run.text = f" {risk_text}"
                            risk_run.font.size = Pt(20)
                            risk_run.font.bold = True
                            risk_run.font.color.rgb = risk_color

                        # 设置段落格式
                        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                        p.space_before = Pt(0)
                        p.space_after = Pt(4)

                    actual_y += 28  # 主文本后的间距

                # 处理描述文本
                desc_p = text_container.find('p', class_='text-sm')
                if desc_p:
                    desc_text = desc_p.get_text(strip=True)
                    if desc_text:
                        # 描述文本的左边距应该与主文本对齐（在图标后面）
                        desc_text_box = pptx_slide.shapes.add_textbox(
                            UnitConverter.px_to_emu(x + 20 + (25 if icon_char else 0)),  # 如果有图标则缩进
                            UnitConverter.px_to_emu(actual_y),
                            UnitConverter.px_to_emu(width - 40 - (25 if icon_char else 0)),  # 相应减少宽度
                            UnitConverter.px_to_emu(40)
                        )
                        desc_frame = desc_text_box.text_frame
                        desc_para = desc_frame.paragraphs[0]
                        desc_para.text = desc_text

                        # 设置字体样式
                        desc_run = desc_para.runs[0] if desc_para.runs else desc_para.add_run()
                        desc_run.font.size = Pt(14)
                        desc_run.font.color.rgb = RGBColor(107, 114, 128)  # 灰色
                        desc_run.font.name = self.font_manager.get_font('p')

                        # 设置段落格式
                        desc_para.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                        desc_para.space_before = Pt(0)
                        desc_para.space_after = Pt(4)

                        actual_y += 22  # 描述文本高度

            # risk-item之间的间距
            actual_y += 12

        logger.info(f"成功处理了{len(risk_items)}个risk-item")

    def _convert_grid_risk_card(self, card, pptx_slide, shape_converter, x, y, width):
        """
        转换网格中的risk-card（带红色边框和特殊样式）

        Args:
            card: risk-card元素
            pptx_slide: PPTX幻灯片
            shape_converter: 形状转换器
            x: X坐标
            y: Y坐标
            width: 宽度

        Returns:
            下一个元素的Y坐标
        """
        logger.info("处理网格中的risk-card")

        # 调用完整的risk-card转换方法，但传入指定的位置和宽度
        # 临时保存原始的x_base和card_width
        original_x_base = 80
        original_width = 1760

        # 使用ContentHeightCalculator动态计算risk-card高度
        # risk-card内容类似data-card结构，使用data-card计算方法
        card_height = self.height_calculator.calculate_data_card_height(card, width)

        # 获取CSS样式
        card_style = self.css_parser.get_class_style('risk-card') or {}

        # 添加背景（渐变效果）
        bg_color_str = card_style.get('background', 'linear-gradient(135deg, rgba(239, 68, 68, 0.08) 0%, rgba(239, 68, 68, 0.02) 100%)')

        # 创建矩形背景
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.dml.color import RGBColor
        from pptx.util import Pt
        from src.utils.unit_converter import UnitConverter
        from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

        bg_shape = pptx_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            UnitConverter.px_to_emu(x),
            UnitConverter.px_to_emu(y),
            UnitConverter.px_to_emu(width),
            UnitConverter.px_to_emu(card_height)
        )
        bg_shape.fill.solid()

        # 解析渐变背景色，使用最深的颜色
        if 'rgba(239, 68, 68' in bg_color_str:
            # 红色系风险
            bg_rgb = RGBColor(254, 242, 242)  # 非常浅的红色
        elif 'rgba(251, 146' in bg_color_str:
            # 橙色系风险
            bg_rgb = RGBColor(255, 251, 235)  # 非常浅的橙色
        elif 'rgba(250, 204' in bg_color_str:
            # 黄色系风险
            bg_rgb = RGBColor(254, 252, 232)  # 非常浅的黄色
        else:
            # 默认浅灰色
            bg_rgb = RGBColor(249, 250, 251)

        bg_shape.fill.fore_color.rgb = bg_rgb
        bg_shape.line.fill.background()
        bg_shape.shadow.inherit = False

        # 添加左边框
        border_color_str = card_style.get('border-left-color', '#ef4444')
        from src.utils.color_parser import ColorParser
        border_color = ColorParser.parse_color(border_color_str)
        if not border_color:
            # 根据风险等级确定边框颜色
            border_color = ColorParser.parse_color('#ef4444')  # 默认红色

        border_shape = pptx_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            UnitConverter.px_to_emu(x),
            UnitConverter.px_to_emu(y),
            UnitConverter.px_to_emu(4),
            UnitConverter.px_to_emu(card_height)
        )
        border_shape.fill.solid()
        border_shape.fill.fore_color.rgb = border_color
        border_shape.line.fill.background()

        current_y = y + 15

        # 缩小内容区域以适应网格
        content_width = width - 40
        content_x = x + 20

        # 处理flex布局内容
        flex_container = card.find('div', class_='flex')
        if flex_container:
            # 获取左侧内容区域
            left_div = flex_container.find('div', class_='flex-1')
            if left_div:
                # 处理风险标题
                title_div = left_div.find('div', class_='risk-title')
                if title_div:
                    # 获取图标
                    icon_elem = title_div.find('i')
                    icon_text = ""
                    icon_color = None

                    if icon_elem:
                        icon_classes = icon_elem.get('class', [])
                        # 根据图标类确定颜色
                        if 'severity-critical' in icon_classes:
                            icon_color = ColorParser.parse_color('#dc2626')
                            icon_text = "⚠"
                        elif 'severity-high' in icon_classes:
                            icon_color = ColorParser.parse_color('#ea580c')
                            icon_text = "⚠"
                        elif 'severity-medium' in icon_classes:
                            icon_color = ColorParser.parse_color('#d97706')
                            icon_text = "⚠"
                        else:
                            icon_text = "•"

                    # 获取标题文本
                    title_text = title_div.get_text(strip=True)
                    if icon_text:
                        title_text = title_text.replace(icon_text, "").strip()

                    # 添加标题文本（缩小字体以适应网格）
                    text_left = UnitConverter.px_to_emu(content_x)
                    text_top = UnitConverter.px_to_emu(current_y)

                    if icon_text and icon_color:
                        # 如果有图标，创建两段式文本
                        text_box = pptx_slide.shapes.add_textbox(
                            text_left, text_top,
                            UnitConverter.px_to_emu(content_width - 150), UnitConverter.px_to_emu(30)
                        )
                        text_frame = text_box.text_frame
                        p = text_frame.paragraphs[0]

                        # 图标run
                        icon_run = p.add_run()
                        icon_run.text = icon_text + " "
                        icon_run.font.size = Pt(20)
                        icon_run.font.name = self.font_manager.get_font('body')
                        icon_run.font.color.rgb = icon_color
                        icon_run.font.bold = True

                        # 标题run
                        title_run = p.add_run()
                        title_run.text = title_text
                        title_run.font.size = Pt(20)
                        title_run.font.name = self.font_manager.get_font('body')
                        title_run.font.bold = True
                        title_run.font.color.rgb = RGBColor(51, 51, 51)  # 深灰色
                    else:
                        # 没有图标，直接添加标题
                        text_box = pptx_slide.shapes.add_textbox(
                            text_left, text_top,
                            UnitConverter.px_to_emu(content_width - 150), UnitConverter.px_to_emu(30)
                        )
                        text_frame = text_box.text_frame
                        text_frame.text = title_text

                        for paragraph in text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.size = Pt(20)
                                run.font.name = self.font_manager.get_font('body')
                                run.font.bold = True
                                run.font.color.rgb = RGBColor(51, 51, 51)

                    current_y += 30

                # 处理风险描述（缩小字体）
                desc_div = left_div.find('div', class_='risk-desc')
                if desc_div:
                    desc_text = desc_div.get_text(strip=True)
                    if desc_text:
                        text_box = pptx_slide.shapes.add_textbox(
                            UnitConverter.px_to_emu(content_x),
                            UnitConverter.px_to_emu(current_y),
                            UnitConverter.px_to_emu(content_width - 150),
                            UnitConverter.px_to_emu(40)
                        )
                        text_frame = text_box.text_frame
                        text_frame.text = desc_text

                        for paragraph in text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.size = Pt(16)
                                run.font.name = self.font_manager.get_font('body')
                                run.font.color.rgb = RGBColor(102, 102, 102)  # 灰色

                        current_y += 25

                # 处理标签（缩小尺寸）
                tag_div = left_div.find('div', class_='mt-3')
                if tag_div:
                    span_elem = tag_div.find('span')
                    if span_elem:
                        tag_text = span_elem.get_text(strip=True)
                        tag_classes = span_elem.get('class', [])

                        # 确定标签颜色
                        tag_bg_color = RGBColor(254, 226, 226)  # 浅红色背景
                        tag_text_color = RGBColor(153, 27, 27)  # 深红色文字

                        if 'bg-orange-100' in tag_classes:
                            tag_bg_color = RGBColor(255, 237, 213)  # 浅橙色背景
                            tag_text_color = RGBColor(154, 52, 18)  # 深橙色文字
                        elif 'bg-yellow-100' in tag_classes:
                            tag_bg_color = RGBColor(254, 249, 195)  # 浅黄色背景
                            tag_text_color = RGBColor(120, 53, 15)  # 深黄色文字

                        # 创建标签背景
                        tag_box = pptx_slide.shapes.add_shape(
                            MSO_SHAPE.RECTANGLE,
                            UnitConverter.px_to_emu(content_x),
                            UnitConverter.px_to_emu(current_y),
                            UnitConverter.px_to_emu(80),
                            UnitConverter.px_to_emu(24)
                        )
                        tag_box.fill.solid()
                        tag_box.fill.fore_color.rgb = tag_bg_color
                        tag_box.line.fill.background()

                        # 添加标签文本
                        tag_text_box = pptx_slide.shapes.add_textbox(
                            UnitConverter.px_to_emu(content_x),
                            UnitConverter.px_to_emu(current_y + 2),
                            UnitConverter.px_to_emu(80),
                            UnitConverter.px_to_emu(20)
                        )
                        tag_text_frame = tag_text_box.text_frame
                        tag_text_frame.text = tag_text

                        for paragraph in tag_text_frame.paragraphs:
                            paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                            for run in paragraph.runs:
                                run.font.size = Pt(12)
                                run.font.name = self.font_manager.get_font('body')
                                run.font.color.rgb = tag_text_color
                                run.font.bold = True

            # 获取右侧CVSS分数区域（缩小字体）
            right_div = flex_container.find('div', class_='text-center')
            if right_div:
                # 获取CVSS分数
                score_div = right_div.find('div', class_='cvss-score')
                if score_div:
                    score_text = score_div.get_text(strip=True)

                    # 添加CVSS分数
                    score_box = pptx_slide.shapes.add_textbox(
                        UnitConverter.px_to_emu(x + width - 120),
                        UnitConverter.px_to_emu(y + 40),
                        UnitConverter.px_to_emu(100),
                        UnitConverter.px_to_emu(50)
                    )
                    score_frame = score_box.text_frame
                    score_frame.text = score_text

                    for paragraph in score_frame.paragraphs:
                        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                        for run in paragraph.runs:
                            run.font.size = Pt(36)
                            run.font.name = self.font_manager.get_font('body')
                            run.font.bold = True

                            # 根据分数确定颜色
                            if '10.0' in score_text or '9.' in score_text:
                                run.font.color.rgb = RGBColor(239, 68, 68)  # 红色
                            elif '8.' in score_text or '7.' in score_text:
                                run.font.color.rgb = RGBColor(234, 88, 12)  # 橙色
                            elif '6.' in score_text or '5.' in score_text:
                                run.font.color.rgb = RGBColor(217, 119, 6)  # 黄色
                            else:
                                run.font.color.rgb = RGBColor(107, 114, 128)  # 灰色

                # 获取CVSS标签
                label_div = right_div.find('div', class_='cvss-label')
                if label_div:
                    label_text = label_div.get_text(strip=True)

                    label_box = pptx_slide.shapes.add_textbox(
                        UnitConverter.px_to_emu(x + width - 120),
                        UnitConverter.px_to_emu(y + 90),
                        UnitConverter.px_to_emu(100),
                        UnitConverter.px_to_emu(25)
                    )
                    label_frame = label_box.text_frame
                    label_frame.text = label_text

                    for paragraph in label_frame.paragraphs:
                        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                        for run in paragraph.runs:
                            run.font.size = Pt(14)
                            run.font.name = self.font_manager.get_font('body')
                            run.font.color.rgb = RGBColor(102, 102, 102)

        return y + card_height + 10

    def _convert_grid_stat_card(self, card, pptx_slide, shape_converter, x, y, width):
        """
        转换网格中的stat-card（带背景色）

        Args:
            card: stat-card元素
            pptx_slide: PPTX幻灯片
            shape_converter: 形状转换器
            x: X坐标
            y: Y坐标
            width: 宽度

        Returns:
            下一个元素的Y坐标
        """
        logger.info("处理网格中的stat-card")

        # 使用ContentHeightCalculator动态计算stat-card高度
        # Note: stat-card可能包含bullet-points，使用data-card计算方法作为通用fallback
        height = self.height_calculator.calculate_data_card_height(card, width)
        
        # 添加背景色
        bg_color_str = self.css_parser.get_background_color('.stat-card')
        if bg_color_str:
            from pptx.enum.shapes import MSO_SHAPE
            bg_shape = pptx_slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                UnitConverter.px_to_emu(x),
                UnitConverter.px_to_emu(y),
                UnitConverter.px_to_emu(width),
                UnitConverter.px_to_emu(height)
            )
            bg_shape.fill.solid()
            bg_rgb, alpha = ColorParser.parse_rgba(bg_color_str)
            if bg_rgb:
                if alpha < 1.0:
                    bg_rgb = ColorParser.blend_with_white(bg_rgb, alpha)
                bg_shape.fill.fore_color.rgb = bg_rgb
            bg_shape.line.fill.background()
            logger.info(f"添加stat-card背景色: {bg_color_str}")

        # 添加左边框
        border_left_style = self.css_parser.get_style('.stat-card').get('border-left', '')
        if '4px solid' in border_left_style:
            shape_converter.add_border_left(x, y, height, 4)  # 使用动态计算的height

        # 首先检查是否包含bullet-point结构
        bullet_points = card.find_all('div', class_='bullet-point')
        if bullet_points:
            logger.info(f"stat-card包含{len(bullet_points)}个bullet-point，使用bullet-point处理逻辑")
            # 获取h3标题（如果有）
            h3_elem = card.find('h3')
            # 转换为类似data-card的格式处理，但要传入x和width参数
            return self._convert_grid_card_with_bullet_points(card, pptx_slide, shape_converter, x, y, width, bullet_points, h3_elem)

        # 然后检查是否包含risk-level标签（风险分布）
        risk_levels = card.find_all('span', class_='risk-level')
        if risk_levels:
            logger.info(f"stat-card包含{len(risk_levels)}个risk-level标签，处理为风险分布")

            # 处理h3标题
            h3_elem = card.find('h3')
            current_y = y + 20

            if h3_elem:
                h3_text = h3_elem.get_text(strip=True)
                if h3_text:
                    h3_font_size_pt = self.style_computer.get_font_size_pt(h3_elem)
                    h3_color = self._get_element_color(h3_elem) or ColorParser.get_primary_color()

                    text_left = UnitConverter.px_to_emu(x + 20)
                    text_top = UnitConverter.px_to_emu(current_y)
                    text_box = pptx_slide.shapes.add_textbox(
                        text_left, text_top,
                        UnitConverter.px_to_emu(width - 40), UnitConverter.px_to_emu(30)
                    )
                    text_frame = text_box.text_frame
                    text_frame.text = h3_text
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(h3_font_size_pt)
                            run.font.name = self.font_manager.get_font('h3')
                            # 智能判断是否应该加粗
                            if self._should_be_bold(h3_elem):
                                run.font.bold = True
                            run.font.color.rgb = h3_color

                    current_y += 35

            # 处理risk-level标签
            num_risks = len(risk_levels)
            total_width = width - 40
            risk_width = total_width // num_risks - 20
            current_x = x + 20

            for risk_level in risk_levels:
                risk_text = risk_level.get_text(strip=True)
                risk_classes = risk_level.get('class', [])

                # 获取风险等级的颜色
                risk_color = None
                bg_color = None
                if 'risk-high' in risk_classes:
                    risk_color = ColorParser.parse_color('#dc2626')  # 红色
                    bg_color = RGBColor(252, 231, 229)  # 浅红色背景
                elif 'risk-medium' in risk_classes:
                    risk_color = ColorParser.parse_color('#f59e0b')  # 橙色
                    bg_color = RGBColor(254, 243, 199)  # 浅橙色背景
                elif 'risk-low' in risk_classes:
                    risk_color = ColorParser.parse_color('#3b82f6')  # 蓝色
                    bg_color = RGBColor(239, 246, 255)  # 浅蓝色背景

                # 添加背景形状
                if bg_color:
                    from pptx.enum.shapes import MSO_SHAPE
                    bg_shape = pptx_slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        UnitConverter.px_to_emu(current_x),
                        UnitConverter.px_to_emu(current_y),
                        UnitConverter.px_to_emu(risk_width),
                        UnitConverter.px_to_emu(35)
                    )
                    bg_shape.fill.solid()
                    bg_shape.fill.fore_color.rgb = bg_color
                    bg_shape.line.fill.background()

                # 创建风险等级文本框
                text_left = UnitConverter.px_to_emu(current_x + 5)
                text_top = UnitConverter.px_to_emu(current_y + 5)
                text_box = pptx_slide.shapes.add_textbox(
                    text_left, text_top,
                    UnitConverter.px_to_emu(risk_width - 10), UnitConverter.px_to_emu(25)
                )
                text_frame = text_box.text_frame
                text_frame.text = risk_text
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

                for paragraph in text_frame.paragraphs:
                    paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(20)
                        run.font.bold = True
                        run.font.name = self.font_manager.get_font('body')

                        # 应用风险等级颜色
                        if risk_color:
                            run.font.color.rgb = risk_color

                # 移动到下一个位置
                current_x += risk_width + 20

            return y + height  # 使用动态计算的height

        # 直接提取所有文本内容，不跳过flex容器
        all_content = []

        # 方法1：提取h3和p标签
        h3_elem = card.find('h3')
        if h3_elem:
            h3_text = h3_elem.get_text(strip=True)
            if h3_text:
                all_content.append(('h3', h3_text))

        # 提取所有p标签
        for p in card.find_all('p'):
            p_text = p.get_text(strip=True)
            if p_text:
                all_content.append(('p', p_text))

        # 如果没有找到内容，使用更通用的方法
        if not all_content:
            logger.info("使用通用方法提取stat-card内容")
            # 遍历所有后代元素
            for elem in card.descendants:
                if elem.name in ['h1', 'h2', 'h3', 'h4', 'p', 'span', 'div']:
                    # 只提取没有子块级元素的文本节点
                    if not elem.find_all(['h1', 'h2', 'h3', 'h4', 'p', 'div']):
                        text = elem.get_text(strip=True)
                        if text and len(text) > 1:
                            # 判断元素类型
                            if elem.name == 'h3':
                                all_content.append(('h3', text))
                            else:
                                all_content.append(('text', text))

        logger.info(f"stat-card提取到{len(all_content)}个内容项")

        # 渲染内容
        current_y = y + 20
        for elem_type, text in all_content[:5]:  # 最多5项
            if text:
                # 根据类型设置样式
                if elem_type == 'h3':
                    font_size = 24
                    is_bold = True
                else:
                    font_size = 20
                    is_bold = False

                # 创建文本框
                text_left = UnitConverter.px_to_emu(x + 20)
                text_top = UnitConverter.px_to_emu(current_y)
                text_box = pptx_slide.shapes.add_textbox(
                    text_left, text_top,
                    UnitConverter.px_to_emu(width - 40), UnitConverter.px_to_emu(30)
                )
                text_frame = text_box.text_frame
                text_frame.text = text
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

                for paragraph in text_frame.paragraphs:
                    paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                    for run in paragraph.runs:
                        run.font.size = Pt(font_size)
                        run.font.name = self.font_manager.get_font('body')
                        run.font.color.rgb = ColorParser.get_primary_color()
                        if is_bold:
                            run.font.bold = True

                current_y += 35

        return y + height  # 使用动态计算的height

        # 查找内部的flex容器
        flex_container = card.find('div', class_='flex')
        if flex_container:
            # 收集所有内容（标题和数字）
            content_elements = []

            # 查找左侧内容区域（包含标题和数字的div）
            content_div = flex_container.find('div')
            if content_div:
                # 收集标题
                h3_elem = content_div.find('h3')
                if h3_elem:
                    content_elements.append(('h3', h3_elem))

                # 收集所有p标签（数字）
                for p_elem in content_div.find_all('p'):
                    content_elements.append(('p', p_elem))

            # 计算内容的总高度
            total_content_height = 0
            element_heights = []

            for elem_type, elem in content_elements:
                if elem_type == 'h3':
                    height = 30  # 标题高度
                else:
                    # 检查是否是大数字
                    p_classes = elem.get('class', [])
                    is_large_number = any(cls in p_classes for cls in ['text-4xl', 'text-3xl', 'text-2xl'])
                    height = 50 if is_large_number else 25

                element_heights.append(height)
                total_content_height += height

            # 使用动态计算的card高度（使用data-card计算方法）
            card_height = self.height_calculator.calculate_data_card_height(card, width)
            
            # 计算垂直起始位置（垂直居中）
            start_y = y + (card_height - total_content_height) // 2
            if start_y < y + 15:
                start_y = y + 15  # 保证最小内边距

            # 渲染内容
            current_y = start_y
            for idx, (elem_type, elem) in enumerate(content_elements):
                text = elem.get_text(strip=True)
                if not text:
                    continue

                text_left = UnitConverter.px_to_emu(x + 20)
                text_top = UnitConverter.px_to_emu(current_y)
                height = element_heights[idx]

                if elem_type == 'h3':
                    # 处理标题
                    text_box = pptx_slide.shapes.add_textbox(
                        text_left, text_top,
                        UnitConverter.px_to_emu(width - 80), UnitConverter.px_to_emu(height)
                    )
                    text_frame = text_box.text_frame
                    text_frame.text = text
                    text_frame.word_wrap = True
                    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

                    # 获取字体大小和颜色
                    font_size_pt = self.style_computer.get_font_size_pt(elem)
                    element_color = self._get_element_color(elem)

                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(font_size_pt)
                            run.font.name = self.font_manager.get_font('h3')
                            # 智能判断是否应该加粗
                            if self._should_be_bold(elem):
                                run.font.bold = True

                            # 应用颜色
                            if element_color:
                                run.font.color.rgb = element_color
                            else:
                                # 检查类名
                                classes = elem.get('class', [])
                                if 'text-gray-600' in classes:
                                    run.font.color.rgb = RGBColor(102, 102, 102)  # 灰色

                else:
                    # 处理数字
                    p_classes = elem.get('class', [])
                    is_large_number = any(cls in p_classes for cls in ['text-4xl', 'text-3xl', 'text-2xl'])

                    text_box = pptx_slide.shapes.add_textbox(
                        text_left, text_top,
                        UnitConverter.px_to_emu(width - 80), UnitConverter.px_to_emu(height)
                    )
                    text_frame = text_box.text_frame
                    text_frame.text = text
                    text_frame.word_wrap = True
                    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

                    # 获取字体大小和颜色
                    font_size_pt = self.style_computer.get_font_size_pt(elem)
                    element_color = self._get_element_color(elem)

                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(font_size_pt)

                            if is_large_number:
                                run.font.bold = True

                            run.font.name = self.font_manager.get_font('body')

                            # 应用颜色 - 检查Tailwind CSS颜色类
                            color_found = False
                            for cls in p_classes:
                                # 处理Tailwind CSS颜色类
                                if cls == 'primary-color':
                                    run.font.color.rgb = ColorParser.parse_color('rgb(10, 66, 117)')
                                    color_found = True
                                    break
                                elif cls.startswith('text-') and hasattr(self.css_parser, 'tailwind_colors'):
                                    color_str = self.css_parser.tailwind_colors.get(cls)
                                    if color_str:
                                        color_rgb = ColorParser.parse_color(color_str)
                                        if color_rgb:
                                            run.font.color.rgb = color_rgb
                                            color_found = True
                                            break

                            # 如果element_color存在，优先使用
                            if element_color and not color_found:
                                run.font.color.rgb = element_color
                            elif not color_found:
                                # 默认文本颜色
                                run.font.color.rgb = ColorParser.get_text_color()

                current_y += height

            # 处理右侧图标
            icon_elem = flex_container.find('i')
            if icon_elem:
                icon_classes = icon_elem.get('class', [])
                icon_char = self._get_icon_char(icon_classes)

                if icon_char:
                    # 获取图标字体大小
                    icon_font_size_px = 36  # 默认值
                    icon_font_size_pt = 27  # 默认值

                    # 检查是否有text-4xl等字体大小类
                    icon_classes = icon_elem.get('class', [])
                    for cls in icon_classes:
                        if cls.startswith('text-') and hasattr(self.css_parser, 'tailwind_font_sizes'):
                            font_size_str = self.css_parser.tailwind_font_sizes.get(cls)
                            if font_size_str:
                                icon_font_size_px = int(font_size_str.replace('px', ''))
                                icon_font_size_pt = UnitConverter.font_size_px_to_pt(icon_font_size_px)
                                break

                    # 图标框尺寸基于字体大小
                    icon_box_size = icon_font_size_px + 4  # 稍微留点边距
                    icon_left = UnitConverter.px_to_emu(x + width - icon_box_size - 10)
                    icon_top = UnitConverter.px_to_emu(y + (card_height - icon_box_size) // 2)  # 垂直居中（使用动态计算的card_height）

                    icon_box = pptx_slide.shapes.add_textbox(
                        icon_left, icon_top,
                        UnitConverter.px_to_emu(icon_box_size), UnitConverter.px_to_emu(icon_box_size)
                    )
                    icon_frame = icon_box.text_frame
                    icon_frame.text = icon_char
                    icon_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

                    for paragraph in icon_frame.paragraphs:
                        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                        for run in paragraph.runs:
                            run.font.size = Pt(icon_font_size_pt)
                            run.font.name = self.font_manager.get_font('body')

                            # 图标颜色
                            icon_color = self._get_element_color(icon_elem)
                            if icon_color:
                                run.font.color.rgb = icon_color
                            elif 'primary-color' in icon_classes:
                                run.font.color.rgb = ColorParser.get_primary_color()
                            elif 'text-orange-600' in icon_classes:
                                run.font.color.rgb = ColorParser.get_color_by_name('orange')
                            else:
                                run.font.color.rgb = RGBColor(200, 200, 200)  # 浅色（opacity-30效果）

        else:
            # 降级处理：查找所有文本内容
            text_elements = []
            for elem in card.descendants:
                if hasattr(elem, 'name') and elem.name in ['p', 'h1', 'h2', 'h3', 'h4']:
                    # 只提取没有子块级元素的文本节点
                    if not elem.find_all(['div', 'p', 'h1', 'h2', 'h3']):
                        text = elem.get_text(strip=True)
                        if text and len(text) > 0:
                            text_elements.append(elem)

            # 初始化current_y
            current_y = y + 20

            # 渲染文本
            for elem in text_elements[:5]:  # 最多5个元素
                text = elem.get_text(strip=True)
                if text:
                    text_left = UnitConverter.px_to_emu(x + 20)
                    text_top = UnitConverter.px_to_emu(current_y)

                    # 根据元素类型和字体大小计算高度
                    font_size_pt = self.style_computer.get_font_size_pt(elem)
                    if elem.name == 'h3':
                        height = 40
                    elif font_size_pt and font_size_pt > 30:  # text-4xl 等大字体
                        height = 50
                    elif font_size_pt and font_size_pt > 20:  # text-lg 等中等字体
                        height = 35
                    else:
                        height = 30

                    text_box = pptx_slide.shapes.add_textbox(
                        text_left, text_top,
                        UnitConverter.px_to_emu(width - 40), UnitConverter.px_to_emu(height)
                    )
                    text_frame = text_box.text_frame
                    text_frame.text = text
                    text_frame.word_wrap = True
                    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            # 使用样式计算器获取字体大小
                            run.font.size = Pt(font_size_pt) if font_size_pt else Pt(16)

                            # 设置字体
                            if elem.name == 'h3':
                                run.font.name = self.font_manager.get_font('h3')
                                run.font.bold = True
                            else:
                                run.font.name = self.font_manager.get_font('body')

                            # 处理文字颜色
                            elem_classes = elem.get('class', [])
                            element_color = self._get_element_color(elem)

                            if element_color:
                                run.font.color.rgb = element_color
                            else:
                                # 检查特定的颜色类
                                if 'primary-color' in elem_classes:
                                    run.font.color.rgb = ColorParser.get_primary_color()
                                elif 'text-red-600' in elem_classes:
                                    run.font.color.rgb = RGBColor(220, 38, 38)  # 红色
                                elif 'text-green-600' in elem_classes:
                                    run.font.color.rgb = RGBColor(22, 163, 74)  # 绿色
                                elif 'text-gray-800' in elem_classes:
                                    run.font.color.rgb = RGBColor(31, 41, 55)  # 深灰色
                                elif 'text-gray-600' in elem_classes:
                                    run.font.color.rgb = RGBColor(75, 85, 99)  # 中灰色
                                else:
                                    # 默认文本颜色
                                    run.font.color.rgb = ColorParser.get_text_color()

                    current_y += height + 10  # 增加间距

        return y + 180  # 返回固定高度

    def _get_element_relative_position(self, element, container):
        """
        获取元素相对于容器的位置

        Args:
            element: 元素
            container: 容器元素

        Returns:
            (x, y) 相对位置
        """
        # 初始化相对位置
        rel_x = 0
        rel_y = 0

        # 获取元素的样式
        style_str = element.get('style', '')
        classes = element.get('class', [])

        # 解析margin
        if style_str:
            # 解析margin-top
            import re
            margin_match = re.search(r'margin-top:\s*(\d+)px', style_str)
            if margin_match:
                rel_y += int(margin_match.group(1))

            # 解析margin-bottom
            margin_match = re.search(r'margin-bottom:\s*(\d+)px', style_str)
            if margin_match:
                # margin-bottom会在后续处理
                pass

        # 根据class判断位置
        if isinstance(classes, str):
            classes = classes.split()

        # Bootstrap/Tailwind margin classes
        for cls in classes:
            if cls.startswith('mb-') or cls.startswith('margin-bottom-'):
                # 提取margin-bottom值
                try:
                    value = int(cls.replace('mb-', '').replace('margin-bottom-', ''))
                    # Tailwind默认单位是0.25rem (4px)
                    if 'mb-' in cls:
                        rel_y += value * 4
                except:
                    pass
            elif cls.startswith('mt-') or cls.startswith('margin-top-'):
                # 提取margin-top值
                try:
                    value = int(cls.replace('mt-', '').replace('margin-top-', ''))
                    if 'mt-' in cls:
                        rel_y += value * 4
                except:
                    pass

        return rel_x, rel_y

    def _determine_title_text_alignment(self, title_elem):
        """
        智能检测标题的文本对齐方式

        Args:
            title_elem: 标题元素

        Returns:
            PP_PARAGRAPH_ALIGNMENT 枚举值
        """
        from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

        # 1. 检查内联样式
        style_str = title_elem.get('style', '')
        if 'text-align' in style_str:
            import re
            align_match = re.search(r'text-align:\s*(\w+)', style_str)
            if align_match:
                align_value = align_match.group(1).lower()
                if align_value == 'center':
                    return PP_PARAGRAPH_ALIGNMENT.CENTER
                elif align_value == 'right':
                    return PP_PARAGRAPH_ALIGNMENT.RIGHT
                elif align_value == 'justify':
                    return PP_PARAGRAPH_ALIGNMENT.JUSTIFY

        # 2. 检查CSS类
        classes = title_elem.get('class', [])
        if isinstance(classes, str):
            classes = classes.split()

        for cls in classes:
            if cls == 'text-center':
                return PP_PARAGRAPH_ALIGNMENT.CENTER
            elif cls == 'text-right':
                return PP_PARAGRAPH_ALIGNMENT.RIGHT
            elif cls == 'text-justify':
                return PP_PARAGRAPH_ALIGNMENT.JUSTIFY
            elif cls == 'text-left':
                return PP_PARAGRAPH_ALIGNMENT.LEFT

        # 3. 检查父容器的对齐设置
        parent = title_elem.parent
        while parent:
            parent_style = parent.get('style', '')
            if 'text-align' in parent_style:
                import re
                align_match = re.search(r'text-align:\s*(\w+)', parent_style)
                if align_match:
                    align_value = align_match.group(1).lower()
                    if align_value == 'center':
                        return PP_PARAGRAPH_ALIGNMENT.CENTER
                    elif align_value == 'right':
                        return PP_PARAGRAPH_ALIGNMENT.RIGHT
                    elif align_value == 'justify':
                        return PP_PARAGRAPH_ALIGNMENT.JUSTIFY

            parent_classes = parent.get('class', [])
            if isinstance(parent_classes, str):
                parent_classes = parent_classes.split()

            for cls in parent_classes:
                if cls == 'text-center':
                    return PP_PARAGRAPH_ALIGNMENT.CENTER
                elif cls == 'text-right':
                    return PP_PARAGRAPH_ALIGNMENT.RIGHT
                elif cls == 'text-justify':
                    return PP_PARAGRAPH_ALIGNMENT.JUSTIFY
                elif cls == 'text-left':
                    return PP_PARAGRAPH_ALIGNMENT.LEFT

            parent = parent.parent

        # 4. 检查CSS计算样式
        # 尝试从CSS解析器获取样式
        if hasattr(self, 'css_parser'):
            # 获取元素的样式
            selector = title_elem.name
            if classes:
                selector += '.' + '.'.join(classes)

            style = self.css_parser.get_style(selector)
            if style and 'text-align' in style:
                align_value = style['text-align'].lower()
                if align_value == 'center':
                    return PP_PARAGRAPH_ALIGNMENT.CENTER
                elif align_value == 'right':
                    return PP_PARAGRAPH_ALIGNMENT.RIGHT
                elif align_value == 'justify':
                    return PP_PARAGRAPH_ALIGNMENT.JUSTIFY

        # 5. 根据上下文推断对齐方式
        # 检查是否在flex容器中
        flex_container = title_elem.find_parent(class_='flex')
        if flex_container:
            # 检查flex容器的对齐类
            flex_classes = flex_container.get('class', [])
            if isinstance(flex_classes, str):
                flex_classes = flex_classes.split()

            # 如果有justify-between，可能是多列布局，默认左对齐
            if 'justify-between' in flex_classes or 'justify-around' in flex_classes or 'justify-evenly' in flex_classes:
                return PP_PARAGRAPH_ALIGNMENT.LEFT

        # 6. 默认值：左对齐（大多数图表标题的默认选择）
        return PP_PARAGRAPH_ALIGNMENT.LEFT

    def _convert_flex_charts_container(self, container, pptx_slide, y_start, shape_converter):
        """
        转换包含SVG图表的flex容器

        Args:
            container: flex容器元素
            pptx_slide: PPTX幻灯片
            y_start: 起始Y坐标
            shape_converter: 形状转换器

        Returns:
            下一个元素的Y坐标
        """
        logger.info("处理包含SVG图表的flex容器")

        # 初始化SVG转换器
        svg_converter = SvgConverter(pptx_slide, self.css_parser, self.html_path, self.use_stable_chart_capture)
        self.svg_converters.append(svg_converter)  # 记录实例

        # 获取所有直接子元素（应该是图表容器）
        chart_containers = []
        for child in container.children:
            if hasattr(child, 'name') and child.name == 'div':
                chart_containers.append(child)

        if not chart_containers:
            logger.warning("flex容器中未找到图表容器")
            return y_start

        logger.info(f"找到 {len(chart_containers)} 个图表容器")

        # 计算每个图表的宽度和水平位置
        total_width = 1760  # 总可用宽度
        # 从CSS获取间距，优先检查容器类
        gap = self.css_parser.get_gap_size('.gap-6')  # gap-6通常是24px
        # 检查容器inline style
        inline_style = container.get('style', '')
        if 'gap:' in inline_style:
            import re
            gap_match = re.search(r'gap:\s*(\d+)px', inline_style)
            if gap_match:
                gap = int(gap_match.group(1))
                logger.info(f"flex-charts从inline style检测到间距: {gap}px")
        else:
            logger.debug(f"flex-charts使用默认间距: {gap}px")

        # 获取flex布局信息
        container_style = container.get('style', '')
        if 'justify-content' in container_style:
            # 根据justify-content调整布局
            if 'center' in container_style:
                # 居中对齐
                chart_width = 400  # 每个图表的宽度
                total_charts_width = len(chart_containers) * chart_width + (len(chart_containers) - 1) * gap
                start_x = 80 + (total_width - total_charts_width) // 2
            elif 'space-between' in container_style:
                # 两端对齐
                chart_width = (total_width - (len(chart_containers) - 1) * gap) // len(chart_containers)
                start_x = 80
            else:
                # 默认平均分布
                chart_width = (total_width - (len(chart_containers) - 1) * gap) // len(chart_containers)
                total_charts_width = len(chart_containers) * chart_width + (len(chart_containers) - 1) * gap
                start_x = 80 + (total_width - total_charts_width) // 2
        else:
            # 默认平均分布
            chart_width = (total_width - (len(chart_containers) - 1) * gap) // len(chart_containers)
            total_charts_width = len(chart_containers) * chart_width + (len(chart_containers) - 1) * gap
            start_x = 80 + (total_width - total_charts_width) // 2

        current_y = y_start
        max_chart_height = 0

        # 处理每个图表容器（水平布局）
        for i, chart_container in enumerate(chart_containers):
            chart_x = start_x + i * (chart_width + gap)
            chart_y = current_y

            # 首先处理标题（智能识别h3或其他标题元素）
            title_elem = None
            title_text = None

            # 尝试多种标题选择器
            title_selectors = ['h3', 'h2', '.chart-title', '.title', 'h4']
            for selector in title_selectors:
                title_elem = chart_container.select_one(selector)
                if title_elem:
                    title_text = title_elem.get_text(strip=True)
                    if title_text:
                        logger.info(f"找到图表标题 ({selector}): {title_text}")
                        break

            if not title_elem:
                # 如果没有找到标准标题，尝试查找第一个包含文本的元素
                for child in chart_container.children:
                    if hasattr(child, 'name') and child.name in ['div', 'p']:
                        text = child.get_text(strip=True)
                        if text and len(text) < 50:  # 假设标题长度不超过50个字符
                            # 检查是否包含图表相关关键词
                            if any(keyword in text for keyword in ['分布', '统计', '图表', '分析', '趋势']):
                                title_elem = child
                                title_text = text
                                logger.info(f"通过内容识别找到图表标题: {title_text}")
                                break

            if title_elem and title_text:
                # 标题是容器的第一个元素，应该从容器顶部开始
                title_x = chart_x
                title_y = chart_y

                # 获取字体大小
                font_size_pt = self.style_computer.get_font_size_pt(title_elem)
                if not font_size_pt:
                    # 根据元素类型设置默认字体大小
                    if title_elem.name == 'h2':
                        font_size_pt = 24
                    elif title_elem.name == 'h3':
                        font_size_pt = 20
                    elif title_elem.name == 'h4':
                        font_size_pt = 18
                    else:
                        font_size_pt = 18

                # 计算标题高度（基于字体大小）
                title_height = int(font_size_pt * 1.5)  # 1.5倍行高

                # 智能检测文本对齐方式
                text_alignment = self._determine_title_text_alignment(title_elem)

                # 计算标题的margin-bottom
                title_classes = title_elem.get('class', [])
                if isinstance(title_classes, str):
                    title_classes = title_classes.split()

                margin_bottom = 16  # 默认margin-bottom
                for cls in title_classes:
                    if cls.startswith('mb-'):
                        try:
                            value = int(cls.replace('mb-', ''))
                            margin_bottom = value * 4  # Tailwind单位转换
                            break
                        except:
                            pass
                    elif cls.startswith('margin-bottom'):
                        # 解析内联样式
                        style_str = title_elem.get('style', '')
                        import re
                        mb_match = re.search(r'margin-bottom:\s*(\d+)px', style_str)
                        if mb_match:
                            margin_bottom = int(mb_match.group(1))

                text_box = pptx_slide.shapes.add_textbox(
                    UnitConverter.px_to_emu(title_x),
                    UnitConverter.px_to_emu(title_y),
                    UnitConverter.px_to_emu(chart_width),
                    UnitConverter.px_to_emu(title_height)
                )
                text_frame = text_box.text_frame
                text_frame.text = title_text
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

                for paragraph in text_frame.paragraphs:
                    paragraph.alignment = text_alignment
                    for run in paragraph.runs:
                        run.font.size = Pt(font_size_pt)
                        run.font.bold = True
                        # 根据元素类型选择字体
                        if title_elem.name == 'h2':
                            run.font.name = self.font_manager.get_font('h2')
                        else:
                            run.font.name = self.font_manager.get_font('h3')

                        # 检查颜色类
                        if 'primary-color' in title_classes:
                            run.font.color.rgb = ColorParser.parse_color('rgb(10, 66, 117)')
                        elif 'text-gray-600' in title_classes:
                            run.font.color.rgb = RGBColor(102, 102, 102)
                        else:
                            run.font.color.rgb = ColorParser.get_text_color()

                # 更新SVG的Y位置（标题高度 + margin-bottom）
                chart_y += title_height + margin_bottom
            else:
                logger.warning(f"图表容器 {i+1} 中未找到标题元素")

            # 查找SVG元素
            svg_elem = chart_container.find('svg')

            if svg_elem:
                logger.info(f"处理第 {i+1} 个SVG图表")

                # 转换SVG图表 - 每个容器只有一个SVG，所以索引应该是0
                chart_height = svg_converter.convert_svg(
                    svg_elem,
                    chart_container,
                    chart_x,
                    chart_y,
                    chart_width,
                    0  # 每个容器只有一个SVG，索引总是0
                )

                # 更新最大高度（包含标题）
                max_chart_height = max(max_chart_height, chart_y + chart_height - current_y)
            else:
                logger.warning(f"第 {i+1} 个图表容器中未找到SVG元素")
                max_chart_height = max(max_chart_height, 50)

        # 返回下一个元素的Y坐标（加上图表高度和间距）
        return current_y + max_chart_height + 40

    def _convert_content_container(self, container, pptx_slide, y_start, shape_converter):
        """
        转换内容容器（flex-1 overflow-hidden），处理所有子容器

        Args:
            container: 内容容器元素
            pptx_slide: PPTX幻灯片
            y_start: 起始Y坐标
            shape_converter: 形状转换器

        Returns:
            下一个元素的Y坐标
        """
        logger.info("处理内容容器(flex-1 overflow-hidden)")

        current_y = y_start

        # 获取所有直接子元素（跳过文本节点）
        children = []
        for child in container.children:
            if hasattr(child, 'name') and child.name:
                children.append(child)

        logger.info(f"找到 {len(children)} 个子容器")

        # 处理每个子容器
        for i, child in enumerate(children):
            if i > 0:
                # 从CSS读取margin-bottom值
                child_classes = child.get('class', [])
                margin_bottom = 20  # 默认20px
                if 'data-card' in child_classes:
                    constraints = self.css_parser.get_height_constraints('.data-card')
                    margin_bottom = constraints.get('margin_bottom', 20)
                elif 'stat-card' in child_classes:
                    constraints = self.css_parser.get_height_constraints('.stat-card')
                    margin_bottom = constraints.get('margin_bottom', 20)
                current_y += margin_bottom  # 使用CSS定义的间距

            # 递归调用_process_container处理每个子容器
            current_y = self._process_container(child, pptx_slide, current_y, shape_converter)

        return current_y + 20

    def _convert_flex_container(self, container, pptx_slide, y_start, shape_converter):
        """
        转换flex容器

        Args:
            container: flex容器元素
            pptx_slide: PPTX幻灯片
            y_start: 起始Y坐标
            shape_converter: 形状转换器

        Returns:
            下一个元素的Y坐标
        """
        # 获取所有子元素
        children = []
        for child in container.children:
            if hasattr(child, 'name') and child.name:
                children.append(child)

        current_y = y_start
        for child in children:
            child_classes = child.get('class', [])

            # 优先检测网格布局
            if 'grid' in child_classes:
                current_y = self._convert_grid_container(child, pptx_slide, current_y, shape_converter)
            elif 'data-card' in child_classes:
                current_y = self._convert_data_card(child, pptx_slide, shape_converter, current_y)
            elif 'stat-card' in child_classes:
                current_y = self._convert_stat_card(child, pptx_slide, current_y)
            else:
                # 降级处理
                current_y = self._convert_generic_card(child, pptx_slide, current_y, card_type='flex-item')

            # 添加间距
            current_y += 20

        return current_y

    def _convert_numbered_list_group(self, container, pptx_slide, y_start) -> int:
        """
        转换包含多个数字列表项的容器（如flex-1包含多个toc-item）

        Args:
            container: 容器元素
            pptx_slide: PPTX幻灯片
            y_start: 起始Y坐标

        Returns:
            下一个元素的Y坐标
        """
        logger.info("处理数字列表组容器")

        # 初始化文本转换器
        text_converter = TextConverter(pptx_slide, self.css_parser)

        # 获取所有toc-item
        toc_items = container.find_all('div', class_='toc-item')
        current_y = y_start

        for toc_item in toc_items:
            number_elem = toc_item.find('div', class_='toc-number')
            text_elem = toc_item.find('div', class_='toc-title')

            if number_elem and text_elem:
                numbered_item = {
                    'type': 'toc',
                    'container': toc_item,
                    'number_elem': number_elem,
                    'text_elem': text_elem,
                    'number': number_elem.get_text(strip=True),
                    'text': text_elem.get_text(strip=True)
                }
                current_y = text_converter.convert_numbered_list(numbered_item, 80, current_y)

        return current_y

    def _convert_cover_container(self, container, pptx_slide, y_start):
        """
        转换封面页容器（cover-content, cover-info）
        不添加背景，直接处理内容

        Args:
            container: 封面页容器元素
            pptx_slide: PPTX幻灯片
            y_start: 起始Y坐标

        Returns:
            下一个元素的Y坐标
        """
        logger.info(f"处理封面页容器: {container.get('class', [])}")

        # 获取容器类名以确定布局
        container_classes = container.get('class', [])
        text_converter = TextConverter(pptx_slide, self.css_parser)

        # 处理容器内的所有p标签
        paragraphs = container.find_all('p')
        current_y = y_start

        for p in paragraphs:
            text = p.get_text(strip=True)
            if not text:
                continue

            # 获取p标签的类名
            p_classes = p.get('class', [])

            # 封面页的段落需要居中对齐
            # 创建文本框
            left = UnitConverter.px_to_emu(80)
            top = UnitConverter.px_to_emu(current_y)
            width = UnitConverter.px_to_emu(1760)
            height = UnitConverter.px_to_emu(40)  # 默认高度

            text_box = pptx_slide.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame
            text_frame.text = text
            text_frame.word_wrap = True
            text_frame.margin_top = 0
            text_frame.margin_bottom = 0
            text_frame.margin_left = 0

            # 设置字体样式
            style_computer = get_style_computer(self.css_parser)
            font_manager = get_font_manager(self.css_parser)

            # 获取字体大小
            p_font_size_pt = style_computer.get_font_size_pt(p)

            for paragraph in text_frame.paragraphs:
                # 居中对齐
                paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(p_font_size_pt)
                    run.font.name = font_manager.get_font('p')
                    # 检查是否有primary-color类
                    if 'primary-color' in p_classes:
                        run.font.color.rgb = ColorParser.get_primary_color()

            # 计算实际高度并更新Y坐标
            p_font_size_px = UnitConverter.pt_to_px(p_font_size_pt)
            line_height = int(p_font_size_px * 1.5)

            # 根据容器类型调整间距
            if 'cover-content' in container_classes:
                # cover-content内的段落间距较小
                current_y += line_height
            elif 'cover-info' in container_classes:
                # cover-info内的段落间距为mb-6 = 24px
                if p != paragraphs[-1]:  # 不是最后一个段落
                    current_y += line_height + 24
                else:
                    current_y += line_height

            logger.info(f"添加封面页段落: {text}")

        # 处理装饰图标（如果有的话）
        icon = container.find('i', class_='fa-shield-alt')
        if icon:
            logger.info("找到封面页盾牌图标")
            # 图标居中对齐
            icon_size = 80  # text-5xl = 5rem = 80px
            icon_left = UnitConverter.px_to_emu(960 - icon_size // 2)  # 居中
            icon_top = UnitConverter.px_to_emu(current_y)
            icon_width = UnitConverter.px_to_emu(icon_size)
            icon_height = UnitConverter.px_to_emu(icon_size)

            # 创建图标文本框
            icon_box = pptx_slide.shapes.add_textbox(
                icon_left, icon_top, icon_width, icon_height
            )
            icon_frame = icon_box.text_frame
            icon_frame.text = "🛡️"  # 使用盾牌emoji
            icon_frame.margin_top = 0
            icon_frame.margin_bottom = 0
            icon_frame.margin_left = 0

            # 设置图标样式
            for paragraph in icon_frame.paragraphs:
                paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(icon_size * 0.75)  # 80px = 60pt
                    run.font.color.rgb = ColorParser.get_primary_color()
                    run.font.name = "Arial"

            logger.info("添加封面页盾牌图标")
            # 添加图标后的间距
            current_y += icon_size + 40

        return current_y

    def _convert_centered_container(self, container, pptx_slide, y_start, shape_converter):
        """
        转换垂直居中的flex容器（flex-col justify-center）
        内容在可用空间内垂直居中

        Args:
            container: 居中容器元素
            pptx_slide: PPTX幻灯片
            y_start: 起始Y坐标
            shape_converter: 形状转换器

        Returns:
            下一个元素的Y坐标
        """
        logger.info("处理垂直居中的flex容器")

        # 获取所有直接子元素
        children = []
        for child in container.children:
            if hasattr(child, 'name') and child.name:
                children.append(child)

        if not children:
            return y_start

        # 计算所有元素的总高度
        total_height = 0
        for child in children:
            child_classes = child.get('class', [])

            # 估算每个元素的高度
            if 'data-card' in child_classes:
                total_height += 100  # data-card高度
            elif 'grid' in child_classes:
                # grid包含两个stat-card
                total_height += 220  # grid高度
            else:
                # 普通div（p标签）
                p_elem = child.find('p')
                if p_elem:
                    p_classes = p_elem.get('class', [])
                    if 'text-2xl' in p_classes:
                        total_height += 50
                    elif 'text-xl' in p_classes:
                        total_height += 40
                    else:
                        total_height += 50
                else:
                    total_height += 50

            # 添加间距
            spacing_value = self._get_spacing_value_for_mb(child_classes)
            total_height += spacing_value
            total_height += 20  # 默认元素间距

        # 计算可用高度（从y_start到页面底部，留出页码空间）
        available_height = 1080 - y_start - 100  # 留出底部空间

        # 计算起始Y坐标（垂直居中）
        if total_height < available_height:
            # 内容在可用空间内垂直居中
            current_y = y_start + (available_height - total_height) // 2
            logger.info(f"内容垂直居中: 总高度={total_height}px, 可用高度={available_height}px, 起始Y={current_y}px")
        else:
            # 内容太高，从顶部开始
            current_y = y_start
            logger.info(f"内容过高，从顶部开始: 总高度={total_height}px")

        # 顺序处理每个子元素，保持HTML结构和间距
        for child in children:
            child_classes = child.get('class', [])

            # 处理上边距（mb-*）
            spacing_value = self._get_spacing_value_for_mb(child_classes)
            if spacing_value > 0 and current_y > y_start:
                current_y += spacing_value

            # 根据子元素类型调用相应的处理方法
            if 'data-card' in child_classes:
                logger.info(f"处理data-card: {child_classes}")
                current_y = self._convert_centered_data_card(child, pptx_slide, current_y)
            elif 'grid' in child_classes:
                logger.info(f"处理grid布局: {child_classes}")
                current_y = self._convert_grid_container(child, pptx_slide, current_y, shape_converter)
            elif 'stat-card' in child_classes:
                logger.info(f"处理stat-card: {child_classes}")
                current_y = self._convert_stat_card(child, pptx_slide, current_y)
            else:
                # 处理普通div（如text-center）
                logger.info(f"处理普通div: {child_classes}")
                current_y = self._convert_simple_div(child, pptx_slide, current_y)

            # 添加默认间距
            current_y += 20

        return current_y

    def _get_spacing_value_for_mb(self, classes):
        """
        获取Tailwind margin-bottom类的像素值

        Args:
            classes: CSS类列表

        Returns:
            int: 间距像素值
        """
        spacing_map = {
            'mb-1': 4, 'mb-2': 8, 'mb-3': 12, 'mb-4': 16,
            'mb-5': 20, 'mb-6': 24, 'mb-8': 32, 'mb-10': 40,
            'mb-12': 48, 'mb-16': 64, 'mb-20': 80
        }

        for cls in classes:
            if cls in spacing_map:
                return spacing_map[cls]
        return 0

    def _convert_simple_div(self, div, pptx_slide, y_start):
        """
        处理简单的div（如text-center的问答提示和结尾语）
        增强处理：检查是否有嵌套的data-card或其他特殊元素

        Args:
            div: div元素
            pptx_slide: PPTX幻灯片
            y_start: 起始Y坐标

        Returns:
            下一个元素的Y坐标
        """
        # 首先检查是否有嵌套的data-card
        nested_data_cards = div.find_all('div', class_='data-card')
        if nested_data_cards:
            logger.info(f"普通div中发现{len(nested_data_cards)}个嵌套的data-card")
            current_y = y_start

            # 获取space-y间距
            div_classes = div.get('class', [])
            space_y_spacing = self._get_spacing_value_for_space_y(div_classes)

            for i, card in enumerate(nested_data_cards):
                current_y = self._convert_centered_data_card(card, pptx_slide, current_y)
                # 添加间距，但最后一个不添加
                if i < len(nested_data_cards) - 1 and space_y_spacing > 0:
                    current_y += space_y_spacing

            return current_y

        # 查找p标签
        p_elem = div.find('p')
        if p_elem:
            text = p_elem.get_text(strip=True)
            if text:
                # 获取样式
                div_classes = div.get('class', [])
                p_classes = p_elem.get('class', [])

                # 字体大小
                font_size = 25  # 默认
                if 'text-2xl' in p_classes:
                    font_size = 24
                elif 'text-xl' in p_classes:
                    font_size = 20
                elif 'text-3xl' in p_classes:
                    font_size = 30

                # 创建文本框
                text_box = pptx_slide.shapes.add_textbox(
                    UnitConverter.px_to_emu(80),
                    UnitConverter.px_to_emu(y_start),
                    UnitConverter.px_to_emu(1760),
                    UnitConverter.px_to_emu(50)
                )

                # 设置文本和样式
                text_frame = text_box.text_frame
                text_frame.text = text
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

                for paragraph in text_frame.paragraphs:
                    # 如果有text-center类，居中对齐
                    if 'text-center' in div_classes or 'text-center' in p_classes:
                        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                    else:
                        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

                    for run in paragraph.runs:
                        run.font.size = Pt(font_size)
                        run.font.name = self.font_manager.get_font('body')

                        # 颜色处理
                        if 'text-gray-600' in p_classes:
                            run.font.color.rgb = RGBColor(102, 102, 102)
                        elif 'primary-color' in p_classes:
                            run.font.color.rgb = ColorParser.get_primary_color()

        return y_start + 60

    def _get_spacing_value_for_space_y(self, classes):
        """
        获取Tailwind space-y类的像素值

        Args:
            classes: CSS类列表

        Returns:
            int: 间距像素值
        """
        spacing_map = {
            'space-y-1': 4, 'space-y-2': 8, 'space-y-3': 12, 'space-y-4': 16,
            'space-y-5': 20, 'space-y-6': 24, 'space-y-8': 32, 'space-y-10': 40,
            'space-y-12': 48, 'space-y-16': 64, 'space-y-20': 80
        }

        for cls in classes:
            if cls in spacing_map:
                return spacing_map[cls]
        return 0

    def _convert_centered_data_card(self, card, pptx_slide, y_start: int) -> int:
        """
        转换居中容器中的data-card（智能识别是否需要左边框）

        Args:
            card: data-card元素
            pptx_slide: PPTX幻灯片
            y_start: 起始Y坐标

        Returns:
            下一个元素的Y坐标
        """
        logger.info("处理居中容器中的data-card")
        from src.utils.color_parser import ColorParser

        # 检查是否有max-w-2xl类，如果有则限制宽度
        card_classes = card.get('class', [])
        if 'max-w-2xl' in card_classes:
            # max-w-2xl在Tailwind中是42rem = 672px
            card_width = 672
            # 居中定位
            x_base = (1920 - card_width) // 2
        else:
            # 默认宽度，但留有左右边距
            card_width = 600  # 适中的宽度
            x_base = (1920 - card_width) // 2

        current_y = y_start + 10

        # 智能识别是否需要左边框
        # 检查CSS定义中是否有border-left
        border_style = self.css_parser.get_style('.data-card').get('border-left', '')
        has_left_border = bool(border_style)
        
        # 添加背景
        bg_color_str = self.css_parser.get_background_color('.data-card')
        if bg_color_str:
            from pptx.enum.shapes import MSO_SHAPE
            bg_shape = pptx_slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                UnitConverter.px_to_emu(x_base),
                UnitConverter.px_to_emu(y_start),
                UnitConverter.px_to_emu(card_width),
                UnitConverter.px_to_emu(80)  # 估算高度
            )
            bg_shape.fill.solid()
            bg_rgb, alpha = ColorParser.parse_rgba(bg_color_str)
            if bg_rgb:
                if alpha < 1.0:
                    bg_rgb = ColorParser.blend_with_white(bg_rgb, alpha)
                bg_shape.fill.fore_color.rgb = bg_rgb
            bg_shape.line.fill.background()
            bg_shape.shadow.inherit = False

        # 如果需要左边框，添加左边框
        if has_left_border:
            # 解析边框样式
            border_width = 4  # 默认4px
            # 从border-left样式中提取颜色
            border_color_str = border_style.split('solid ')[-1].strip(')') if 'solid' in border_style else 'rgb(10, 66, 117)'
            # 确保颜色格式正确
            if not border_color_str.startswith('rgb'):
                border_color_str = f"rgb({border_color_str})"
            border_color = ColorParser.parse_color(border_color_str)
            from pptx.enum.shapes import MSO_SHAPE
            border_shape = pptx_slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                UnitConverter.px_to_emu(x_base),
                UnitConverter.px_to_emu(y_start),
                UnitConverter.px_to_emu(border_width),
                UnitConverter.px_to_emu(80)
            )
            border_shape.fill.solid()
            border_shape.fill.fore_color.rgb = border_color
            border_shape.line.fill.background()

        # 提取并渲染文本内容
        text_elements = []
        for elem in card.descendants:
            if hasattr(elem, 'name') and elem.name in ['p', 'h1', 'h2', 'h3', 'h4', 'div', 'span']:
                # 只提取没有子块级元素的文本节点
                if not elem.find_all(['div', 'p', 'h1', 'h2', 'h3']):
                    text = elem.get_text(strip=True)
                    if text and len(text) > 2:
                        text_elements.append(elem)

        # 渲染文本（如果有左边框，文本需要稍微右移）
        text_left_offset = 20 if has_left_border else 20
        text_left_offset += 8 if has_left_border else 0  # 左边框额外留出空间

        for elem in text_elements[:5]:  # 最多5个元素
            text = elem.get_text(strip=True)
            if text:
                # 获取元素的class以判断样式
                elem_classes = elem.get('class', [])
                
                # 文本框宽度要比卡片宽度小一些，留有内边距
                text_width = card_width - 40 - (8 if has_left_border else 0)  # 如果有左边框，减少文本宽度
                text_left = UnitConverter.px_to_emu(x_base + text_left_offset)
                text_top = UnitConverter.px_to_emu(current_y)
                text_box = pptx_slide.shapes.add_textbox(
                    text_left, text_top,
                    UnitConverter.px_to_emu(text_width), UnitConverter.px_to_emu(30)
                )
                text_frame = text_box.text_frame
                text_frame.text = text
                text_frame.word_wrap = True
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

                # 检查元素是否有text-center类或父容器有text-center
                parent = elem.parent
                has_text_center = 'text-center' in elem_classes

                # 检查父容器是否有text-center类
                while parent and not has_text_center:
                    parent_classes = parent.get('class', [])
                    if isinstance(parent_classes, str):
                        parent_classes = parent_classes.split()
                    if 'text-center' in parent_classes:
                        has_text_center = True
                        break
                    parent = parent.parent

                # 判断是否需要加粗（stat-value通常加粗）
                is_bold = 'stat-value' in elem_classes or 'font-bold' in elem_classes or 'font-semibold' in elem_classes
                is_primary = 'stat-value' in elem_classes or 'primary-color' in elem_classes
                is_gray = 'stat-label' in elem_classes or 'text-gray-600' in elem_classes

                for paragraph in text_frame.paragraphs:
                    # 如果有text-center类或者元素本身是居中的，则设置居中对齐
                    if has_text_center:
                        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                    else:
                        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
                    for run in paragraph.runs:
                        # 使用样式计算器获取字体大小
                        font_size_px = self.style_computer.get_font_size_pt(elem)
                        run.font.size = Pt(font_size_px) if font_size_px else Pt(16)
                        run.font.name = self.font_manager.get_font('body')
                        
                        # 应用加粗
                        if is_bold:
                            run.font.bold = True

                        # 应用颜色
                        if is_primary:
                            run.font.color.rgb = ColorParser.get_primary_color()
                        elif is_gray:
                            run.font.color.rgb = RGBColor(102, 102, 102)  # #666
                        elif 'primary-color' in elem_classes:
                            run.font.color.rgb = ColorParser.get_primary_color()
                        else:
                            element_color = self._get_element_color(elem)
                            if element_color:
                                run.font.color.rgb = element_color

                # 根据元素类型计算合适的间距
                if 'stat-value' in elem_classes:
                    # stat-value后的间距较小（CSS: margin-bottom: 8px）
                    current_y += 58
                elif 'stat-label' in elem_classes:
                    # stat-label后的间距
                    current_y += 36
                elif 'mt-3' in elem_classes:
                    # p标签带mt-3
                    current_y += 35
                else:
                    # 默认间距
                    current_y += 35

        return current_y + 10

    def _convert_stats_container(self, container, pptx_slide, y_start: int) -> int:
        """
        转换统计卡片容器 (.stats-container)

        Returns:
            下一个元素的Y坐标
        """
        stat_boxes = container.find_all('div', class_='stat-box')
        num_boxes = len(stat_boxes)

        if num_boxes == 0:
            return y_start

        # 动态获取列数：优先从inline style，其次从CSS规则
        num_columns = 4  # 默认值

        # 1. 检查inline style属性
        inline_style = container.get('style', '')
        if 'grid-template-columns' in inline_style:
            # 解析inline style中的grid-template-columns
            import re
            repeat_match = re.search(r'repeat\((\d+),', inline_style)
            if repeat_match:
                num_columns = int(repeat_match.group(1))
                logger.info(f"从inline style检测到列数: {num_columns}列")
            else:
                fr_count = len(re.findall(r'1fr', inline_style))
                if fr_count > 0:
                    num_columns = fr_count
                    logger.info(f"从inline style检测到列数: {num_columns}列")
        else:
            # 2. 从CSS规则获取
            num_columns = self.css_parser.get_grid_columns('.stats-container')
            logger.info(f"从CSS规则检测到列数: {num_columns}列")

        # 根据列数动态计算box宽度
        # 总宽度 = 1920 - 2*80(左右边距) = 1760
        # box_width = (1760 - (num_columns-1) * gap) / num_columns
        gap = self.css_parser.get_gap_size('.stats-container')
        total_width = 1760
        box_width = int((total_width - (num_columns - 1) * gap) / num_columns)
        x_start = 80
        
        # 动态计算每个box的高度
        box_heights = []
        for box in stat_boxes:
            box_height = self.height_calculator.calculate_stat_box_height(box, box_width)
            box_heights.append(box_height)
        
        logger.info(f"计算box尺寸: 宽度={box_width}px, 间距={gap}px")
        logger.info(f"box高度范围: {min(box_heights)}px - {max(box_heights)}px")

        # 计算每一行的最大高度
        num_rows = (num_boxes + num_columns - 1) // num_columns
        row_max_heights = []
        for row in range(num_rows):
            row_start = row * num_columns
            row_end = min(row_start + num_columns, num_boxes)
            row_max_height = max(box_heights[row_start:row_end])
            row_max_heights.append(row_max_height)
        
        # 记录每一行的起始Y坐标
        row_y_positions = [y_start]
        for row in range(1, num_rows):
            prev_row_y = row_y_positions[row - 1]
            prev_row_height = row_max_heights[row - 1]
            row_y_positions.append(prev_row_y + prev_row_height + gap)
        
        for idx, box in enumerate(stat_boxes):
            col = idx % num_columns
            row = idx // num_columns

            x = x_start + col * (box_width + gap)
            y = row_y_positions[row]
            box_height = box_heights[idx]

            # 添加背景（使用动态计算的box_height）
            shape_converter = ShapeConverter(pptx_slide, self.css_parser)
            shape_converter.add_stat_box_background(x, y, box_width, box_height)

            # 提取内容
            icon = box.find('i')
            title_elem = box.find('div', class_='stat-title')
            h2 = box.find('h2')
            # p标签将在下面统一处理

            # 智能判断布局方向：检查CSS的align-items设置
            layout_direction = self._determine_layout_direction(box)

            # 智能判断文字对齐方式
            text_alignment = self._determine_text_alignment(box)

            if layout_direction == 'horizontal':
                # 水平布局：图标在左，文字在右
                # 根据CSS样式计算间距：padding: 20px, icon margin-right: 20px
                icon_x = x + 20  # 左padding
                content_x = icon_x + 36 + 20  # icon_x + icon_width + margin-right
                content_width = box_width - 40 - 36 - 20  # box_width - 左padding - icon_width - margin-right

                # 添加图标（左侧）
                if icon:
                    icon_classes = icon.get('class', [])
                    icon_char = self._get_icon_char(icon_classes)

                    # 图标垂直居中（根据CSS font-size: 36px）
                    icon_height = 36
                    icon_top = y + (box_height - icon_height) // 2  # 垂直居中计算
                    icon_left = UnitConverter.px_to_emu(icon_x)
                    icon_top = UnitConverter.px_to_emu(icon_top)
                    icon_box = pptx_slide.shapes.add_textbox(
                        icon_left, icon_top,
                        UnitConverter.px_to_emu(36), UnitConverter.px_to_emu(icon_height)
                    )
                    icon_frame = icon_box.text_frame
                    icon_frame.text = icon_char
                    icon_frame.vertical_anchor = MSO_ANCHOR.MIDDLE  # 垂直居中
                    for paragraph in icon_frame.paragraphs:
                        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                        for run in paragraph.runs:
                            run.font.size = Pt(36)
                            run.font.color.rgb = ColorParser.get_primary_color()
                            run.font.name = self.font_manager.get_font('body')

                # 添加文字内容（右侧），也垂直居中
                content_height = 0
                if title_elem:
                    title_text = title_elem.get_text(strip=True)
                    title_font_size_pt = self.style_computer.get_font_size_pt(title_elem)
                    title_height = int(title_font_size_pt * 1.5)  # 估算标题高度
                    content_height += title_height + 5  # margin-bottom: 5px

                if h2:
                    h2_text = h2.get_text(strip=True)
                    h2_font_size_pt = self.style_computer.get_font_size_pt(h2)
                    h2_height = int(h2_font_size_pt * 1.5)  # 估算h2高度
                    content_height += h2_height + 5

                # 计算所有p标签的总高度（包括第一个p标签）
                all_p_tags = box.find_all('p')
                for p_tag in all_p_tags:
                    p_text = p_tag.get_text(strip=True)
                    if p_text:
                        p_font_size_pt = self.style_computer.get_font_size_pt(p_tag)
                        # 计算p标签的行数（估算每行80个字符）
                        p_lines = max(1, (len(p_text) + 79) // 80)
                        p_height = p_lines * int(p_font_size_pt * 1.5)
                        content_height += p_height + 5  # 5px间距

                # 垂直居中文字内容
                content_start_y = y + (box_height - content_height) // 2
                current_y = content_start_y

                # 添加标题
                if title_elem:
                    title_text = title_elem.get_text(strip=True)
                    title_left = UnitConverter.px_to_emu(content_x)
                    title_top = UnitConverter.px_to_emu(current_y)
                    title_font_size_pt = self.style_computer.get_font_size_pt(title_elem)
                    title_height = int(title_font_size_pt * 1.5)
                    title_box = pptx_slide.shapes.add_textbox(
                        title_left, title_top,
                        UnitConverter.px_to_emu(content_width), UnitConverter.px_to_emu(title_height)
                    )
                    title_frame = title_box.text_frame
                    title_frame.text = title_text
                    title_frame.word_wrap = True
                    title_frame.vertical_anchor = MSO_ANCHOR.TOP  # 顶部对齐，确保精确定位
                    for paragraph in title_frame.paragraphs:
                        paragraph.alignment = text_alignment
                        for run in paragraph.runs:
                            run.font.size = Pt(title_font_size_pt)
                            run.font.color.rgb = ColorParser.get_primary_color()
                            run.font.name = self.font_manager.get_font('body')

                    current_y += int(self.style_computer.get_font_size_pt(title_elem) * 1.5) + 5

                # 添加主数据
                if h2:
                    h2_text = h2.get_text(strip=True)
                    h2_left = UnitConverter.px_to_emu(content_x)
                    h2_top = UnitConverter.px_to_emu(current_y)
                    h2_font_size_pt = self.style_computer.get_font_size_pt(h2)
                    h2_height = int(h2_font_size_pt * 1.5)
                    h2_box = pptx_slide.shapes.add_textbox(
                        h2_left, h2_top,
                        UnitConverter.px_to_emu(content_width), UnitConverter.px_to_emu(h2_height)
                    )
                    h2_frame = h2_box.text_frame
                    h2_frame.text = h2_text
                    h2_frame.vertical_anchor = MSO_ANCHOR.TOP  # 顶部对齐
                    for paragraph in h2_frame.paragraphs:
                        paragraph.alignment = text_alignment
                        for run in paragraph.runs:
                            run.font.size = Pt(h2_font_size_pt)
                            run.font.bold = True
                            run.font.color.rgb = ColorParser.get_primary_color()
                            run.font.name = self.font_manager.get_font('body')

                    current_y += int(self.style_computer.get_font_size_pt(h2) * 1.5) + 5

                # 添加描述（统一处理所有p标签）
                all_p_tags = box.find_all('p')
                for p_tag in all_p_tags:
                    p_text = p_tag.get_text(strip=True)
                    if p_text:
                        p_font_size_pt = self.style_computer.get_font_size_pt(p_tag)
                        # 更精确的行数计算：每行大约80个字符
                        p_lines = max(1, (len(p_text) + 79) // 80)
                        p_height = p_lines * int(p_font_size_pt * 1.5)

                        p_left = UnitConverter.px_to_emu(content_x)
                        p_top = UnitConverter.px_to_emu(current_y)
                        p_box = pptx_slide.shapes.add_textbox(
                            p_left, p_top,
                            UnitConverter.px_to_emu(content_width), UnitConverter.px_to_emu(p_height)
                        )
                        p_frame = p_box.text_frame
                        p_frame.text = p_text
                        p_frame.word_wrap = True
                        p_frame.vertical_anchor = MSO_ANCHOR.TOP  # 顶部对齐
                        for paragraph in p_frame.paragraphs:
                            paragraph.alignment = text_alignment
                            for run in paragraph.runs:
                                run.font.size = Pt(p_font_size_pt)
                                run.font.name = self.font_manager.get_font('body')

                        current_y += p_height + 5  # 间距

            else:
                # 垂直布局：图标在上，文字在下（原有逻辑，但优化间距）
                current_y = y + 25  # 增加顶部间距，避免重合
                if icon:
                    icon_classes = icon.get('class', [])
                    icon_char = self._get_icon_char(icon_classes)

                    # 图标居中
                    icon_left = UnitConverter.px_to_emu(x + box_width // 2 - 25)
                    icon_top = UnitConverter.px_to_emu(current_y)
                    icon_box = pptx_slide.shapes.add_textbox(
                        icon_left, icon_top,
                        UnitConverter.px_to_emu(50), UnitConverter.px_to_emu(40)
                    )
                    icon_frame = icon_box.text_frame
                    icon_frame.text = icon_char
                    icon_frame.vertical_anchor = 1  # 居中
                    for paragraph in icon_frame.paragraphs:
                        paragraph.alignment = 2  # PP_ALIGN.CENTER
                        for run in paragraph.runs:
                            run.font.size = Pt(36)
                            run.font.color.rgb = ColorParser.get_primary_color()
                            run.font.name = self.font_manager.get_font('body')

                    current_y += 50  # 增加图标与文字间距

                # 添加标题
                if title_elem:
                    title_text = title_elem.get_text(strip=True)
                    title_left = UnitConverter.px_to_emu(x + 15)
                    title_top = UnitConverter.px_to_emu(current_y)
                    title_box = pptx_slide.shapes.add_textbox(
                        title_left, title_top,
                        UnitConverter.px_to_emu(box_width - 30), UnitConverter.px_to_emu(25)
                    )
                    title_frame = title_box.text_frame
                    title_frame.text = title_text
                    title_frame.word_wrap = True
                    for paragraph in title_frame.paragraphs:
                        paragraph.alignment = text_alignment
                        for run in paragraph.runs:
                            title_font_size_pt = self.style_computer.get_font_size_pt(title_elem)
                            run.font.size = Pt(title_font_size_pt)
                            run.font.color.rgb = ColorParser.get_primary_color()
                            run.font.name = self.font_manager.get_font('body')

                    current_y += 30

                # 添加主数据
                if h2:
                    h2_text = h2.get_text(strip=True)
                    h2_left = UnitConverter.px_to_emu(x + 15)
                    h2_top = UnitConverter.px_to_emu(current_y)
                    h2_box = pptx_slide.shapes.add_textbox(
                        h2_left, h2_top,
                        UnitConverter.px_to_emu(box_width - 30), UnitConverter.px_to_emu(40)
                    )
                    h2_frame = h2_box.text_frame
                    h2_frame.text = h2_text
                    for paragraph in h2_frame.paragraphs:
                        paragraph.alignment = text_alignment
                        for run in paragraph.runs:
                            h2_font_size_pt = self.style_computer.get_font_size_pt(h2)
                            run.font.size = Pt(h2_font_size_pt)
                            run.font.bold = True
                            run.font.color.rgb = ColorParser.get_primary_color()
                            run.font.name = self.font_manager.get_font('body')

                    current_y += 45

                # 添加描述（统一处理所有p标签）
                all_p_tags = box.find_all('p')
                for p_tag in all_p_tags:
                    p_text = p_tag.get_text(strip=True)
                    if p_text:
                        p_font_size_pt = self.style_computer.get_font_size_pt(p_tag)
                        # 更精确的行数计算：每行大约80个字符
                        p_lines = max(1, (len(p_text) + 79) // 80)
                        p_height = p_lines * int(p_font_size_pt * 1.5)

                        p_left = UnitConverter.px_to_emu(x + 15)
                        p_top = UnitConverter.px_to_emu(current_y)
                        p_box = pptx_slide.shapes.add_textbox(
                            p_left, p_top,
                            UnitConverter.px_to_emu(box_width - 30), UnitConverter.px_to_emu(p_height)
                        )
                        p_frame = p_box.text_frame
                        p_frame.text = p_text
                        p_frame.word_wrap = True
                        p_frame.vertical_anchor = MSO_ANCHOR.TOP  # 顶部对齐
                        for paragraph in p_frame.paragraphs:
                            paragraph.alignment = text_alignment
                            for run in paragraph.runs:
                                run.font.size = Pt(p_font_size_pt)
                                run.font.name = self.font_manager.get_font('body')

                        current_y += p_height + 5  # 间距

        # 计算下一个元素的Y坐标
        # 实际高度 = 所有行的高度之和 + 行间距之和
        # 正确公式：sum(row_heights) + (num_rows - 1) * gap
        actual_height = sum(row_max_heights) + (num_rows - 1) * gap if num_rows > 0 else 0

        logger.info(f"stats-container高度计算: 行数={num_rows}, gap={gap}px, 总高度={actual_height}px")
        logger.info(f"  各行高度: {row_max_heights}")

        return y_start + actual_height

    def _convert_stat_card(self, card, pptx_slide, y_start: int) -> int:
        """转换统计卡片(.stat-card) - 支持多种内部结构"""

        logger.info(f"开始处理stat-card，y_start={y_start}")

        # 0. 检查是否包含bullet-point结构
        bullet_points = card.find_all('div', class_='bullet-point')
        if bullet_points:
            logger.info(f"stat-card包含{len(bullet_points)}个bullet-point，使用bullet-point处理逻辑")
            # 获取h3标题（如果有）
            h3_elem = card.find('h3')
            # 转换为类似data-card的格式处理
            return self._convert_card_with_bullet_points(card, pptx_slide, y_start, bullet_points, h3_elem)

        # 1. 检查是否包含目录布局 (toc-item)
        toc_items = card.find_all('div', class_='toc-item')
        if toc_items:
            logger.info("stat-card包含toc-item目录结构，处理目录布局")
            return self._convert_toc_layout(card, toc_items, pptx_slide, y_start)

        # 1. 检查是否包含stats-container (stat-box容器类型)
        stats_container = card.find('div', class_='stats-container')
        if stats_container:
            logger.info("stat-card包含stats-container,处理嵌套的stat-box结构")

            # 获取所有的stat-box
            stat_boxes = stats_container.find_all('div', class_='stat-box')
            num_boxes = len(stat_boxes)

            # 动态获取列数 - 优先检查内联样式
            num_columns = 3  # 默认3列（slide01.html使用3列）
            inline_style = stats_container.get('style', '')
            if 'grid-template-columns' in inline_style:
                import re
                # 查找 repeat(n, 1fr) 或直接的 1fr 1fr 1fr 格式
                repeat_match = re.search(r'repeat\((\d+),', inline_style)
                if repeat_match:
                    num_columns = int(repeat_match.group(1))
                else:
                    fr_count = len(re.findall(r'1fr', inline_style))
                    if fr_count > 0:
                        num_columns = fr_count
                logger.info(f"从内联样式解析出列数: {num_columns}")
            else:
                # 从CSS类获取列数
                num_columns = self.css_parser.get_grid_columns('.stats-container')
                logger.info(f"从CSS类解析出列数: {num_columns}")

            # 从CSS读取约束
            constraints = self.css_parser.get_height_constraints('.stat-card')
            stat_card_padding_top = constraints.get('padding_top', 20)
            stat_card_padding_bottom = constraints.get('padding_bottom', 20)
            stats_container_gap = self.css_parser.get_gap_size('.stats-container')

            # 动态计算每个stat-box的高度（使用ContentHeightCalculator）
            total_width = 1760
            box_width = int((total_width - (num_columns - 1) * stats_container_gap) / num_columns)
            
            box_heights = []
            for box in stat_boxes:
                box_height = self.height_calculator.calculate_stat_box_height(box, box_width)
                box_heights.append(box_height)

            # 计算stats-container的实际高度
            num_rows = (num_boxes + num_columns - 1) // num_columns
            row_max_heights = []
            for row in range(num_rows):
                row_start = row * num_columns
                row_end = min(row_start + num_columns, num_boxes)
                row_max_height = max(box_heights[row_start:row_end])
                row_max_heights.append(row_max_height)
            
            stats_container_height = sum(row_max_heights) + (num_rows - 1) * stats_container_gap if num_rows > 0 else 0

            # 计算stat-card总高度（包括自身padding）
            title_elem = card.find('p', class_='primary-color')
            has_title = title_elem is not None
            title_height = 0
            if has_title:
                title_font_size_pt = self.style_computer.get_font_size_pt(title_elem)
                title_height = int(title_font_size_pt * 1.5) + 12  # 字体高度 + margin-bottom

            card_height = stat_card_padding_top + title_height + stats_container_height + stat_card_padding_bottom

            logger.info(f"stat-card动态高度计算: boxes={num_boxes}, columns={num_columns}, rows={num_rows}")
            logger.info(f"stat-card高度组成: padding={stat_card_padding_top+stat_card_padding_bottom}px, "
                       f"标题={title_height}px, stats-container={stats_container_height}px, 总高度={card_height}px")
            logger.info(f"  各行最大高度: {row_max_heights}")

            # 添加stat-card背景
            bg_color_str = self.css_parser.get_background_color('.stat-card')
            if bg_color_str:
                from pptx.enum.shapes import MSO_SHAPE
                bg_shape = pptx_slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    UnitConverter.px_to_emu(80),
                    UnitConverter.px_to_emu(y_start),
                    UnitConverter.px_to_emu(1760),
                    UnitConverter.px_to_emu(card_height)
                )
                bg_shape.fill.solid()
                bg_rgb, alpha = ColorParser.parse_rgba(bg_color_str)
                if bg_rgb:
                    if alpha < 1.0:
                        bg_rgb = ColorParser.blend_with_white(bg_rgb, alpha)
                    bg_shape.fill.fore_color.rgb = bg_rgb
                bg_shape.line.fill.background()
                bg_shape.shadow.inherit = False
                logger.info(f"添加stat-card背景色: {bg_color_str}, 高度={card_height}px")

            y_start += 15  # 顶部padding

            # 添加标题(如果有)
            p_elem = card.find('p', class_='primary-color', recursive=False)
            if not p_elem:
                # 尝试在第一层查找
                for child in card.children:
                    if hasattr(child, 'get') and 'primary-color' in child.get('class', []):
                        p_elem = child
                        break

            if p_elem and p_elem.name == 'p':
                text = p_elem.get_text(strip=True)
                if text:
                    text_left = UnitConverter.px_to_emu(95)
                    text_top = UnitConverter.px_to_emu(y_start)
                    text_box = pptx_slide.shapes.add_textbox(
                        text_left, text_top,
                        UnitConverter.px_to_emu(1730), UnitConverter.px_to_emu(30)
                    )
                    text_frame = text_box.text_frame
                    text_frame.text = text
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            # 使用样式计算器获取正确的字体大小
                            p_font_size_pt = self.style_computer.get_font_size_pt(p_elem)
                            run.font.size = Pt(p_font_size_pt)
                            run.font.color.rgb = ColorParser.get_primary_color()
                            run.font.name = self.font_manager.get_font('body')

                    y_start += 35

            # 处理嵌套的stats-container
            next_y = self._convert_stats_container(stats_container, pptx_slide, y_start)
            return next_y + 20  # 底部padding

        # 2. 检查是否包含timeline (时间线类型)
        timeline = card.find('div', class_='timeline')
        if timeline:
            logger.info("stat-card包含timeline,处理时间线结构")

            # 计算stat-card总高度（用于背景）
            timeline_items = timeline.find_all('div', class_='timeline-item')
            num_items = len(timeline_items)
            # 每个timeline-item约85px, 加上标题35px和padding 30px
            card_height = num_items * 85 + 65

            # 添加stat-card背景
            shape_converter = ShapeConverter(pptx_slide, self.css_parser)
            # 从CSS获取背景颜色
            bg_color = self.css_parser.get_background_color('.stat-card')
            if bg_color:
                # 添加带颜色的背景矩形
                from pptx.enum.shapes import MSO_SHAPE
                bg_shape = pptx_slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    UnitConverter.px_to_emu(80),
                    UnitConverter.px_to_emu(y_start),
                    UnitConverter.px_to_emu(1760),
                    UnitConverter.px_to_emu(card_height)
                )
                bg_shape.fill.solid()
                # 解析背景颜色（支持rgba透明度）
                bg_rgb, alpha = ColorParser.parse_rgba(bg_color)
                if bg_rgb:
                    # 如果有透明度，与白色混合
                    if alpha < 1.0:
                        bg_rgb = ColorParser.blend_with_white(bg_rgb, alpha)
                        logger.info(f"混合透明度: alpha={alpha}, 结果=RGB({bg_rgb[0]}, {bg_rgb[1]}, {bg_rgb[2]})")
                    bg_shape.fill.fore_color.rgb = bg_rgb
                bg_shape.line.fill.background()  # 无边框
                # 移除阴影效果
                bg_shape.shadow.inherit = False
                logger.info(f"添加stat-card背景色: {bg_color}, 高度={card_height}px")

            y_start += 15  # 顶部padding

            # 添加标题(如果有)
            p_elem = card.find('p', class_='primary-color')
            if p_elem:
                text = p_elem.get_text(strip=True)
                if text:
                    text_left = UnitConverter.px_to_emu(95)  # 左侧padding
                    text_top = UnitConverter.px_to_emu(y_start)
                    text_box = pptx_slide.shapes.add_textbox(
                        text_left, text_top,
                        UnitConverter.px_to_emu(1730), UnitConverter.px_to_emu(30)
                    )
                    text_frame = text_box.text_frame
                    text_frame.text = text
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            # 使用样式计算器获取正确的字体大小
                            title_font_size_pt = self.style_computer.get_font_size_pt(p_elem)
                            run.font.size = Pt(title_font_size_pt)
                            run.font.color.rgb = ColorParser.get_primary_color()
                            run.font.name = self.font_manager.get_font('body')

                    y_start += 35

            # 处理timeline
            timeline_converter = TimelineConverter(pptx_slide, self.css_parser)
            next_y = timeline_converter.convert_timeline(timeline, x=95, y=y_start, width=1730)

            return next_y + 35  # 时间线后留一些间距

        # 3. 检查是否包含canvas (图表类型)
        canvas = card.find('canvas')
        if canvas:
            logger.info("stat-card包含canvas,处理图表")

            # 从CSS读取约束
            constraints = self.css_parser.get_height_constraints('.stat-card')
            stat_card_padding_top = constraints.get('padding_top', 20)
            stat_card_padding_bottom = constraints.get('padding_bottom', 20)

            # 标题高度
            title_elem = card.find('p', class_='primary-color')
            has_title = title_elem is not None
            title_height = 0
            if has_title:
                title_font_size_pt = self.style_computer.get_font_size_pt(title_elem)
                title_height = int(title_font_size_pt * 1.5) + 12  # 字体高度 + margin-bottom

            # canvas高度 - 从CSS读取chart-container的约束，如果没有则使用默认值
            chart_constraints = self.css_parser.get_height_constraints('.chart-container')
            canvas_height = chart_constraints.get('min_height', 220)  # 默认220px
            if canvas_height < 180:
                canvas_height = 220  # 确保最小高度
            logger.debug(f"canvas高度: {canvas_height}px (从CSS chart-container约束)")

            # stat-card总高度
            card_height = stat_card_padding_top + title_height + canvas_height + stat_card_padding_bottom

            logger.info(f"stat-card(canvas)高度计算: padding={stat_card_padding_top+stat_card_padding_bottom}px, "
                       f"标题={title_height}px, canvas={canvas_height}px, 总高度={card_height}px")

            bg_color_str = self.css_parser.get_background_color('.stat-card')
            if bg_color_str:
                from pptx.enum.shapes import MSO_SHAPE
                bg_shape = pptx_slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    UnitConverter.px_to_emu(80),
                    UnitConverter.px_to_emu(y_start),
                    UnitConverter.px_to_emu(1760),
                    UnitConverter.px_to_emu(card_height)
                )
                bg_shape.fill.solid()
                bg_rgb, alpha = ColorParser.parse_rgba(bg_color_str)
                if bg_rgb:
                    if alpha < 1.0:
                        bg_rgb = ColorParser.blend_with_white(bg_rgb, alpha)
                    bg_shape.fill.fore_color.rgb = bg_rgb
                bg_shape.line.fill.background()
                bg_shape.shadow.inherit = False
                logger.info(f"添加stat-card背景色: {bg_color_str}")

            y_start += 15  # 顶部padding

            # 添加标题文本(如果有)
            p_elem = card.find('p', class_='primary-color')
            if p_elem:
                text = p_elem.get_text(strip=True)
                if text:
                    text_left = UnitConverter.px_to_emu(95)
                    text_top = UnitConverter.px_to_emu(y_start)
                    text_box = pptx_slide.shapes.add_textbox(
                        text_left, text_top,
                        UnitConverter.px_to_emu(1730), UnitConverter.px_to_emu(30)
                    )
                    text_frame = text_box.text_frame
                    text_frame.text = text
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            # 使用样式计算器获取正确的字体大小
                            title_font_size_pt = self.style_computer.get_font_size_pt(p_elem)
                            run.font.size = Pt(title_font_size_pt)
                            run.font.color.rgb = ColorParser.get_primary_color()
                            run.font.name = self.font_manager.get_font('body')

                    y_start += 35

            # 处理canvas图表
            chart_converter = ChartConverter(pptx_slide, self.css_parser, self.html_path)
            success = chart_converter.convert_chart(
                canvas,
                x=95,
                y=y_start,
                width=1730,
                height=220,
                use_screenshot=ChartCapture.is_available()
            )

            if not success:
                logger.warning("图表转换失败,已显示占位文本")

            return y_start + 240

        # 4. 检查是否包含新的HTML结构（h3 + p + p）
        h3_elem = card.find('h3')
        p_elements = card.find_all('p')
        h3_text = h3_elem.get_text(strip=True) if h3_elem else ""

        if h3_elem and len(p_elements) >= 2:
            logger.info(f"stat-card包含h3 + p + p结构，处理为数据卡片 (h3={h3_text})")
            return self._convert_modern_stat_card(card, pptx_slide, y_start)

        # 4.1 检查是否包含复杂结构（h3 + flex容器等）
        # 查找所有flex容器（不仅仅是直接子元素）
        flex_containers = card.find_all('div', class_='flex')
        logger.info(f"stat-card找到{len(flex_containers)}个flex容器")

        # 检查是否有flex容器包含risk-level标签（风险分布）
        risk_level_found = False
        risk_level_count = 0
        for flex_container in flex_containers:
            risk_levels = flex_container.find_all('span', class_='risk-level')
            risk_level_count += len(risk_levels)
            if risk_levels:
                risk_level_found = True
                logger.info(f"flex容器包含{len(risk_levels)}个risk-level标签")

        if risk_level_found:
            logger.info(f"stat-card包含风险等级标签（共{risk_level_count}个），使用增强处理 (h3={h3_text})")
            return self._convert_enhanced_stat_card(card, pptx_slide, y_start)

        # 5. 通用降级处理 - 提取所有文本内容
        logger.info("stat-card不包含已知结构,使用通用文本提取")
        return self._convert_generic_card(card, pptx_slide, y_start, card_type='stat-card')

    def _convert_card_with_bullet_points(self, card, pptx_slide, y_start: int, bullet_points, h3_elem=None) -> int:
        """
        转换包含bullet-point的卡片（适用于stat-card和data-card）

        Args:
            card: 卡片元素
            pptx_slide: PPTX幻灯片
            y_start: 起始Y坐标
            bullet_points: bullet-point元素列表
            h3_elem: h3标题元素（可选）

        Returns:
            下一个元素的Y坐标
        """
        logger.info(f"处理包含{len(bullet_points)}个bullet-point的卡片")
        x_base = 80

        # 添加背景色
        if 'stat-card' in card.get('class', []):
            bg_color_str = self.css_parser.get_background_color('.stat-card')
        else:
            bg_color_str = 'rgba(10, 66, 117, 0.03)'  # data-card默认背景色

        if bg_color_str:
            from pptx.enum.shapes import MSO_SHAPE
            # 计算高度：标题 + bullet-point列表
            estimated_height = 50  # 顶部padding
            if h3_elem:
                estimated_height += 50  # h3标题高度
            estimated_height += len(bullet_points) * 35 + 20  # bullet-point列表高度

            bg_shape = pptx_slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                UnitConverter.px_to_emu(x_base),
                UnitConverter.px_to_emu(y_start),
                UnitConverter.px_to_emu(1760),
                UnitConverter.px_to_emu(estimated_height)
            )
            bg_shape.fill.solid()
            bg_rgb, alpha = ColorParser.parse_rgba(bg_color_str)
            if bg_rgb:
                if alpha < 1.0:
                    bg_rgb = ColorParser.blend_with_white(bg_rgb, alpha)
                bg_shape.fill.fore_color.rgb = bg_rgb
            bg_shape.line.fill.background()
            logger.info(f"添加卡片背景色，高度={estimated_height}px")

        current_y = y_start + 20  # 顶部padding

        # 处理h3标题
        if h3_elem:
            h3_text = h3_elem.get_text(strip=True)
            if h3_text:
                # 获取h3的字体大小和颜色
                h3_font_size_pt = self.style_computer.get_font_size_pt(h3_elem)
                h3_color = self._get_element_color(h3_elem) or ColorParser.get_primary_color()

                text_left = UnitConverter.px_to_emu(x_base + 20)
                text_top = UnitConverter.px_to_emu(current_y)
                text_box = pptx_slide.shapes.add_textbox(
                    text_left, text_top,
                    UnitConverter.px_to_emu(1720), UnitConverter.px_to_emu(30)
                )
                text_frame = text_box.text_frame
                text_frame.text = h3_text
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        font_size_pt = self.style_computer.get_font_size_pt(h3_elem)
                        run.font.size = Pt(font_size_pt)
                        h3_color = self._get_element_color(h3_elem)
                        if h3_color:
                            run.font.color.rgb = h3_color
                        else:
                            run.font.color.rgb = ColorParser.get_primary_color()
                        run.font.name = self.font_manager.get_font('body')
                        # 检查是否加粗
                        h3_classes = h3_elem.get('class', [])
                        if 'font-bold' in h3_classes:
                            run.font.bold = True

                current_y += 50  # 标题后间距

        # 处理bullet-point列表
        self._process_bullet_points(bullet_points, card, pptx_slide, x_base, current_y, 1760, 0)

        # 计算总高度
        total_height = current_y - y_start + len(bullet_points) * 35 + 20

        # 如果是data-card，需要添加左边框
        if 'data-card' in card.get('class', []):
            from src.converters.shape_converter import ShapeConverter
            shape_converter = ShapeConverter(pptx_slide, self.css_parser)
            shape_converter.add_border_left(x_base, y_start, total_height, 4)

        return y_start + total_height + 10

    def _convert_grid_card_with_bullet_points(self, card, pptx_slide, shape_converter, x, y, width, bullet_points, h3_elem=None) -> int:
        """
        转换网格中包含bullet-point的卡片

        Args:
            card: 卡片元素
            pptx_slide: PPTX幻灯片
            shape_converter: 形状转换器
            x: X坐标
            y: Y坐标
            width: 宽度
            bullet_points: bullet-point元素列表
            h3_elem: h3标题元素（可选）

        Returns:
            下一个元素的Y坐标
        """
        logger.info(f"处理网格中包含{len(bullet_points)}个bullet-point的卡片")

        # 添加背景色
        if 'stat-card' in card.get('class', []):
            bg_color_str = self.css_parser.get_background_color('.stat-card')
        else:
            bg_color_str = 'rgba(10, 66, 117, 0.03)'  # data-card默认背景色

        if bg_color_str:
            from pptx.enum.shapes import MSO_SHAPE
            # 计算高度：标题 + bullet-point列表
            estimated_height = 50  # 顶部padding
            if h3_elem:
                estimated_height += 50  # h3标题高度
            estimated_height += len(bullet_points) * 35 + 20  # bullet-point列表高度

            bg_shape = pptx_slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                UnitConverter.px_to_emu(x),
                UnitConverter.px_to_emu(y),
                UnitConverter.px_to_emu(width),
                UnitConverter.px_to_emu(estimated_height)
            )
            bg_shape.fill.solid()
            bg_rgb, alpha = ColorParser.parse_rgba(bg_color_str)
            if bg_rgb:
                if alpha < 1.0:
                    bg_rgb = ColorParser.blend_with_white(bg_rgb, alpha)
                bg_shape.fill.fore_color.rgb = bg_rgb
            bg_shape.line.fill.background()
            logger.info(f"添加网格卡片背景色，高度={estimated_height}px")

        # 添加左边框（如果是stat-card）
        if 'stat-card' in card.get('class', []):
            border_left_style = self.css_parser.get_style('.stat-card').get('border-left', '')
            if '4px solid' in border_left_style:
                shape_converter.add_border_left(x, y, estimated_height, 4)

        current_y = y + 20  # 顶部padding

        # 处理h3标题
        if h3_elem:
            h3_text = h3_elem.get_text(strip=True)
            if h3_text:
                # 获取h3的字体大小和颜色
                h3_font_size_pt = self.style_computer.get_font_size_pt(h3_elem)
                h3_color = self._get_element_color(h3_elem) or ColorParser.get_primary_color()

                text_left = UnitConverter.px_to_emu(x + 20)
                text_top = UnitConverter.px_to_emu(current_y)
                text_box = pptx_slide.shapes.add_textbox(
                    text_left, text_top,
                    UnitConverter.px_to_emu(width - 40),
                    UnitConverter.px_to_emu(30)
                )
                text_frame = text_box.text_frame
                text_frame.text = h3_text
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        font_size_pt = self.style_computer.get_font_size_pt(h3_elem)
                        run.font.size = Pt(font_size_pt)
                        h3_color = self._get_element_color(h3_elem)
                        if h3_color:
                            run.font.color.rgb = h3_color
                        else:
                            run.font.color.rgb = ColorParser.get_primary_color()
                        run.font.name = self.font_manager.get_font('body')
                        # 检查是否加粗
                        h3_classes = h3_elem.get('class', [])
                        if 'font-bold' in h3_classes:
                            run.font.bold = True

                current_y += 50  # 标题后间距

        # 处理bullet-point列表
        self._process_bullet_points(bullet_points, card, pptx_slide, x, current_y, width, 0)

        # 计算总高度
        total_height = current_y - y + len(bullet_points) * 35 + 20

        return y + total_height + 10

    def _convert_modern_stat_card(self, card, pptx_slide, y_start: int) -> int:
        """
        转换现代样式的stat-card（h3 + p + p结构）

        Args:
            card: stat-card元素
            pptx_slide: PPTX幻灯片
            y_start: 起始Y坐标

        Returns:
            下一个元素的Y坐标
        """
        logger.info("处理现代样式stat-card")
        x_base = 80

        # 添加背景
        bg_color_str = self.css_parser.get_background_color('.stat-card')
        if bg_color_str:
            from pptx.enum.shapes import MSO_SHAPE
            bg_shape = pptx_slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                UnitConverter.px_to_emu(x_base),
                UnitConverter.px_to_emu(y_start),
                UnitConverter.px_to_emu(1760),
                UnitConverter.px_to_emu(200)  # 估算高度
            )
            bg_shape.fill.solid()
            bg_rgb, alpha = ColorParser.parse_rgba(bg_color_str)
            if bg_rgb:
                if alpha < 1.0:
                    bg_rgb = ColorParser.blend_with_white(bg_rgb, alpha)
                bg_shape.fill.fore_color.rgb = bg_rgb
            bg_shape.line.fill.background()
            bg_shape.shadow.inherit = False

        current_y = y_start + 20  # 顶部padding

        # 处理h3标题
        h3_elem = card.find('h3')
        if h3_elem:
            h3_text = h3_elem.get_text(strip=True)
            if h3_text:
                # 获取h3的字体大小和颜色
                h3_font_size_pt = self.style_computer.get_font_size_pt(h3_elem)
                h3_color = self._get_element_color(h3_elem) or ColorParser.get_primary_color()

                text_left = UnitConverter.px_to_emu(x_base + 20)
                text_top = UnitConverter.px_to_emu(current_y)
                text_box = pptx_slide.shapes.add_textbox(
                    text_left, text_top,
                    UnitConverter.px_to_emu(1720), UnitConverter.px_to_emu(30)
                )
                text_frame = text_box.text_frame
                text_frame.text = h3_text
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        # 使用样式计算器获取字体大小和颜色
                        font_size_pt = self.style_computer.get_font_size_pt(h3_elem)
                        run.font.size = Pt(font_size_pt)

                        # 获取h3元素的颜色
                        h3_color = self._get_element_color(h3_elem)
                        if h3_color:
                            run.font.color.rgb = h3_color
                        else:
                            # 默认使用主题色
                            run.font.color.rgb = ColorParser.get_primary_color()

                        run.font.name = self.font_manager.get_font('body')
                        # 检查是否加粗
                        h3_classes = h3_elem.get('class', [])
                        if 'font-bold' in h3_classes:
                            run.font.bold = True

                current_y += 40

        # 处理第一个p标签（主要数据）
        p_elements = card.find_all('p')
        if len(p_elements) >= 1:
            p1_text = p_elements[0].get_text(strip=True)
            if p1_text:
                # 获取p1的字体大小和颜色
                p1_font_size_pt = self.style_computer.get_font_size_pt(p_elements[0])
                p1_color = self._get_element_color(p_elements[0])

                text_left = UnitConverter.px_to_emu(x_base + 20)
                text_top = UnitConverter.px_to_emu(current_y)
                text_box = pptx_slide.shapes.add_textbox(
                    text_left, text_top,
                    UnitConverter.px_to_emu(1720), UnitConverter.px_to_emu(40)
                )
                text_frame = text_box.text_frame
                text_frame.text = p1_text
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(p1_font_size_pt)
                        # 应用颜色
                        if p1_color:
                            run.font.color.rgb = p1_color
                        else:
                            # 如果没有特定颜色，使用默认文本颜色
                            run.font.color.rgb = ColorParser.get_text_color()
                        # 检查是否加粗
                        p1_classes = p_elements[0].get('class', [])
                        if 'font-bold' in p1_classes:
                            run.font.bold = True
                        run.font.name = self.font_manager.get_font('body')

                current_y += 50

        # 处理第二个p标签（描述）
        if len(p_elements) >= 2:
            p2_text = p_elements[1].get_text(strip=True)
            if p2_text:
                # 获取p2的字体大小和颜色
                p2_font_size_pt = self.style_computer.get_font_size_pt(p_elements[1])
                p2_color = self._get_element_color(p_elements[1])

                text_left = UnitConverter.px_to_emu(x_base + 20)
                text_top = UnitConverter.px_to_emu(current_y)
                text_box = pptx_slide.shapes.add_textbox(
                    text_left, text_top,
                    UnitConverter.px_to_emu(1720), UnitConverter.px_to_emu(30)
                )
                text_frame = text_box.text_frame
                text_frame.text = p2_text
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(p2_font_size_pt)
                        # 应用颜色
                        if p2_color:
                            run.font.color.rgb = p2_color
                        else:
                            # 如果没有特定颜色，使用默认文本颜色
                            run.font.color.rgb = ColorParser.get_text_color()
                        run.font.name = self.font_manager.get_font('body')

                current_y += 40

        return y_start + 200  # 返回估算的卡片高度

    def _convert_enhanced_stat_card(self, card, pptx_slide, y_start: int) -> int:
        """
        转换增强样式stat-card（支持复杂内容结构，如flex布局、风险等级标签等）

        Args:
            card: stat-card元素
            pptx_slide: PPTX幻灯片
            y_start: 起始Y坐标

        Returns:
            下一个元素的Y坐标
        """
        logger.info("处理增强样式stat-card")
        x_base = 80

        # 添加背景
        bg_color_str = self.css_parser.get_background_color('.stat-card')
        if bg_color_str:
            from pptx.enum.shapes import MSO_SHAPE
            bg_shape = pptx_slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                UnitConverter.px_to_emu(x_base),
                UnitConverter.px_to_emu(y_start),
                UnitConverter.px_to_emu(1760),
                UnitConverter.px_to_emu(180)  # 估算高度
            )
            bg_shape.fill.solid()
            bg_rgb, alpha = ColorParser.parse_rgba(bg_color_str)
            if bg_rgb:
                if alpha < 1.0:
                    bg_rgb = ColorParser.blend_with_white(bg_rgb, alpha)
                bg_shape.fill.fore_color.rgb = bg_rgb
            bg_shape.line.fill.background()
            bg_shape.shadow.inherit = False

        # 添加左边框
        border_left_style = self.css_parser.get_style('.stat-card').get('border-left', '')
        if '4px solid' in border_left_style:
            shape_converter = ShapeConverter(pptx_slide, self.css_parser)
            shape_converter.add_border_left(x_base, y_start, 180, 4)

        current_y = y_start + 20  # 顶部padding

        # 处理h3标题
        h3_elem = card.find('h3')
        if h3_elem:
            h3_text = h3_elem.get_text(strip=True)
            if h3_text:
                h3_font_size_pt = self.style_computer.get_font_size_pt(h3_elem)
                h3_color = self._get_element_color(h3_elem) or ColorParser.get_primary_color()

                text_left = UnitConverter.px_to_emu(x_base + 20)
                text_top = UnitConverter.px_to_emu(current_y)
                text_box = pptx_slide.shapes.add_textbox(
                    text_left, text_top,
                    UnitConverter.px_to_emu(1720), UnitConverter.px_to_emu(30)
                )
                text_frame = text_box.text_frame
                text_frame.text = h3_text
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        font_size_pt = self.style_computer.get_font_size_pt(h3_elem)
                        run.font.size = Pt(font_size_pt)
                        h3_color = self._get_element_color(h3_elem)
                        if h3_color:
                            run.font.color.rgb = h3_color
                        else:
                            run.font.color.rgb = ColorParser.get_primary_color()
                        run.font.name = self.font_manager.get_font('body')
                        # 检查是否加粗
                        h3_classes = h3_elem.get('class', [])
                        if 'font-bold' in h3_classes or 'text-2xl' in h3_classes:
                            run.font.bold = True

                current_y += 35

        # 处理复杂内容（查找所有flex容器内的内容）
        flex_containers = card.find_all('div', class_='flex')
        risk_levels = []
        for flex_container in flex_containers:
            # 查找所有包含risk-level的span
            risk_levels.extend(flex_container.find_all('span', class_='risk-level'))

        if risk_levels:

            # 计算每个风险标签的宽度和间距
            num_risks = len(risk_levels)
            if num_risks > 0:
                # 水平排列风险等级标签
                total_width = 1720  # 可用宽度
                risk_width = total_width // num_risks - 20  # 每个标签宽度，留出间距
                current_x = x_base + 20

                for risk_level in risk_levels:
                    risk_text = risk_level.get_text(strip=True)
                    risk_classes = risk_level.get('class', [])

                    # 获取风险等级的颜色
                    risk_color = None
                    bg_color = None
                    if 'risk-high' in risk_classes:
                        risk_color = ColorParser.parse_color('#dc2626')  # 红色
                        bg_color = RGBColor(252, 231, 229)  # 浅红色背景
                    elif 'risk-medium' in risk_classes:
                        risk_color = ColorParser.parse_color('#f59e0b')  # 橙色
                        bg_color = RGBColor(254, 243, 199)  # 浅橙色背景
                    elif 'risk-low' in risk_classes:
                        risk_color = ColorParser.parse_color('#3b82f6')  # 蓝色
                        bg_color = RGBColor(239, 246, 255)  # 浅蓝色背景

                    # 添加背景形状
                    if bg_color:
                        from pptx.enum.shapes import MSO_SHAPE
                        bg_shape = pptx_slide.shapes.add_shape(
                            MSO_SHAPE.ROUNDED_RECTANGLE,
                            UnitConverter.px_to_emu(current_x),
                            UnitConverter.px_to_emu(current_y),
                            UnitConverter.px_to_emu(risk_width),
                            UnitConverter.px_to_emu(35)
                        )
                        bg_shape.fill.solid()
                        bg_shape.fill.fore_color.rgb = bg_color
                        bg_shape.line.fill.background()

                    # 创建风险等级文本框
                    text_left = UnitConverter.px_to_emu(current_x + 5)
                    text_top = UnitConverter.px_to_emu(current_y + 5)
                    text_box = pptx_slide.shapes.add_textbox(
                        text_left, text_top,
                        UnitConverter.px_to_emu(risk_width - 10), UnitConverter.px_to_emu(25)
                    )
                    text_frame = text_box.text_frame
                    text_frame.text = risk_text
                    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

                    for paragraph in text_frame.paragraphs:
                        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                        for run in paragraph.runs:
                            run.font.size = Pt(20)
                            run.font.bold = True
                            run.font.name = self.font_manager.get_font('body')

                            # 应用风险等级颜色
                            if risk_color:
                                run.font.color.rgb = risk_color

                    # 移动到下一个位置
                    current_x += risk_width + 20

                current_y += 45
        else:
            # 处理其他p标签
            p_elements = card.find_all('p')
            for p_elem in p_elements:
                p_text = p_elem.get_text(strip=True)
                if p_text:
                    p_font_size_pt = self.style_computer.get_font_size_pt(p_elem)
                    p_color = self._get_element_color(p_elem)

                    text_left = UnitConverter.px_to_emu(x_base + 20)
                    text_top = UnitConverter.px_to_emu(current_y)
                    text_box = pptx_slide.shapes.add_textbox(
                        text_left, text_top,
                        UnitConverter.px_to_emu(1720), UnitConverter.px_to_emu(25)
                    )
                    text_frame = text_box.text_frame
                    text_frame.text = p_text
                    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(p_font_size_pt)
                            if p_color:
                                run.font.color.rgb = p_color
                            run.font.name = self.font_manager.get_font('body')

                    current_y += 30

        return y_start + 180  # 返回估算的卡片高度

    def _convert_generic_card(self, card, pptx_slide, y_start: int, card_type: str = 'card') -> int:
        """
        通用卡片内容转换 - 降级处理未知结构

        提取所有文本内容，按段落渲染，保持基本样式

        Args:
            card: 卡片元素
            pptx_slide: PPTX幻灯片
            y_start: 起始Y坐标
            card_type: 卡片类型（用于样式区分）

        Returns:
            下一个元素的Y坐标
        """
        logger.info(f"使用通用渲染器处理{card_type}")

        x_base = 80
        current_y = y_start

        # 提取所有段落元素 (p, h1, h2, h3, div等)
        text_elements = []

        # 查找所有文本容器
        for elem in card.descendants:
            if hasattr(elem, 'name') and elem.name in ['p', 'h1', 'h2', 'h3', 'h4', 'div', 'span']:
                # 只提取没有子块级元素的文本节点
                if not elem.find_all(['div', 'p', 'h1', 'h2', 'h3']):
                    text = elem.get_text(strip=True)
                    if text and len(text) > 2:  # 过滤空文本和单字符
                        # 检查是否有特殊样式
                        classes = elem.get('class', [])
                        is_primary = 'primary-color' in classes or 'stat-value' in classes
                        is_bold = 'font-bold' in classes or 'font-semibold' in classes or 'stat-value' in classes or elem.name in ['h1', 'h2', 'h3', 'h4']
                        is_stat_value = 'stat-value' in classes
                        is_stat_label = 'stat-label' in classes
                        # 检查是否有其他颜色类
                        has_color_class = any(cls.startswith('text-') for cls in classes)

                        text_elements.append({
                            'text': text,
                            'tag': elem.name,
                            'is_primary': is_primary,
                            'is_bold': is_bold,
                            'is_stat_value': is_stat_value,
                            'is_stat_label': is_stat_label,
                            'has_color_class': has_color_class,
                            'element': elem  # 保存元素引用以获取颜色
                        })

        # 去重（避免嵌套元素重复提取）
        seen_texts = set()
        unique_elements = []
        for elem in text_elements:
            if elem['text'] not in seen_texts:
                seen_texts.add(elem['text'])
                unique_elements.append(elem)

        logger.info(f"提取了 {len(unique_elements)} 个文本段落")

        # 暂时记录背景和边框配置，等实际内容渲染完成后再添加
        shape_converter = ShapeConverter(pptx_slide, self.css_parser)
        
        # 记录起始位置
        content_start_y = current_y
        
        # 根据卡片类型准备背景色
        bg_color = None
        needs_border = False
        
        if 'stat-card' in card_type:
            bg_color = self.css_parser.get_background_color('.stat-card')
            current_y += 15  # 顶部padding
        elif 'data-card' in card_type:
            needs_border = True
            current_y += 10
        elif 'stat-box' in card_type:
            bg_color = self.css_parser.get_background_color('.stat-box')
            current_y += 15
        elif 'strategy-card' in card_type:
            bg_color = self.css_parser.get_background_color('.strategy-card')
            needs_border = True
            current_y += 10

        # 渲染文本（unique_elements已在前面提取）
        for elem in unique_elements[:10]:  # 最多渲染10个段落，避免过长
            text = elem['text']
            is_primary = elem['is_primary']
            is_bold = elem['is_bold']
            is_stat_value = elem.get('is_stat_value', False)
            is_stat_label = elem.get('is_stat_label', False)
            tag = elem['tag']
            element = elem.get('element')  # 获取原始元素引用

            # 根据标签和类确定字体大小
            if is_stat_value:
                font_size = 42  # stat-value使用42px字体
            elif is_stat_label:
                font_size = 20  # stat-label使用20px字体
            elif tag in ['h1', 'h2']:
                font_size = 24
            elif tag == 'h3':
                font_size = 20
            elif is_primary:
                font_size = 20
            else:
                font_size = 16

            # 计算文本高度（粗略估算）
            lines = (len(text) // 80) + 1
            text_height = max(30, lines * 25)

            text_left = UnitConverter.px_to_emu(x_base + 20)
            text_top = UnitConverter.px_to_emu(current_y)
            text_box = pptx_slide.shapes.add_textbox(
                text_left, text_top,
                UnitConverter.px_to_emu(1720), UnitConverter.px_to_emu(text_height)
            )
            text_frame = text_box.text_frame
            text_frame.text = text
            text_frame.word_wrap = True

            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(font_size)
                    if is_bold:
                        run.font.bold = True
                    if is_primary:
                        run.font.color.rgb = ColorParser.get_primary_color()
                    elif is_stat_label:
                        # stat-label使用灰色
                        run.font.color.rgb = RGBColor(102, 102, 102)  # #666
                    elif element:
                        # 检查是否有其他颜色类
                        color = self._get_element_color(element)
                        if color:
                            run.font.color.rgb = color
                    run.font.name = self.font_manager.get_font('body')

            # 根据元素类型计算合适的间距
            if is_stat_value:
                # stat-value后的间距较小（CSS: margin-bottom: 8px）
                current_y += 58
            elif is_stat_label:
                # stat-label后的间距
                current_y += 36
            else:
                # 默认间距
                current_y += text_height + 10

        # 计算实际内容高度
        actual_height = current_y - content_start_y + 20
        
        # 现在根据实际高度添加背景色
        if bg_color:
            from pptx.enum.shapes import MSO_SHAPE
            bg_shape = pptx_slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                UnitConverter.px_to_emu(x_base),
                UnitConverter.px_to_emu(content_start_y),
                UnitConverter.px_to_emu(1760),
                UnitConverter.px_to_emu(actual_height)
            )
            bg_shape.fill.solid()
            bg_rgb, alpha = ColorParser.parse_rgba(bg_color)
            if bg_rgb:
                if alpha < 1.0:
                    bg_rgb = ColorParser.blend_with_white(bg_rgb, alpha)
                bg_shape.fill.fore_color.rgb = bg_rgb
            bg_shape.line.fill.background()
            bg_shape.shadow.inherit = False
            
            # 将背景移到最底层，避免覆盖文本
            sp_tree = pptx_slide.shapes._spTree
            bg_element = bg_shape._element
            sp_tree.remove(bg_element)
            sp_tree.insert(2, bg_element)
            
            logger.info(f"添加{card_type}背景色, 高度={actual_height}px (实际渲染，已移至底层)")
        
        # 添加边框（如果需要）
        if needs_border:
            shape_converter.add_border_left(x_base, content_start_y, actual_height, 4)

        return current_y + 20

    def _convert_strategy_card(self, card, pptx_slide, y_start: int) -> int:
        """
        转换策略卡片(.strategy-card)

        处理action-item结构：圆形数字图标 + 标题 + 描述
        """

        logger.info("处理strategy-card")
        x_base = 80

        action_items = card.find_all('div', class_='action-item')

        # 使用ContentHeightCalculator动态计算card高度
        card_height = self.height_calculator.calculate_strategy_card_height(card)
        
        # 从CSS读取padding用于布局
        constraints = self.css_parser.get_height_constraints('.strategy-card')
        strategy_card_padding = constraints.get('padding_top', 10)

        bg_color_str = self.css_parser.get_background_color('.strategy-card')
        if bg_color_str:
            from pptx.enum.shapes import MSO_SHAPE
            bg_shape = pptx_slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                UnitConverter.px_to_emu(x_base),
                UnitConverter.px_to_emu(y_start),
                UnitConverter.px_to_emu(1760),
                UnitConverter.px_to_emu(card_height)
            )
            bg_shape.fill.solid()
            bg_rgb, alpha = ColorParser.parse_rgba(bg_color_str)
            if bg_rgb:
                if alpha < 1.0:
                    bg_rgb = ColorParser.blend_with_white(bg_rgb, alpha)
                bg_shape.fill.fore_color.rgb = bg_rgb
            bg_shape.line.fill.background()
            bg_shape.shadow.inherit = False

        # 添加左边框
        shape_converter = ShapeConverter(pptx_slide, self.css_parser)
        shape_converter.add_border_left(x_base, y_start, card_height, 4)

        current_y = y_start + 15

        # 添加标题
        p_elem = card.find('p', class_='primary-color')
        if p_elem:
            text = p_elem.get_text(strip=True)
            if text:
                text_left = UnitConverter.px_to_emu(x_base + 20)
                text_top = UnitConverter.px_to_emu(current_y)
                text_box = pptx_slide.shapes.add_textbox(
                    text_left, text_top,
                    UnitConverter.px_to_emu(1720), UnitConverter.px_to_emu(30)
                )
                text_frame = text_box.text_frame
                text_frame.text = text
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        # 使用样式计算器获取正确的字体大小
                        title_font_size_pt = self.style_computer.get_font_size_pt(p_elem)
                        run.font.size = Pt(title_font_size_pt)
                        run.font.color.rgb = ColorParser.get_primary_color()
                        run.font.name = self.font_manager.get_font('body')

                current_y += 40

        # 处理每个action-item
        for item in action_items:
            # 获取数字
            number_elem = item.find('div', class_='action-number')
            number_text = number_elem.get_text(strip=True) if number_elem else "•"

            # 获取action-content
            content_elem = item.find('div', class_='action-content')
            if not content_elem:
                continue

            # 获取标题和描述
            title_elem = content_elem.find('div', class_='action-title')
            title_text = title_elem.get_text(strip=True) if title_elem else ""

            desc_elem = content_elem.find('p')
            desc_text = desc_elem.get_text(strip=True) if desc_elem else ""

            # 渲染圆形数字图标
            from pptx.enum.shapes import MSO_SHAPE
            from pptx.enum.text import MSO_ANCHOR

            circle_size = 28
            circle_left = UnitConverter.px_to_emu(x_base + 20)
            circle_top = UnitConverter.px_to_emu(current_y)
            circle = pptx_slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                circle_left,
                circle_top,
                UnitConverter.px_to_emu(circle_size),
                UnitConverter.px_to_emu(circle_size)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = ColorParser.get_primary_color()
            circle.line.fill.background()

            # 在圆形内添加数字文本
            circle_text_frame = circle.text_frame
            circle_text_frame.text = number_text
            circle_text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in circle_text_frame.paragraphs:
                paragraph.alignment = 2  # PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(14)
                    run.font.color.rgb = ColorParser.parse_color('#FFFFFF')
                    run.font.name = self.font_manager.get_font('body')
                    run.font.bold = True

            # 渲染标题（右侧）
            title_left = UnitConverter.px_to_emu(x_base + 60)
            title_top = UnitConverter.px_to_emu(current_y)
            if title_text:
                title_box = pptx_slide.shapes.add_textbox(
                    title_left, title_top,
                    UnitConverter.px_to_emu(1680), UnitConverter.px_to_emu(25)
                )
                title_frame = title_box.text_frame
                title_frame.text = title_text
                title_frame.word_wrap = True
                for paragraph in title_frame.paragraphs:
                    for run in paragraph.runs:
                        # 使用样式计算器获取正确的字体大小
                        title_font_size_pt = self.style_computer.get_font_size_pt(title_elem)
                        run.font.size = Pt(title_font_size_pt)
                        run.font.bold = True
                        run.font.color.rgb = ColorParser.get_primary_color()
                        run.font.name = self.font_manager.get_font('body')

                current_y += 28

            # 渲染描述（缩进）
            if desc_text:
                desc_left = UnitConverter.px_to_emu(x_base + 60)
                desc_top = UnitConverter.px_to_emu(current_y)
                desc_box = pptx_slide.shapes.add_textbox(
                    desc_left, desc_top,
                    UnitConverter.px_to_emu(1680), UnitConverter.px_to_emu(40)
                )
                desc_frame = desc_box.text_frame
                desc_frame.text = desc_text
                desc_frame.word_wrap = True
                for paragraph in desc_frame.paragraphs:
                    for run in paragraph.runs:
                        # 使用样式计算器获取正确的字体大小
                        desc_font_size_pt = self.style_computer.get_font_size_pt(desc_elem)
                        run.font.size = Pt(desc_font_size_pt)
                        run.font.name = self.font_manager.get_font('body')

                current_y += 50
            else:
                current_y += 35

        return current_y + 20

    def _convert_risk_card(self, card, pptx_slide, shape_converter, y_start: int) -> int:
        """
        转换风险卡片 (risk-card)

        Args:
            card: risk-card元素
            pptx_slide: PPTX幻灯片
            shape_converter: 形状转换器
            y_start: 起始Y坐标

        Returns:
            下一个元素的Y坐标
        """
        logger.info("处理risk-card风险卡片")

        x_base = 80
        card_width = 1760
        
        # 使用ContentHeightCalculator动态计算risk-card高度
        card_height = self.height_calculator.calculate_data_card_height(card, card_width)

        # 获取CSS样式
        card_style = self.css_parser.get_class_style('risk-card') or {}

        # 添加背景（渐变效果）
        bg_color_str = card_style.get('background', 'linear-gradient(135deg, rgba(239, 68, 68, 0.08) 0%, rgba(239, 68, 68, 0.02) 100%)')

        # 创建矩形背景
        bg_shape = pptx_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            UnitConverter.px_to_emu(x_base),
            UnitConverter.px_to_emu(y_start),
            UnitConverter.px_to_emu(card_width),
            UnitConverter.px_to_emu(card_height)
        )
        bg_shape.fill.solid()

        # 解析渐变背景色，使用最深的颜色
        if 'rgba(239, 68, 68' in bg_color_str:
            # 红色系风险
            bg_rgb = RGBColor(254, 242, 242)  # 非常浅的红色
        elif 'rgba(251, 146' in bg_color_str:
            # 橙色系风险
            bg_rgb = RGBColor(255, 251, 235)  # 非常浅的橙色
        elif 'rgba(250, 204' in bg_color_str:
            # 黄色系风险
            bg_rgb = RGBColor(254, 252, 232)  # 非常浅的黄色
        else:
            # 默认浅灰色
            bg_rgb = RGBColor(249, 250, 251)

        bg_shape.fill.fore_color.rgb = bg_rgb
        bg_shape.line.fill.background()
        bg_shape.shadow.inherit = False

        # 添加左边框
        border_color_str = card_style.get('border-left-color', '#ef4444')
        border_color = ColorParser.parse_color(border_color_str)
        if not border_color:
            # 根据风险等级确定边框颜色
            border_color = ColorParser.parse_color('#ef4444')  # 默认红色

        border_shape = pptx_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            UnitConverter.px_to_emu(x_base),
            UnitConverter.px_to_emu(y_start),
            UnitConverter.px_to_emu(4),
            UnitConverter.px_to_emu(card_height)
        )
        border_shape.fill.solid()
        border_shape.fill.fore_color.rgb = border_color
        border_shape.line.fill.background()

        current_y = y_start + 20

        # 处理flex布局内容
        flex_container = card.find('div', class_='flex')
        if flex_container:
            # 获取左侧内容区域
            left_div = flex_container.find('div', class_='flex-1')
            if left_div:
                # 处理风险标题
                title_div = left_div.find('div', class_='risk-title')
                if title_div:
                    # 获取图标
                    icon_elem = title_div.find('i')
                    icon_text = ""
                    icon_color = None

                    if icon_elem:
                        icon_classes = icon_elem.get('class', [])
                        # 根据图标类确定颜色
                        if 'severity-critical' in icon_classes:
                            icon_color = ColorParser.parse_color('#dc2626')
                            icon_text = "⚠"
                        elif 'severity-high' in icon_classes:
                            icon_color = ColorParser.parse_color('#ea580c')
                            icon_text = "⚠"
                        elif 'severity-medium' in icon_classes:
                            icon_color = ColorParser.parse_color('#d97706')
                            icon_text = "⚠"
                        else:
                            icon_text = "•"

                    # 获取标题文本
                    title_text = title_div.get_text(strip=True)
                    if icon_text:
                        title_text = title_text.replace(icon_text, "").strip()

                    # 添加标题文本
                    text_left = UnitConverter.px_to_emu(x_base + 20)
                    text_top = UnitConverter.px_to_emu(current_y)

                    if icon_text and icon_color:
                        # 如果有图标，创建两段式文本
                        text_box = pptx_slide.shapes.add_textbox(
                            text_left, text_top,
                            UnitConverter.px_to_emu(1200), UnitConverter.px_to_emu(35)
                        )
                        text_frame = text_box.text_frame
                        p = text_frame.paragraphs[0]

                        # 图标run
                        icon_run = p.add_run()
                        icon_run.text = icon_text + " "
                        icon_run.font.size = Pt(26)
                        icon_run.font.name = self.font_manager.get_font('body')
                        icon_run.font.color.rgb = icon_color
                        icon_run.font.bold = True

                        # 标题run
                        title_run = p.add_run()
                        title_run.text = title_text
                        title_run.font.size = Pt(26)
                        title_run.font.name = self.font_manager.get_font('body')
                        title_run.font.bold = True
                        title_run.font.color.rgb = RGBColor(51, 51, 51)  # 深灰色
                    else:
                        # 没有图标，直接添加标题
                        text_box = pptx_slide.shapes.add_textbox(
                            text_left, text_top,
                            UnitConverter.px_to_emu(1200), UnitConverter.px_to_emu(35)
                        )
                        text_frame = text_box.text_frame
                        text_frame.text = title_text

                        for paragraph in text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.size = Pt(26)
                                run.font.name = self.font_manager.get_font('body')
                                run.font.bold = True
                                run.font.color.rgb = RGBColor(51, 51, 51)

                    current_y += 40

                # 处理风险描述
                desc_div = left_div.find('div', class_='risk-desc')
                if desc_div:
                    desc_text = desc_div.get_text(strip=True)
                    if desc_text:
                        text_box = pptx_slide.shapes.add_textbox(
                            UnitConverter.px_to_emu(x_base + 20),
                            UnitConverter.px_to_emu(current_y),
                            UnitConverter.px_to_emu(1000),
                            UnitConverter.px_to_emu(30)
                        )
                        text_frame = text_box.text_frame
                        text_frame.text = desc_text

                        for paragraph in text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.size = Pt(22)
                                run.font.name = self.font_manager.get_font('body')
                                run.font.color.rgb = RGBColor(102, 102, 102)  # 灰色

                        current_y += 35

                # 处理标签
                tag_div = left_div.find('div', class_='mt-3')
                if tag_div:
                    span_elem = tag_div.find('span')
                    if span_elem:
                        tag_text = span_elem.get_text(strip=True)
                        tag_classes = span_elem.get('class', [])

                        # 确定标签颜色
                        tag_bg_color = RGBColor(254, 226, 226)  # 浅红色背景
                        tag_text_color = RGBColor(153, 27, 27)  # 深红色文字

                        if 'bg-orange-100' in tag_classes:
                            tag_bg_color = RGBColor(255, 237, 213)  # 浅橙色背景
                            tag_text_color = RGBColor(154, 52, 18)  # 深橙色文字
                        elif 'bg-yellow-100' in tag_classes:
                            tag_bg_color = RGBColor(254, 249, 195)  # 浅黄色背景
                            tag_text_color = RGBColor(120, 53, 15)  # 深黄色文字

                        # 创建标签背景
                        tag_box = pptx_slide.shapes.add_shape(
                            MSO_SHAPE.RECTANGLE,
                            UnitConverter.px_to_emu(x_base + 20),
                            UnitConverter.px_to_emu(current_y),
                            UnitConverter.px_to_emu(120),
                            UnitConverter.px_to_emu(30)
                        )
                        tag_box.fill.solid()
                        tag_box.fill.fore_color.rgb = tag_bg_color
                        tag_box.line.fill.background()

                        # 添加标签文本
                        tag_text_box = pptx_slide.shapes.add_textbox(
                            UnitConverter.px_to_emu(x_base + 20),
                            UnitConverter.px_to_emu(current_y + 3),
                            UnitConverter.px_to_emu(120),
                            UnitConverter.px_to_emu(24)
                        )
                        tag_text_frame = tag_text_box.text_frame
                        tag_text_frame.text = tag_text

                        for paragraph in tag_text_frame.paragraphs:
                            paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                            for run in paragraph.runs:
                                run.font.size = Pt(14)
                                run.font.name = self.font_manager.get_font('body')
                                run.font.color.rgb = tag_text_color
                                run.font.bold = True

            # 获取右侧CVSS分数区域
            right_div = flex_container.find('div', class_='text-center')
            if right_div:
                # 获取CVSS分数
                score_div = right_div.find('div', class_='cvss-score')
                if score_div:
                    score_text = score_div.get_text(strip=True)

                    # 添加CVSS分数
                    score_box = pptx_slide.shapes.add_textbox(
                        UnitConverter.px_to_emu(x_base + 1400),
                        UnitConverter.px_to_emu(y_start + 50),
                        UnitConverter.px_to_emu(300),
                        UnitConverter.px_to_emu(60)
                    )
                    score_frame = score_box.text_frame
                    score_frame.text = score_text

                    for paragraph in score_frame.paragraphs:
                        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                        for run in paragraph.runs:
                            run.font.size = Pt(48)
                            run.font.name = self.font_manager.get_font('body')
                            run.font.bold = True

                            # 根据分数确定颜色
                            if '10.0' in score_text or '9.' in score_text:
                                run.font.color.rgb = RGBColor(239, 68, 68)  # 红色
                            elif '8.' in score_text or '7.' in score_text:
                                run.font.color.rgb = RGBColor(234, 88, 12)  # 橙色
                            elif '6.' in score_text or '5.' in score_text:
                                run.font.color.rgb = RGBColor(217, 119, 6)  # 黄色
                            else:
                                run.font.color.rgb = RGBColor(107, 114, 128)  # 灰色

                # 获取CVSS标签
                label_div = right_div.find('div', class_='cvss-label')
                if label_div:
                    label_text = label_div.get_text(strip=True)

                    label_box = pptx_slide.shapes.add_textbox(
                        UnitConverter.px_to_emu(x_base + 1400),
                        UnitConverter.px_to_emu(y_start + 110),
                        UnitConverter.px_to_emu(300),
                        UnitConverter.px_to_emu(30)
                    )
                    label_frame = label_box.text_frame
                    label_frame.text = label_text

                    for paragraph in label_frame.paragraphs:
                        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                        for run in paragraph.runs:
                            run.font.size = Pt(18)
                            run.font.name = self.font_manager.get_font('body')
                            run.font.color.rgb = RGBColor(102, 102, 102)

        return y_start + card_height + 20

    def _convert_toc_layout(self, card, toc_items, pptx_slide, y_start: int) -> int:
        """
        转换目录布局 (toc-item)

        处理左右两栏的目录布局，每项包含数字编号和文本

        Args:
            card: stat-card容器
            toc_items: 目录项列表
            pptx_slide: PPTX幻灯片
            y_start: 起始Y坐标

        Returns:
            下一个元素的Y坐标
        """
        logger.info("处理目录布局(toc-item)")
        x_base = 80

        # 检查grid布局列数（支持Tailwind CSS网格类）
        grid_container = card.find('div', class_='grid')
        if grid_container:
            grid_classes = grid_container.get('class', [])
            num_columns = 2  # 默认2列

            # 检查Tailwind CSS网格列类
            for cls in grid_classes:
                if cls.startswith('grid-cols-') and hasattr(self.css_parser, 'tailwind_grid_columns'):
                    columns = self.css_parser.tailwind_grid_columns.get(cls)
                    if columns:
                        num_columns = columns
                        logger.info(f"检测到Tailwind网格列类: {cls} -> {num_columns}列")
                        break
        else:
            num_columns = 2  # 默认2列

        logger.info(f"检测到目录布局，{num_columns}列，{len(toc_items)}个目录项")

        # 添加stat-card背景
        card_height = len(toc_items) // num_columns * 60 + 80  # 估算高度
        if len(toc_items) % num_columns > 0:
            card_height += 60

        bg_color_str = self.css_parser.get_background_color('.stat-card')
        if bg_color_str:
            from pptx.enum.shapes import MSO_SHAPE
            bg_shape = pptx_slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                UnitConverter.px_to_emu(x_base),
                UnitConverter.px_to_emu(y_start),
                UnitConverter.px_to_emu(1760),
                UnitConverter.px_to_emu(card_height)
            )
            bg_shape.fill.solid()
            bg_rgb, alpha = ColorParser.parse_rgba(bg_color_str)
            if bg_rgb:
                if alpha < 1.0:
                    bg_rgb = ColorParser.blend_with_white(bg_rgb, alpha)
                bg_shape.fill.fore_color.rgb = bg_rgb
            bg_shape.line.fill.background()
            bg_shape.shadow.inherit = False
            logger.info(f"添加目录卡片背景，高度={card_height}px")

        current_y = y_start + 20

        # 处理目录项
        for idx, toc_item in enumerate(toc_items):
            # 计算位置（网格布局）
            col = idx % num_columns
            row = idx // num_columns
            item_x = x_base + 20 + col * 880  # 每列宽度880px
            item_y = current_y + row * 60  # 每项高度60px

            # 提取数字和文本
            number_elem = toc_item.find('div', class_='toc-number')
            text_elem = toc_item.find('div', class_='toc-text')

            if number_elem and text_elem:
                number_text = number_elem.get_text(strip=True)
                text_content = text_elem.get_text(strip=True)

                # 获取字体大小
                number_font_size = self.style_computer.get_font_size_pt(number_elem)
                text_font_size = self.style_computer.get_font_size_pt(text_elem)

                # 添加数字
                number_left = UnitConverter.px_to_emu(item_x)
                number_top = UnitConverter.px_to_emu(item_y)
                number_box = pptx_slide.shapes.add_textbox(
                    number_left, number_top,
                    UnitConverter.px_to_emu(40), UnitConverter.px_to_emu(30)
                )
                number_frame = number_box.text_frame
                number_frame.text = number_text
                number_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                for paragraph in number_frame.paragraphs:
                    paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(number_font_size)
                        run.font.bold = True
                        run.font.color.rgb = ColorParser.get_primary_color()
                        run.font.name = self.font_manager.get_font('body')

                # 添加文本
                text_left = UnitConverter.px_to_emu(item_x + 50)
                text_top = UnitConverter.px_to_emu(item_y)
                text_box = pptx_slide.shapes.add_textbox(
                    text_left, text_top,
                    UnitConverter.px_to_emu(800), UnitConverter.px_to_emu(30)
                )
                text_frame = text_box.text_frame
                text_frame.text = text_content
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(text_font_size)
                        run.font.name = self.font_manager.get_font('body')

        return y_start + card_height + 10

    def _convert_grid_svg_chart(self, child, pptx_slide, shape_converter, x, y, width):
        """
        转换网格中包含SVG图表的子元素（如slide_003.html）

        Args:
            child: 子元素（包含h3和chart-container）
            pptx_slide: PPTX幻灯片
            shape_converter: 形状转换器
            x: X坐标
            y: Y坐标
            width: 宽度

        Returns:
            下一个元素的Y坐标
        """
        logger.info(f"处理网格中的SVG图表子元素，x={x}, y={y}, width={width}")

        current_y = y

        # 处理h3标题（如果有）
        h3_elem = child.find('h3')
        if h3_elem:
            h3_text = h3_elem.get_text(strip=True)
            if h3_text:
                from src.converters.text_converter import TextConverter
                from src.utils.style_computer import StyleComputer
                from src.utils.font_manager import FontManager

                # 创建临时text_converter用于处理h3标题
                temp_text_converter = TextConverter(pptx_slide, self.css_parser)

                # 转换h3标题（使用convert_paragraph方法）
                current_y = temp_text_converter.convert_paragraph(
                    h3_elem,
                    x,
                    current_y,
                    width
                )
                logger.info(f"添加h3标题: {h3_text}")
                current_y += 15  # 标题后的间距

        # 查找chart-container
        chart_container = child.find('div', class_='chart-container')
        if chart_container:
            # 查找SVG元素
            svg_elem = chart_container.find('svg')
            if svg_elem:
                logger.info("找到SVG元素，开始转换")

                # 初始化SVG转换器
                from src.converters.svg_converter import SvgConverter
                svg_converter = SvgConverter(pptx_slide, self.css_parser, self.html_path, self.use_stable_chart_capture)
                self.svg_converters.append(svg_converter)  # 记录实例

                # 计算SVG图表的位置和尺寸
                svg_x = x  # 使用网格项的x坐标
                svg_y = current_y
                svg_width = width  # 使用网格项的宽度

                # 转换SVG
                chart_height = svg_converter.convert_svg(
                    svg_elem,
                    chart_container,
                    svg_x,
                    svg_y,
                    svg_width,
                    0  # SVG在容器中的索引
                )

                if chart_height > 0:
                    current_y += chart_height
                    logger.info(f"SVG转换成功，高度={chart_height}px")
                else:
                    logger.error("SVG转换失败")
                    current_y += 200  # 默认高度
            else:
                logger.warning("chart-container中未找到SVG元素")
                current_y += 200
        else:
            logger.warning("未找到chart-container")
            current_y += 200

        # 查找并处理后续的data-card（如果有）
        data_cards = child.find_all('div', class_='data-card')
        for card in data_cards:
            # 使用现有的data-card处理逻辑
            card_y = self._convert_grid_data_card_at_position(
                card,
                pptx_slide,
                shape_converter,
                x,
                current_y,
                width
            )
            current_y = card_y + 10  # 卡片间距

        return current_y + 10  # 返回最终位置

    def _convert_grid_data_card_at_position(self, card, pptx_slide, shape_converter, x, y, width):
        """
        在指定位置转换网格中的data-card

        Args:
            card: data-card元素
            pptx_slide: PPTX幻灯片
            shape_converter: 形状转换器
            x: X坐标
            y: Y坐标
            width: 宽度

        Returns:
            下一个元素的Y坐标
        """
        logger.info(f"在指定位置处理data-card，x={x}, y={y}")

        # 使用现有的data-card处理逻辑，但在指定位置
        from src.converters.text_converter import TextConverter
        from src.utils.style_computer import StyleComputer
        from src.utils.font_manager import FontManager
        from pptx.enum.shapes import MSO_SHAPE

        # 创建临时text_converter
        temp_text_converter = TextConverter(pptx_slide, self.css_parser)

        current_y = y + 15  # 顶部padding

        # 处理data-card内容
        # 首先查找p标签
        p_elem = card.find('p')
        if p_elem:
            # 处理p标签内的内容
            current_y = temp_text_converter.convert_paragraph(
                p_elem,
                x + 20,  # 左边距（因为有左边框）
                current_y,
                width - 40
            )
        else:
            # 降级处理：查找bullet-point
            bullet_points = card.find_all('div', class_='bullet-point')
            if bullet_points:
                for bp in bullet_points:
                    text_elem = bp.find('p')
                    if text_elem:
                        current_y = temp_text_converter.convert_paragraph(
                            text_elem,
                            x + 20,
                            current_y,
                            width - 40
                        )
                        current_y += 8  # bullet-point间距

        # 添加左边框
        shape_converter.add_border_left(x, y, current_y - y + 15, 4)

        return current_y

    def _convert_bottom_info(self, bottom_container, pptx_slide, y_start: int) -> int:
        """
        转换底部信息布局

        处理包含bullet-point的flex容器中的底部信息
        HTML中是水平排列，PPTX中也应该水平排列

        Args:
            bottom_container: 底部信息容器
            pptx_slide: PPTX幻灯片
            y_start: 起始Y坐标

        Returns:
            下一个元素的Y坐标
        """
        logger.info("处理底部信息布局（水平排列）")
        x_base = 80
        current_y = y_start

        # 查找所有bullet-point
        bullet_points = bottom_container.find_all('div', class_='bullet-point')

        # 水平排列：计算每个bullet-point的宽度
        total_width = 1760  # 可用总宽度
        # 从CSS获取间距，bottom-info容器可能使用gap-10 (40px)
        gap = self.css_parser.get_gap_size('.gap-10')  # gap-10 = 40px
        # 检查容器inline style
        inline_style = bottom_container.get('style', '')
        if 'gap:' in inline_style:
            import re
            gap_match = re.search(r'gap:\s*(\d+)px', inline_style)
            if gap_match:
                gap = int(gap_match.group(1))
                logger.info(f"bottom-info从inline style检测到间距: {gap}px")
        else:
            logger.debug(f"bottom-info使用默认间距: {gap}px")
        item_width = total_width // len(bullet_points)  # 每项平均分配宽度

        for idx, bullet_point in enumerate(bullet_points):
            icon_elem = bullet_point.find('i')
            p_elem = bullet_point.find('p')

            if icon_elem and p_elem:
                # 获取图标字符和颜色
                icon_classes = icon_elem.get('class', [])
                icon_char = self._get_icon_char(icon_classes)

                # 根据图标类确定颜色
                icon_color = ColorParser.get_primary_color()
                if 'text-red-600' in icon_classes:
                    icon_color = RGBColor(220, 38, 38)  # 红色
                elif 'text-green-600' in icon_classes:
                    icon_color = RGBColor(34, 197, 94)  # 绿色
                elif 'text-blue-600' in icon_classes:
                    icon_color = RGBColor(59, 130, 246)  # 蓝色
                elif 'primary-color' in icon_classes:
                    icon_color = ColorParser.get_primary_color()

                # 检查是否包含priority-tag
                priority_tag = p_elem.find('span', class_='priority-tag')
                main_text = ""
                tag_text = ""
                tag_color = None

                if priority_tag:
                    # 提取标签文本
                    tag_text = priority_tag.get_text(strip=True)
                    tag_classes = priority_tag.get('class', [])

                    # 确定标签颜色
                    if 'priority-high' in tag_classes:
                        tag_color = RGBColor(239, 68, 68)  # 红色
                    elif 'priority-medium' in tag_classes:
                        tag_color = RGBColor(251, 146, 60)  # 橙色
                    elif 'priority-low' in tag_classes:
                        tag_color = RGBColor(209, 177, 0)  # 黄色

                    # 移除标签后获取主文本
                    priority_tag.extract()
                    main_text = p_elem.get_text(strip=True)
                else:
                    main_text = p_elem.get_text(strip=True)

                # 计算水平位置
                item_x = x_base + idx * (item_width + gap)

                # 添加图标（在文本左侧）
                icon_left = UnitConverter.px_to_emu(item_x)
                icon_top = UnitConverter.px_to_emu(current_y)
                icon_box = pptx_slide.shapes.add_textbox(
                    icon_left, icon_top,
                    UnitConverter.px_to_emu(30), UnitConverter.px_to_emu(25)
                )
                icon_frame = icon_box.text_frame
                icon_frame.text = icon_char
                icon_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                for paragraph in icon_frame.paragraphs:
                    paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                    for run in paragraph.runs:
                        icon_font_size_pt = self._get_font_size_pt(p_elem, default_px=20)
                        run.font.size = Pt(icon_font_size_pt)
                        run.font.color.rgb = icon_color
                        run.font.name = self.font_manager.get_font('body')

                # 添加文本（在图标右侧）
                text_left = UnitConverter.px_to_emu(item_x + 40)
                text_top = UnitConverter.px_to_emu(current_y)
                text_width = item_width - 40  # 减去图标占用的宽度

                # 获取字体大小
                font_size_pt = self._get_font_size_pt(p_elem, default_px=20)

                # 检查文本中是否有strong标签
                strong_elem = p_elem.find('strong')

                if strong_elem or tag_text:
                    # 创建文本框
                    text_box = pptx_slide.shapes.add_textbox(
                        text_left, text_top,
                        UnitConverter.px_to_emu(text_width), UnitConverter.px_to_emu(25)
                    )
                    text_frame = text_box.text_frame
                    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                    text_frame.clear()
                    p = text_frame.paragraphs[0]

                    # 处理带strong的文本
                    if strong_elem:
                        strong_text = strong_elem.get_text(strip=True)
                        # 移除strong标签获取剩余文本
                        strong_copy = strong_elem.extract()
                        remaining_text = p_elem.get_text(strip=True)

                        # 添加加粗部分
                        if strong_text:
                            strong_run = p.add_run()
                            strong_run.text = strong_text
                            strong_run.font.size = Pt(font_size_pt)
                            strong_run.font.bold = True
                            strong_run.font.name = self.font_manager.get_font('body')

                        # 添加剩余部分
                        if remaining_text:
                            normal_run = p.add_run()
                            normal_run.text = remaining_text
                            normal_run.font.size = Pt(font_size_pt)
                            normal_run.font.name = self.font_manager.get_font('body')

                        # 添加标签文本（如果存在）
                        if tag_text and tag_color:
                            tag_run = p.add_run()
                            tag_run.text = " " + tag_text
                            tag_font_size_pt = max(10, font_size_pt - 2)
                            tag_run.font.size = Pt(tag_font_size_pt)
                            tag_run.font.color.rgb = tag_color
                            tag_run.font.bold = True
                            tag_run.font.name = self.font_manager.get_font('body')
                    else:
                        # 普通文本 + 标签
                        if main_text:
                            normal_run = p.add_run()
                            normal_run.text = main_text
                            normal_run.font.size = Pt(font_size_pt)
                            normal_run.font.name = self.font_manager.get_font('body')

                        # 添加标签文本（如果存在）
                        if tag_text and tag_color:
                            tag_run = p.add_run()
                            tag_run.text = " " + tag_text
                            tag_font_size_pt = max(10, font_size_pt - 2)
                            tag_run.font.size = Pt(tag_font_size_pt)
                            tag_run.font.color.rgb = tag_color
                            tag_run.font.bold = True
                            tag_run.font.name = self.font_manager.get_font('body')
                else:
                    # 普通文本
                    text_box = pptx_slide.shapes.add_textbox(
                        text_left, text_top,
                        UnitConverter.px_to_emu(text_width), UnitConverter.px_to_emu(50)  # 增加高度以支持换行
                    )
                    text_frame = text_box.text_frame
                    text_frame.text = main_text
                    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                    text_frame.word_wrap = True
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(font_size_pt)
                            run.font.name = self.font_manager.get_font('body')

        return current_y + 30

    def _convert_data_card(self, card, pptx_slide, shape_converter, y_start: int) -> int:
        """转换数据卡片(.data-card)"""

        # 防止重复处理：检查是否已经在其他容器中处理过
        # if hasattr(card, '_processed'):
        #     logger.info("data-card已处理过，跳过")
        #     return y_start
        # card._processed = True

        x_base = 80

        # 检查data-card内是否包含网格布局
        grid_container = card.find('div', class_='grid')
        if grid_container and grid_container.find_all('div', class_='bullet-point'):
            # 处理包含bullet-point的网格布局
            logger.info(f"data-card内发现grid和bullet-point，使用网格布局处理")
            return self._convert_data_card_grid_layout(card, grid_container, pptx_slide, shape_converter, y_start)
        else:
            logger.info(f"data-card使用标准处理流程（grid: {'是' if grid_container else '否'}, bullet-point: {'是' if grid_container and grid_container.find_all('div', class_='bullet-point') else '否'}）")

        # 智能判断data-card是否应该有背景色
        # 从当前HTML的CSS解析器获取实际的背景色定义
        bg_color_str = self.css_parser.get_background_color('.data-card')
        should_add_bg = False
        
        # 注意：背景色和左边框需要在计算完实际内容高度后统一添加
        # 暂时记录背景色配置，稍后根据实际内容高度添加
        should_add_bg = False
        if bg_color_str and bg_color_str != 'transparent' and bg_color_str != 'none':
            should_add_bg = True
            logger.info(f"data-card应该添加背景色: {bg_color_str}")
        else:
            logger.info(f"data-card没有定义背景色，只添加左边框")

        # 初始化当前Y坐标
        current_y = y_start + 10

        # 检查是否包含cve-card，如果有则跳过标题处理，让专门的CVE方法处理
        cve_cards = card.find_all('div', class_='cve-card')

        # 初始化标题变量（用于后面的检查）
        title_elem = None
        title_text = None

        if not cve_cards:
            # === 修复：简化的标题和内容处理逻辑 ===
            # 1. 首先查找并处理标题（查找h3标签或primary-color的p标签）
            title_elem = card.find('h3')

            # 如果没找到h3，再查找primary-color的p标签
            if not title_elem:
                title_elem = card.find('p', class_='primary-color')

            if title_elem:
                title_text = title_elem.get_text(strip=True)
                if title_text:
                    # 渲染标题
                    text_left = UnitConverter.px_to_emu(x_base + 20)
                    text_top = UnitConverter.px_to_emu(current_y)
                    text_box = pptx_slide.shapes.add_textbox(
                        text_left, text_top,
                        UnitConverter.px_to_emu(1720), UnitConverter.px_to_emu(30)
                    )
                    text_frame = text_box.text_frame
                    text_frame.text = title_text
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            # 使用样式计算器获取正确的字体大小
                            font_size_pt = self.style_computer.get_font_size_pt(title_elem)
                            run.font.size = Pt(font_size_pt)

                            # 设置颜色和字体
                            if title_elem.name == 'h3':
                                # h3标签使用主题色
                                run.font.color.rgb = ColorParser.get_primary_color()
                                run.font.bold = True
                            else:
                                # p标签根据class设置颜色
                                run.font.color.rgb = ColorParser.get_primary_color()

                            run.font.name = self.font_manager.get_font('body')

                    current_y += 40  # 标题后间距
                    logger.info(f"渲染data-card标题: {title_text}")

        # 2. 处理普通段落内容（明确排除标题元素、bullet-point内的元素和cve-card内的元素）
        content_paragraphs = []
        all_paragraphs = card.find_all('p')

        for p in all_paragraphs:
            # 新增：检查是否在bullet-point里
            parent = p.parent
            is_in_bullet_point = False
            is_in_cve_card = False
            while parent and parent != card:
                if parent.get('class'):
                    if 'bullet-point' in parent.get('class', []):
                        is_in_bullet_point = True
                        break
                    elif 'cve-card' in parent.get('class', []):
                        is_in_cve_card = True
                        break
                parent = parent.parent

            if is_in_bullet_point or is_in_cve_card:
                continue

            # 方法1：检查是否有primary-color类
            if 'primary-color' in p.get('class', []):
                continue

            # 方法2：如果是同一个元素对象，也跳过（防止类检查失败的情况）
            if title_elem and p is title_elem:
                continue

            # 方法3：如果文本内容完全相同，也跳过（最后保险）
            p_text = p.get_text(strip=True)
            if title_text and p_text == title_text:
                continue

            # 通过所有检查的才是真正的内容段落
            if p_text:
                content_paragraphs.append(p)

        logger.info(f"data-card段落过滤: 找到{len(all_paragraphs)}个p标签，排除标题后{len(content_paragraphs)}个普通段落")

        # 3. 渲染内容段落
        for p in content_paragraphs:
            text = p.get_text(strip=True)
            if text:
                text_left = UnitConverter.px_to_emu(x_base + 20)
                text_top = UnitConverter.px_to_emu(current_y)
                text_box = pptx_slide.shapes.add_textbox(
                    text_left, text_top,
                    UnitConverter.px_to_emu(1720), UnitConverter.px_to_emu(30)
                )
                text_frame = text_box.text_frame
                text_frame.text = text
                text_frame.word_wrap = True
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        # 使用样式计算器获取正确的字体大小
                        font_size_px = self.style_computer.get_font_size_pt(p)
                        run.font.size = Pt(font_size_px)
                        run.font.name = self.font_manager.get_font('body')

                current_y += 35  # 段落后间距
                # 过滤特殊字符，避免Windows控制台乱码
                clean_text = text[:30].replace('•', '*').replace('•', '*')
                logger.info(f"渲染data-card内容: {clean_text}...")  # 只记录前30个字符

        # 进度条
        progress_bars = card.find_all('div', class_='progress-container')
        progress_y = current_y
        for progress in progress_bars:
            label_div = progress.find('div', class_='progress-label')
            if label_div:
                spans = label_div.find_all('span')
                if len(spans) >= 2:
                    label_text = spans[0].get_text(strip=True)
                    percent_text = spans[1].get_text(strip=True)
                    percentage = UnitConverter.normalize_percentage(percent_text)

                    shape_converter.add_progress_bar(
                        label_text, percentage, x_base + 20, progress_y, 1700
                    )
                    progress_y += 60

        # 列表项 (bullet-point)
        bullet_points = card.find_all('div', class_='bullet-point')

        # 风险项目 (risk-item)
        risk_items = card.find_all('div', class_='risk-item')

        # === 修复：正确判断是否已有内容 ===
        # 不仅要检查progress-bar和bullet-point，还要检查是否已经处理了标题和段落
        has_title_or_content = title_elem is not None or len(content_paragraphs) > 0
        has_special_content = len(progress_bars) > 0 or len(bullet_points) > 0 or len(risk_items) > 0
        has_content = has_title_or_content or has_special_content

        logger.info(f"data-card内容检查: 标题={'是' if title_elem else '否'}, "
                   f"内容段落数={len(content_paragraphs)}, 进度条数={len(progress_bars)}, "
                   f"列表项数={len(bullet_points)}, 风险项数={len(risk_items)}, 总已有内容={'是' if has_content else '否'}")

        # 如果没有其他内容，从y_start开始处理bullet-point
        if not has_title_or_content and not progress_bars:
            progress_y = y_start + 10
        else:
            progress_y = current_y + 10 if current_y > y_start else current_y

        # 处理风险项目
        for risk_item in risk_items:
            icon_elem = risk_item.find('i')
            content_div = risk_item.find('div')

            if icon_elem and content_div:
                # 获取图标
                icon_classes = icon_elem.get('class', [])
                icon_char = self._get_icon_char(icon_classes)

                # 获取所有p标签
                p_tags = content_div.find_all('p')

                # 先处理第一个p标签的内容，计算实际需要的宽度
                total_text_width = 0
                has_inline_risk_level = False
                text_elements = []

                if len(p_tags) > 0:
                    first_p = p_tags[0]

                    # 智能识别内联元素位置的规则
                    # 如果span.risk-level紧跟在strong后面，应该在同一行显示
                    for elem in first_p.children:
                        if hasattr(elem, 'name'):
                            if elem.name == 'span' and 'risk-level' in elem.get('class', []):
                                # 检查前一个兄弟元素是否是strong
                                prev_sibling = elem.previous_sibling
                                if prev_sibling and hasattr(prev_sibling, 'name') and prev_sibling.name == 'strong':
                                    has_inline_risk_level = True
                                    break

                    # 计算图标宽度
                    icon_width = 0
                    if icon_char:
                        icon_width = self._calculate_text_width(icon_char + " ", Pt(22))

                    # 遍历first_p的所有直接子元素，计算文本宽度
                    current_x = icon_width
                    for elem in first_p.children:
                        if hasattr(elem, 'name'):
                            if elem.name == 'strong':
                                strong_text = elem.get_text(strip=True)
                                if strong_text:
                                    text_width = self._calculate_text_width(strong_text, Pt(22))
                                    text_elements.append({
                                        'type': 'strong',
                                        'text': strong_text,
                                        'font_size': Pt(22),
                                        'x_start': current_x,
                                        'x_end': current_x + text_width
                                    })
                                    current_x += text_width
                            elif elem.name == 'span' and 'risk-level' in elem.get('class', []):
                                risk_text = elem.get_text(strip=True)
                                risk_classes = elem.get('class', [])
                                # 在strong和risk-level之间添加空格
                                space_width = self._calculate_text_width(" ", Pt(22))
                                current_x += space_width
                                text_width = self._calculate_text_width(risk_text, Pt(20))
                                text_elements.append({
                                    'type': 'risk',
                                    'text': risk_text,
                                    'classes': risk_classes,
                                    'font_size': Pt(20),
                                    'x_start': current_x,
                                    'x_end': current_x + text_width,
                                    'inline': True  # 标记为内联元素
                                })
                                current_x += text_width
                        else:
                            # 处理文本节点
                            text_content = str(elem).strip()
                            if text_content:
                                text_width = self._calculate_text_width(text_content, Pt(22))
                                text_elements.append({
                                    'type': 'text',
                                    'text': text_content,
                                    'font_size': Pt(22),
                                    'x_start': current_x,
                                    'x_end': current_x + text_width
                                })
                                current_x += text_width

                    total_text_width = current_x

                # 计算文本框宽度（自适应文本长度）
                min_width = 400  # 最小宽度
                max_width = 1720  # 最大宽度
                box_width = max(min_width, min(total_text_width + 40, max_width))  # 加40px的padding

                # 计算所需的高度
                # 第一行：strong + risk-level（如果有内联）
                first_line_height = 35 if has_inline_risk_level else 30
                # 其他行：每个p标签占一行
                other_lines_height = (len(p_tags) - 1) * 28  # 每个额外的p标签28px
                total_height = first_line_height + other_lines_height + 20  # 20px padding

                # 创建自适应大小的文本框
                # 注意：risk_left 已经在前面设置为 UnitConverter.px_to_emu(x + 20)
                # 不要重复设置，保留正确的x坐标
                risk_top = UnitConverter.px_to_emu(progress_y)
                risk_box = pptx_slide.shapes.add_textbox(
                    risk_left, risk_top,
                    UnitConverter.px_to_emu(box_width), UnitConverter.px_to_emu(total_height)
                )
                risk_frame = risk_box.text_frame
                risk_frame.clear()
                risk_frame.margin_left = 0
                risk_frame.margin_right = 0
                risk_frame.margin_top = 0
                risk_frame.margin_bottom = 0
                p = risk_frame.paragraphs[0]

                # 添加图标
                if icon_char:
                    icon_run = p.add_run()
                    icon_run.text = icon_char + " "
                    icon_run.font.size = Pt(22)
                    icon_run.font.color.rgb = ColorParser.get_primary_color()
                    icon_run.font.name = self.font_manager.get_font('body')

                # 处理第一个p标签（可能包含strong和risk-level）
                if len(p_tags) > 0:
                    first_p = p_tags[0]

                    # 遍历first_p的所有直接子元素
                    for elem in first_p.children:
                        if hasattr(elem, 'name'):
                            if elem.name == 'strong':
                                strong_text = elem.get_text(strip=True)
                                if strong_text:
                                    strong_run = p.add_run()
                                    strong_run.text = strong_text
                                    strong_run.font.size = Pt(22)
                                    strong_run.font.bold = True
                                    strong_run.font.name = self.font_manager.get_font('body')
                                    strong_run.font.color.rgb = RGBColor(0, 0, 0)  # 黑色

                                    # 检查下一个元素是否是risk-level，如果是则添加空格
                                    next_sibling = elem.next_sibling
                                    if next_sibling and hasattr(next_sibling, 'name') and next_sibling.name == 'span' and 'risk-level' in next_sibling.get('class', []):
                                        strong_run.text += " "
                            elif elem.name == 'span' and 'risk-level' in elem.get('class', []):
                                risk_text = elem.get_text(strip=True)
                                risk_classes = elem.get('class', [])

                                # 获取风险等级的颜色和背景色
                                risk_color = None
                                bg_color = None
                                if 'risk-high' in risk_classes:
                                    risk_color = ColorParser.parse_color('#dc2626')  # 红色
                                    bg_color = RGBColor(252, 231, 229)  # 浅红色背景
                                elif 'risk-medium' in risk_classes:
                                    risk_color = ColorParser.parse_color('#f59e0b')  # 橙色
                                    bg_color = RGBColor(254, 243, 199)  # 浅橙色背景
                                elif 'risk-low' in risk_classes:
                                    risk_color = ColorParser.parse_color('#3b82f6')  # 蓝色
                                    bg_color = RGBColor(239, 246, 255)  # 浅蓝色背景
                                elif 'CVSS' in risk_text:
                                    # CVSS分数也使用特殊颜色
                                    if '10.0' in risk_text:
                                        risk_color = ColorParser.parse_color('#dc2626')  # 红色
                                        bg_color = RGBColor(252, 231, 229)  # 浅红色背景
                                    elif '9.8' in risk_text:
                                        risk_color = ColorParser.parse_color('#dc2626')  # 红色
                                        bg_color = RGBColor(252, 231, 229)  # 浅红色背景
                                    elif '8.6' in risk_text:
                                        risk_color = ColorParser.parse_color('#f59e0b')  # 橙色
                                        bg_color = RGBColor(254, 243, 199)  # 浅橙色背景

                                # 不添加到主文本框，而是创建独立的带背景文本框
                                # 这样"高危"会紧跟在strong后面，并有自己的背景
                                if elem_info:
                                    # 计算文本的绝对位置
                                    # elem_info['x_start'] 是相对于文本框的像素位置
                                    text_abs_left = risk_left + UnitConverter.px_to_emu(elem_info['x_start'])
                                    text_abs_top = risk_top + UnitConverter.px_to_emu(5)  # 微调垂直位置

                                    # 计算文本宽度
                                    text_width = elem_info['x_end'] - elem_info['x_start']
                                    bg_width = text_width + 16  # 左右各8px padding
                                    bg_width = max(bg_width, 50)  # 最小宽度50px

                                    # 先创建背景形状（如果有背景色）
                                    if bg_color:
                                        bg_shape = pptx_slide.shapes.add_shape(
                                            MSO_SHAPE.ROUNDED_RECTANGLE,
                                            text_abs_left,
                                            text_abs_top,
                                            UnitConverter.px_to_emu(bg_width),
                                            UnitConverter.px_to_emu(28)
                                        )
                                        bg_shape.fill.solid()
                                        bg_shape.fill.fore_color.rgb = bg_color
                                        bg_shape.line.fill.background()

                                    # 再创建文本框（覆盖在背景上）
                                    risk_text_box = pptx_slide.shapes.add_textbox(
                                        text_abs_left + UnitConverter.px_to_emu(8),  # 左内边距
                                        text_abs_top + UnitConverter.px_to_emu(4),   # 上内边距
                                        UnitConverter.px_to_emu(bg_width - 16),  # 减去padding
                                        UnitConverter.px_to_emu(20)  # 高度
                                    )
                                    risk_text_frame = risk_text_box.text_frame
                                    risk_text_frame.clear()
                                    risk_text_frame.text = risk_text
                                    risk_text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

                                    # 设置文本样式
                                    for paragraph in risk_text_frame.paragraphs:
                                        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                                        for run in paragraph.runs:
                                            run.font.size = Pt(20)
                                            run.font.bold = True
                                            run.font.name = self.font_manager.get_font('body')
                                            if risk_color:
                                                run.font.color.rgb = risk_color
                                            else:
                                                run.font.color.rgb = RGBColor(220, 38, 38)  # 默认红色

                                    logger.info(f"创建独立文本框: {risk_text}")
                                    logger.info(f"  绝对位置: ({UnitConverter.emu_to_px(text_abs_left)}, {UnitConverter.emu_to_px(text_abs_top)})")
                                    logger.info(f"  尺寸: {bg_width}px x 28px")

                                    # 将背景移到下层（这样不会覆盖文本）
                                    # 在python-pptx中，后添加的形状在上层
                                    # 所以我们需要把背景移到前面添加的元素后面
                                    try:
                                        # 获取所有形状，调整z-order
                                        shapes = pptx_slide.shapes
                                        bg_index = len(shapes) - 1  # 最后添加的背景
                                        # 找到文本框的索引
                                        text_index = None
                                        for i in range(len(shapes)):
                                            if shapes[i] == risk_box:
                                                text_index = i
                                                break
                                        # 如果找到了文本框，把背景移到它前面
                                        if text_index is not None and bg_index > text_index:
                                            # 需要重新创建顺序
                                            pass  # python-pptx不支持直接调整z-order
                                    except:
                                        pass

                # 处理第二个p标签（描述信息）
                if len(p_tags) > 1:
                    second_p = p_tags[1]
                    desc_text = second_p.get_text(strip=True)
                    if desc_text:
                        # 添加换行
                        p.add_run().text = "\n"

                        desc_run = p.add_run()
                        desc_run.text = desc_text
                        desc_run.font.size = Pt(18)
                        desc_run.font.name = self.font_manager.get_font('body')
                        desc_run.font.color.rgb = RGBColor(102, 102, 102)  # 灰色

                progress_y += total_height + 10  # 使用计算的高度+间距

        for bullet in bullet_points:
            # 检查是否有图标
            icon_elem = bullet.find('i')
            icon_char = None
            if icon_elem:
                icon_classes = icon_elem.get('class', [])
                icon_char = self._get_icon_char(icon_classes)

            # 检查是否有嵌套的div结构
            nested_div = bullet.find('div')
            if nested_div:
                # 处理嵌套结构: <div class="bullet-point"><i>...</i><div><p>...</p><p>...</p></div></div>
                all_p = nested_div.find_all('p')

                for idx, p in enumerate(all_p):
                    text = p.get_text(strip=True)
                    if not text:
                        continue

                    bullet_left = UnitConverter.px_to_emu(x_base + 20)
                    bullet_top = UnitConverter.px_to_emu(progress_y)

                    # 第一个p加图标,后续p缩进
                    if idx == 0:
                        prefix = f"{icon_char} " if icon_char else "• "
                        bullet_width = 1720
                    else:
                        prefix = "  "
                        bullet_width = 1720

                    bullet_box = pptx_slide.shapes.add_textbox(
                        bullet_left, bullet_top,
                        UnitConverter.px_to_emu(bullet_width), UnitConverter.px_to_emu(50)
                    )
                    bullet_frame = bullet_box.text_frame
                    bullet_frame.clear()
                    paragraph = bullet_frame.paragraphs[0]

                    # 处理冒号后的换行
                    if '：' in text or ':' in text:
                        # 分割文本为两部分
                        if '：' in text:
                            parts = text.split('：', 1)
                            separator = '：'
                        else:
                            parts = text.split(':', 1)
                            separator = ':'

                        if len(parts) == 2:
                            # 添加图标和第一部分（加粗）
                            run1 = paragraph.add_run()
                            run1.text = f"{prefix}{parts[0]}{separator}"
                            run1.font.bold = True
                            run1.font.size = Pt(25)
                            run1.font.name = self.font_manager.get_font('body')

                            # 获取字体大小
                            font_size_pt = self.style_computer.get_font_size_pt(p)
                            if font_size_pt:
                                run1.font.size = Pt(font_size_pt)

                            # 添加换行和第二部分
                            run2 = paragraph.add_run()
                            run2.text = "\n" + parts[1]
                            run2.font.size = Pt(25)
                            run2.font.name = self.font_manager.get_font('body')

                            if font_size_pt:
                                run2.font.size = Pt(font_size_pt)
                        else:
                            run = paragraph.add_run()
                            run.text = f"{prefix}{text}"
                            run.font.size = Pt(25)
                            run.font.name = self.font_manager.get_font('body')
                            # 获取实际字体大小
                            font_size_pt = self.style_computer.get_font_size_pt(p)
                            if font_size_pt:
                                run.font.size = Pt(font_size_pt)
                            # 第一个p加粗
                            if idx == 0:
                                run.font.bold = True
                    else:
                        run = paragraph.add_run()
                        run.text = f"{prefix}{text}"
                        run.font.size = Pt(25)
                        run.font.name = self.font_manager.get_font('body')
                        # 获取实际字体大小
                        font_size_pt = self.style_computer.get_font_size_pt(p)
                        if font_size_pt:
                            run.font.size = Pt(font_size_pt)
                        # 第一个p加粗
                        if idx == 0:
                            run.font.bold = True

                    bullet_frame.word_wrap = True

                    progress_y += 28 if idx == 0 else 50
            else:
                # 处理简单结构: <div class="bullet-point"><i>...</i><p>...</p></div>
                p = bullet.find('p')
                if p:
                    text = p.get_text(strip=True)
                    bullet_left = UnitConverter.px_to_emu(x_base + 20)
                    bullet_top = UnitConverter.px_to_emu(progress_y)
                    bullet_box = pptx_slide.shapes.add_textbox(
                        bullet_left, bullet_top,
                        UnitConverter.px_to_emu(1720), UnitConverter.px_to_emu(30)
                    )
                    # 使用图标或默认圆点
                    prefix = f"{icon_char} " if icon_char else "• "
                    bullet_frame = bullet_box.text_frame
                    bullet_frame.text = f"{prefix}{text}"
                    for paragraph in bullet_frame.paragraphs:
                        for run in paragraph.runs:
                            # 使用样式计算器获取正确的字体大小
                            font_size_px = self.style_computer.get_font_size_pt(p)
                            run.font.size = Pt(font_size_px)
                            run.font.name = self.font_manager.get_font('body')

                    progress_y += 35

        # 检查是否包含cve-card（使用前面已经检测的结果）
        if cve_cards:
            logger.info(f"检测到{len(cve_cards)}个cve-card，使用专门处理")
            return self._convert_cve_card_list(card, pptx_slide, shape_converter, y_start)

        # 如果没有识别到任何已知内容，使用通用降级处理
        if not has_content:
            logger.info("data-card不包含progress-bar或bullet-point,使用通用处理")
            return self._convert_generic_card(card, pptx_slide, y_start, card_type='data-card')

        # 计算实际渲染高度（基于实际内容的精确坐标）
        final_y = progress_y + 20
        actual_height = final_y - y_start

        # 现在根据实际内容高度添加背景色和左边框
        if should_add_bg:
            from pptx.enum.shapes import MSO_SHAPE
            bg_shape = pptx_slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                UnitConverter.px_to_emu(x_base),
                UnitConverter.px_to_emu(y_start),
                UnitConverter.px_to_emu(1760),
                UnitConverter.px_to_emu(actual_height)  # 使用实际渲染高度
            )
            bg_shape.fill.solid()
            bg_rgb, alpha = ColorParser.parse_rgba(bg_color_str)
            if bg_rgb:
                if alpha < 1.0:
                    bg_rgb = ColorParser.blend_with_white(bg_rgb, alpha)
                bg_shape.fill.fore_color.rgb = bg_rgb
            bg_shape.line.fill.background()
            
            # 将背景移到最底层，避免覆盖文本
            # 使用 shapes._spTree 访问实际的形状容器
            sp_tree = pptx_slide.shapes._spTree
            bg_element = bg_shape._element
            # 从当前位置移除
            sp_tree.remove(bg_element)
            # 插入到最前面（在非可视元素之后）
            # spTree 通常包含: nvGrpSpPr, grpSpPr, 然后是各个形状
            # 我们插入到索引2的位置（第一个实际形状位置）
            sp_tree.insert(2, bg_element)
            
            logger.info(f"添加data-card背景色: {bg_color_str}, 高度={actual_height}px (实际渲染，已移至底层)")

        # 添加左边框（使用实际计算的高度）
        shape_converter.add_border_left(x_base, y_start, actual_height, 4)

        logger.info(f"data-card高度计算: 实际高度={actual_height}px, "
                   f"进度条数={len(progress_bars)}, 列表项数={len(bullet_points)}")

        return final_y

    def _convert_cve_card_list(self, card, pptx_slide, shape_converter, y_start: int) -> int:
        """
        转换包含CVE卡片列表的data-card

        Args:
            card: data-card元素
            pptx_slide: PPTX幻灯片
            shape_converter: 形状转换器
            y_start: 起始Y坐标

        Returns:
            下一个元素的Y坐标
        """
        logger.info("开始处理CVE卡片列表")

        x_base = 80
        current_y = y_start
        width = 1760

        # 添加data-card背景色
        bg_color_str = 'rgba(10, 66, 117, 0.03)'
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.dml.color import RGBColor
        from pptx.util import Pt
        from src.utils.unit_converter import UnitConverter
        from src.utils.color_parser import ColorParser

        # 计算总高度
        # 基础padding: 20px上下 = 40px
        total_height = 40

        # 处理h3标题（使用动态字号）
        h3_elem = card.find('h3')
        if h3_elem:
            title_text = h3_elem.get_text(strip=True)
            if title_text:
                # 渲染标题
                text_left = UnitConverter.px_to_emu(x_base + 20)
                text_top = UnitConverter.px_to_emu(current_y)
                text_box = pptx_slide.shapes.add_textbox(
                    text_left, text_top,
                    UnitConverter.px_to_emu(width - 40), UnitConverter.px_to_emu(30)
                )
                text_frame = text_box.text_frame
                text_frame.text = title_text
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        # 使用样式计算器动态获取字体大小
                        font_size_px = self.style_computer.get_font_size_pt(h3_elem)
                        run.font.size = Pt(font_size_px) if font_size_px else Pt(20)
                        run.font.color.rgb = ColorParser.get_primary_color()
                        run.font.bold = True
                        run.font.name = self.font_manager.get_font('body')

                current_y += 40  # 标题后间距
                total_height += 40
                logger.info(f"渲染CVE列表标题: {title_text} (字号: {font_size_px if font_size_px else 20}px)")

        # 处理所有cve-card
        cve_cards = card.find_all('div', class_='cve-card')
        logger.info(f"找到{len(cve_cards)}个CVE卡片")

        for i, cve_card in enumerate(cve_cards):
            # 计算每个cve-card的高度
            card_height = self._convert_single_cve_card(cve_card, pptx_slide, shape_converter,
                                                      x_base, current_y, width)
            current_y = card_height
            total_height += card_height - current_y + 15  # 15px是cve-card之间的间距

            # 最后一个卡片不需要额外间距
            if i == len(cve_cards) - 1:
                total_height -= 15

        # 添加data-card的左边框
        shape_converter.add_border_left(x_base, y_start, total_height, 4)

        logger.info(f"CVE卡片列表处理完成，总高度: {total_height}px")
        return y_start + total_height

    def _convert_single_cve_card(self, card, pptx_slide, shape_converter, x, y, width) -> int:
        """
        转换单个CVE卡片

        Args:
            card: cve-card元素
            pptx_slide: PPTX幻灯片
            shape_converter: 形状转换器
            x: X坐标
            y: Y坐标
            width: 宽度

        Returns:
            下一个元素的Y坐标
        """
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.dml.color import RGBColor
        from pptx.util import Pt
        from src.utils.unit_converter import UnitConverter
        from src.utils.color_parser import ColorParser
        from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

        # CVE卡片的padding: 20px
        padding = 20
        content_width = width - padding * 2
        current_y = y + padding

        # 创建CVE卡片背景（渐变效果用单色代替）
        bg_shape = pptx_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            UnitConverter.px_to_emu(x),
            UnitConverter.px_to_emu(y),
            UnitConverter.px_to_emu(width),
            UnitConverter.px_to_emu(100)  # 初始高度，后续会调整
        )
        bg_shape.fill.solid()
        # 使用浅色背景
        bg_shape.fill.fore_color.rgb = RGBColor(248, 250, 252)  # 相当于rgba(10, 66, 117, 0.05)
        # 不设置边框，后续单独添加左边框
        bg_shape.line.fill.background()

        # 处理主要内容区域
        main_content = card.find('div', class_='flex-1')
        if not main_content:
            main_content = card

        # 处理徽章区域
        badge_area = main_content.find('div', class_='flex')
        if badge_area and 'items-center' in badge_area.get('class', []) and 'mb-2' in badge_area.get('class', []):
            badge_y = current_y
            badge_x = x + padding

            # 处理所有徽章
            badges_processed = 0
            for child in badge_area.children:
                if hasattr(child, 'name') and child.name == 'span':
                    badge_text = child.get_text(strip=True)
                    if badge_text:
                        badge_classes = child.get('class', [])

                        # 确定徽章颜色
                        bg_color = RGBColor(255, 255, 255)  # 默认白色
                        text_color = RGBColor(0, 0, 0)  # 默认黑色

                        if 'critical' in badge_classes:
                            bg_color = RGBColor(252, 231, 229)  # 浅红色
                            text_color = RGBColor(220, 38, 38)
                        elif 'high' in badge_classes:
                            bg_color = RGBColor(254, 243, 199)  # 浅橙色
                            text_color = RGBColor(251, 146, 60)
                        elif 'medium' in badge_classes:
                            bg_color = RGBColor(254, 252, 224)  # 浅黄色
                            text_color = RGBColor(251, 191, 36)
                        elif 'exploited' in badge_classes:
                            bg_color = RGBColor(252, 231, 229)  # 浅红色
                            text_color = RGBColor(220, 38, 38)

                        # 计算徽章宽度
                        badge_width = len(badge_text) * 12 + 24  # 估算宽度

                        # 创建徽章背景
                        badge_bg = pptx_slide.shapes.add_shape(
                            MSO_SHAPE.ROUNDED_RECTANGLE,
                            UnitConverter.px_to_emu(badge_x),
                            UnitConverter.px_to_emu(badge_y),
                            UnitConverter.px_to_emu(badge_width),
                            UnitConverter.px_to_emu(24)
                        )
                        badge_bg.fill.solid()
                        badge_bg.fill.fore_color.rgb = bg_color
                        badge_bg.line.fill.background()

                        # 创建徽章文本
                        badge_text_box = pptx_slide.shapes.add_textbox(
                            UnitConverter.px_to_emu(badge_x),
                            UnitConverter.px_to_emu(badge_y + 2),
                            UnitConverter.px_to_emu(badge_width),
                            UnitConverter.px_to_emu(20)
                        )
                        badge_frame = badge_text_box.text_frame
                        badge_frame.text = badge_text
                        for paragraph in badge_frame.paragraphs:
                            paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                            for run in paragraph.runs:
                                # 使用动态字号，而不是硬编码14px
                                badge_font_size = self.style_computer.get_font_size_pt(child)
                                run.font.size = Pt(badge_font_size) if badge_font_size else Pt(14)
                                run.font.bold = True
                                run.font.color.rgb = text_color
                                run.font.name = self.font_manager.get_font('body')

                        badge_x += badge_width + 10
                        badges_processed += 1

            if badges_processed > 0:
                current_y += 35  # 徽章区域高度

        # 处理漏洞名称
        name_p = main_content.find('p', class_=lambda x: x and 'font-medium' in x)
        if name_p:
            name_text = name_p.get_text(strip=True)
            if name_text:
                name_box = pptx_slide.shapes.add_textbox(
                    UnitConverter.px_to_emu(x + padding),
                    UnitConverter.px_to_emu(current_y),
                    UnitConverter.px_to_emu(content_width),
                    UnitConverter.px_to_emu(30)
                )
                name_frame = name_box.text_frame
                name_frame.text = name_text
                for paragraph in name_frame.paragraphs:
                    for run in paragraph.runs:
                        # 使用动态字号，而不是硬编码18px
                        name_font_size = self.style_computer.get_font_size_pt(name_p)
                        run.font.size = Pt(name_font_size) if name_font_size else Pt(18)
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(0, 0, 0)
                        run.font.name = self.font_manager.get_font('body')

                current_y += 30

        # 处理受影响资产
        asset_p = main_content.find('p', class_=lambda x: x and 'text-gray-600' in x)
        if asset_p:
            asset_text = asset_p.get_text(strip=True)
            if asset_text:
                asset_box = pptx_slide.shapes.add_textbox(
                    UnitConverter.px_to_emu(x + padding),
                    UnitConverter.px_to_emu(current_y),
                    UnitConverter.px_to_emu(content_width),
                    UnitConverter.px_to_emu(25)
                )
                asset_frame = asset_box.text_frame
                asset_frame.text = asset_text
                for paragraph in asset_frame.paragraphs:
                    for run in paragraph.runs:
                        # 使用动态字号，而不是硬编码16px
                        asset_font_size = self.style_computer.get_font_size_pt(asset_p)
                        run.font.size = Pt(asset_font_size) if asset_font_size else Pt(16)
                        run.font.color.rgb = RGBColor(102, 102, 102)
                        run.font.name = self.font_manager.get_font('body')

                current_y += 25

        # 处理右侧图标
        icon_elem = card.find('i')
        if icon_elem:
            icon_classes = icon_elem.get('class', [])
            icon_char = self._get_icon_char(icon_classes)
            if icon_char:
                # 简单处理：在右侧添加图标文本
                icon_box = pptx_slide.shapes.add_textbox(
                    UnitConverter.px_to_emu(x + width - 60),
                    UnitConverter.px_to_emu(y + 30),
                    UnitConverter.px_to_emu(40),
                    UnitConverter.px_to_emu(40)
                )
                icon_frame = icon_box.text_frame
                icon_frame.text = icon_char
                for paragraph in icon_frame.paragraphs:
                    paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(36)
                        run.font.color.rgb = RGBColor(220, 38, 38)  # 红色
                        run.font.name = self.font_manager.get_font('body')

        # 调整背景高度
        card_height = current_y - y + padding
        bg_shape.height = UnitConverter.px_to_emu(card_height)

        # 添加左边框
        shape_converter.add_border_left(x, y, card_height, 4)

        return y + card_height

    def _convert_data_card_grid_layout(self, card, grid_container, pptx_slide, shape_converter, y_start: int) -> int:
        """
        转换data-card内的网格布局（如2x2的bullet-point网格）

        Args:
            card: data-card元素
            grid_container: 网格容器元素
            pptx_slide: PPTX幻灯片
            shape_converter: 形状转换器
            y_start: 起始Y坐标

        Returns:
            下一个元素的Y坐标
        """
        logger.info("处理data-card内的网格布局")
        x_base = 80

        # 获取网格列数
        grid_classes = grid_container.get('class', [])
        num_columns = 2  # 默认2列

        # 检查Tailwind CSS网格列类
        for cls in grid_classes:
            if cls.startswith('grid-cols-') and hasattr(self.css_parser, 'tailwind_grid_columns'):
                columns = self.css_parser.tailwind_grid_columns.get(cls)
                if columns:
                    num_columns = columns
                    logger.info(f"检测到网格列数: {num_columns}")
                    break

        # 获取所有bullet-point
        bullet_points = grid_container.find_all('div', class_='bullet-point')

        # 获取标题 - 优先查找h3标签
        h3_elem = card.find('h3')
        title_elem = None
        title_text = None
        actual_title_elem = None  # 实际的标题元素，用于字体大小提取

        if h3_elem:
            # 优先使用h3标签作为标题
            title_text = h3_elem.get_text(strip=True)
            actual_title_elem = h3_elem
            logger.info(f"找到h3标题: {title_text}")
        else:
            # 兼容旧逻辑，查找p标签
            title_elem = card.find('p', class_='primary-color')
            if title_elem:
                title_text = title_elem.get_text(strip=True)
                actual_title_elem = title_elem  # 使用p标签作为标题元素
                logger.info(f"找到p标签标题: {title_text}")
        current_y = y_start

        # 添加data-card背景
        bg_color_str = 'rgba(10, 66, 117, 0.03)'
        from pptx.enum.shapes import MSO_SHAPE

        # 使用ContentHeightCalculator动态计算卡片高度
        # 但由于网格布局特殊，需要手动计算
        
        # 从CSS获取约束
        constraints = self.css_parser.get_height_constraints('.data-card')
        padding_top = constraints.get('padding_top', 15)
        padding_bottom = constraints.get('padding_bottom', 15)
        
        card_height = padding_top  # 顶部padding
        
        # 计算标题高度
        if title_text and actual_title_elem:
            title_font_size_pt = self.style_computer.get_font_size_pt(actual_title_elem)
            title_height = int(title_font_size_pt * 1.5)  # 标题文字高度
            title_margin_bottom = 12  # margin-bottom from CSS
            title_margin_top = 10  # 顶部间距
            card_height += title_margin_top + title_height + title_margin_bottom
            logger.info(f"标题高度: {title_height}px + margin({title_margin_top+title_margin_bottom}px) = {title_margin_top+title_height+title_margin_bottom}px")

        # 计算网格行数
        num_rows = (len(bullet_points) + num_columns - 1) // num_columns
        logger.info(f"网格布局: {num_columns}列 x {num_rows}行")

        # 计算每个bullet-point的平均高度
        if bullet_points:
            # 动态计算每个bullet-point的高度
            bp_heights = []
            card_width = 1760
            col_width = (card_width - padding_top * 2) // num_columns
            
            for bp in bullet_points:
                p_elem = bp.find('p')
                if p_elem:
                    p_font_size_pt = self.style_computer.get_font_size_pt(p_elem)
                    p_text = p_elem.get_text(strip=True)
                    # 计算文本行数
                    available_width = col_width - 40  # 减去icon和间距
                    num_lines = self.height_calculator._calculate_text_lines(p_text, p_font_size_pt, available_width)
                    bp_height = max(20, num_lines * int(p_font_size_pt * 1.5))  # 至少20px（图标高度）
                    bp_heights.append(bp_height)
            
            # 计算每行的最大高度
            row_heights = []
            for row in range(num_rows):
                row_start = row * num_columns
                row_end = min(row_start + num_columns, len(bp_heights))
                row_max_height = max(bp_heights[row_start:row_end]) if row_start < row_end else 30
                row_heights.append(row_max_height)
            
            # 获取网格gap
            grid_gap = self.css_parser.get_gap_size('.space-y-3')  # space-y-3 = 12px
            if grid_gap == 20:  # 如果返回默认值，使用8px（bullet-point的margin-bottom）
                grid_gap = 8
            
            grid_total_height = sum(row_heights) + (num_rows - 1) * grid_gap if num_rows > 0 else 0
            card_height += grid_total_height
            logger.info(f"网格高度: {grid_total_height}px (行高={row_heights}, 行间距={grid_gap}px)")
        
        # 底部padding和额外间距
        card_height += padding_bottom
        card_height += 10  # 额外的底部间距
        
        logger.info(f"data-card网格布局高度计算: padding_top={padding_top}px, 内容={card_height-padding_top-padding_bottom-10}px, padding_bottom={padding_bottom}px, 总高度={card_height}px")

        # 添加背景
        bg_shape = pptx_slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            UnitConverter.px_to_emu(x_base),
            UnitConverter.px_to_emu(y_start),
            UnitConverter.px_to_emu(1760),
            UnitConverter.px_to_emu(card_height)
        )
        bg_shape.fill.solid()
        bg_rgb, alpha = ColorParser.parse_rgba(bg_color_str)
        if bg_rgb:
            if alpha < 1.0:
                bg_rgb = ColorParser.blend_with_white(bg_rgb, alpha)
            bg_shape.fill.fore_color.rgb = bg_rgb
        bg_shape.line.fill.background()
        logger.info(f"添加data-card网格背景，高度={card_height}px")

        # 渲染标题
        if title_text:
            text_left = UnitConverter.px_to_emu(x_base + 20)
            text_top = UnitConverter.px_to_emu(current_y + 10)
            text_box = pptx_slide.shapes.add_textbox(
                text_left, text_top,
                UnitConverter.px_to_emu(1720), UnitConverter.px_to_emu(30)
            )
            text_frame = text_box.text_frame
            text_frame.text = title_text
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    # 使用实际的标题元素来获取字体大小
                    font_size_pt = self.style_computer.get_font_size_pt(actual_title_elem)
                    run.font.size = Pt(font_size_pt)
                    run.font.color.rgb = ColorParser.get_primary_color()
                    run.font.name = self.font_manager.get_font('body')

                    # 智能判断是否应该加粗
                    if self._should_be_bold(actual_title_elem):
                        run.font.bold = True

            current_y += 50  # 标题后间距

        # 处理网格中的bullet-point
        item_width = 1720 // num_columns  # 每列宽度
        
        # 使用之前计算的row_heights来定位每一行
        row_y_positions = [current_y + 10]  # 第一行的Y位置
        grid_gap = 8  # bullet-point之间的间距
        for row_idx in range(1, num_rows):
            prev_y = row_y_positions[row_idx - 1]
            prev_height = row_heights[row_idx - 1] if row_idx - 1 < len(row_heights) else 30
            row_y_positions.append(prev_y + prev_height + grid_gap)

        for idx, bullet_point in enumerate(bullet_points):
            # 计算网格位置
            col = idx % num_columns
            row = idx // num_columns
            item_x = x_base + 20 + col * item_width
            item_y = row_y_positions[row] if row < len(row_y_positions) else current_y + 10

            # 获取图标和文本
            icon_elem = bullet_point.find('i')
            p_elem = bullet_point.find('p')

            if icon_elem and p_elem:
                # 获取图标字符和颜色
                icon_classes = icon_elem.get('class', [])
                icon_char = self._get_icon_char(icon_classes)

                # 获取颜色
                icon_color = self._get_element_color(icon_elem)
                if not icon_color:
                    # 根据图标类确定颜色
                    if 'text-red-600' in icon_classes:
                        icon_color = RGBColor(220, 38, 38)  # 红色
                    elif 'text-orange-600' in icon_classes:
                        icon_color = ColorParser.get_color_by_name('orange')
                    elif 'text-green-600' in icon_classes:
                        icon_color = RGBColor(34, 197, 94)  # 绿色
                    elif 'text-blue-600' in icon_classes:
                        icon_color = RGBColor(59, 130, 246)  # 蓝色
                    elif 'text-purple-600' in icon_classes:
                        icon_color = ColorParser.get_color_by_name('purple')
                    else:
                        icon_color = ColorParser.get_primary_color()

                # 检查是否包含priority-tag
                priority_tag = p_elem.find('span', class_='priority-tag')
                main_text = ""
                tag_text = ""
                tag_color = None

                if priority_tag:
                    # 提取标签文本
                    tag_text = priority_tag.get_text(strip=True)
                    tag_classes = priority_tag.get('class', [])

                    # 确定标签颜色
                    if 'priority-high' in tag_classes:
                        tag_color = RGBColor(239, 68, 68)  # 红色
                    elif 'priority-medium' in tag_classes:
                        tag_color = RGBColor(251, 146, 60)  # 橙色
                    elif 'priority-low' in tag_classes:
                        tag_color = RGBColor(209, 177, 0)  # 黄色

                    # 移除标签后获取主文本
                    priority_tag.extract()
                    main_text = p_elem.get_text(strip=True)
                else:
                    main_text = p_elem.get_text(strip=True)

                # 获取字体大小
                font_size_pt = self._get_font_size_pt(p_elem, default_px=20)

                # 添加图标
                if icon_char:
                    # 图标和文本都在行内，需要垂直居中
                    # 计算图标和文本的垂直居中位置
                    line_height = 60  # 行高
                    vertical_center = item_y + line_height // 2

                    icon_left = UnitConverter.px_to_emu(item_x)
                    icon_top = UnitConverter.px_to_emu(vertical_center - 15)  # 图标垂直居中
                    icon_box = pptx_slide.shapes.add_textbox(
                        icon_left, icon_top,
                        UnitConverter.px_to_emu(30), UnitConverter.px_to_emu(30)
                    )
                    icon_frame = icon_box.text_frame
                    icon_frame.text = icon_char
                    icon_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                    for paragraph in icon_frame.paragraphs:
                        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                        for run in paragraph.runs:
                            run.font.size = Pt(font_size_pt)
                            run.font.color.rgb = icon_color
                            run.font.name = self.font_manager.get_font('body')

                    # 文本在图标右侧，也要垂直居中
                    text_left = UnitConverter.px_to_emu(item_x + 40)
                    text_width = item_width - 60
                    text_top = UnitConverter.px_to_emu(vertical_center - 15)  # 文本垂直居中
                else:
                    # 没有图标，直接显示文本
                    text_left = UnitConverter.px_to_emu(item_x)
                    text_width = item_width - 20
                    line_height = 60
                    vertical_center = item_y + line_height // 2
                    text_top = UnitConverter.px_to_emu(vertical_center - 15)  # 文本垂直居中

                # 添加文本
                text_box = pptx_slide.shapes.add_textbox(
                    text_left, text_top,
                    UnitConverter.px_to_emu(text_width), UnitConverter.px_to_emu(50)  # 增加高度以支持换行
                )
                text_frame = text_box.text_frame
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                text_frame.word_wrap = True

                # 检查是否有strong标签
                strong_elem = p_elem.find('strong')

                if strong_elem or tag_text:
                    # 清除默认段落
                    text_frame.clear()
                    p = text_frame.paragraphs[0]

                    # 处理带strong的文本
                    if strong_elem:
                        strong_text = strong_elem.get_text(strip=True)
                        # 移除strong标签获取剩余文本
                        strong_copy = strong_elem.extract()
                        remaining_text = p_elem.get_text(strip=True)

                        # 添加加粗部分
                        if strong_text:
                            strong_run = p.add_run()
                            strong_run.text = strong_text
                            strong_run.font.size = Pt(font_size_pt)
                            strong_run.font.bold = True
                            strong_run.font.name = self.font_manager.get_font('body')

                        # 添加剩余部分
                        if remaining_text:
                            normal_run = p.add_run()
                            normal_run.text = remaining_text
                            normal_run.font.size = Pt(font_size_pt)
                            normal_run.font.name = self.font_manager.get_font('body')

                        # 添加标签文本（如果存在）
                        if tag_text and tag_color:
                            tag_run = p.add_run()
                            tag_run.text = " " + tag_text
                            tag_font_size_pt = max(10, font_size_pt - 2)
                            tag_run.font.size = Pt(tag_font_size_pt)
                            tag_run.font.color.rgb = tag_color
                            tag_run.font.bold = True
                            tag_run.font.name = self.font_manager.get_font('body')
                    else:
                        # 普通文本 + 标签
                        if main_text:
                            normal_run = p.add_run()
                            normal_run.text = main_text
                            normal_run.font.size = Pt(font_size_pt)
                            normal_run.font.name = self.font_manager.get_font('body')

                        # 添加标签文本（如果存在）
                        if tag_text and tag_color:
                            tag_run = p.add_run()
                            tag_run.text = " " + tag_text
                            tag_font_size_pt = max(10, font_size_pt - 2)
                            tag_run.font.size = Pt(tag_font_size_pt)
                            tag_run.font.color.rgb = tag_color
                            tag_run.font.bold = True
                            tag_run.font.name = self.font_manager.get_font('body')
                else:
                    # 普通文本
                    text_frame.text = main_text
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(font_size_pt)
                            run.font.name = self.font_manager.get_font('body')

        # 计算实际内容高度
        num_rows = (len(bullet_points) + num_columns - 1) // num_columns
        content_height = 50 + num_rows * 60  # 标题高度 + 内容行数 * 行高
        actual_height = max(content_height, 180)  # 最小高度180px

        # 添加左边框
        shape_converter.add_border_left(x_base, y_start, actual_height, 4)

        return y_start + actual_height + 10

    def _get_icon_char(self, icon_classes: list) -> str:
        """根据FontAwesome类获取对应emoji/Unicode字符"""
        icon_map = {
            # === 网络安全相关 ===
            # 核心安全图标
            'fa-shield': '🛡',
            'fa-shield-alt': '🛡',
            'fa-shield-virus': '🦠',
            'fa-virus-slash': '🦠',
            'fa-virus': '🦠',
            'fa-lock': '🔒',
            'fa-unlock': '🔓',
            'fa-key': '🔑',
            'fa-fingerprint': '👆',
            'fa-user-shield': '🛡',
            'fa-user-lock': '🔐',

            # 威胁和警告
            'fa-exclamation-triangle': '⚠',
            'fa-exclamation-circle': '⚠',
            'fa-exclamation': '❗',
            'fa-warning': '⚠️',
            'fa-bell': '🔔',
            'fa-bug': '🐛',
            'fa-radiation': '☢️',
            'fa-biohazard': '☣️',

            # 检查和确认
            'fa-check': '✓',
            'fa-check-circle': '✓',
            'fa-check-square': '☑',
            'fa-check-double': '✓',

            # === 计算机和硬件 ===
            # 设备
            'fa-laptop': '💻',
            'fa-desktop': '🖥',
            'fa-server': '🖥',
            'fa-mobile': '📱',
            'fa-tablet': '📱',
            'fa-wifi': '📶',
            'fa-network-wired': '🔌',
            'fa-usb': '🔌',
            'fa-plug': '🔌',

            # 存储
            'fa-database': '🗄',
            'fa-hdd': '💾',
            'fa-sd-card': '💾',
            'fa-save': '💾',

            # === 人工智能和机器学习 ===
            'fa-robot': '🤖',
            'fa-brain': '🧠',
            'fa-microchip': '💻',
            'fa-memory': '🧠',
            'fa-cpu': '💻',
            'fa-cloud': '☁',
            'fa-cloud-upload-alt': '☁️',
            'fa-cloud-download-alt': '☁️',

            # === 网络和通信 ===
            'fa-globe': '🌐',
            'fa-globe-americas': '🌎',
            'fa-globe-europe': '🌍',
            'fa-globe-asia': '🌏',
            'fa-wifi': '📶',
            'fa-signal': '📶',
            'fa-satellite': '🛰️',
            'fa-ethernet': '🔌',
            'fa-router': '📡',

            # === 法律法规和合规 ===
            'fa-balance-scale': '⚖️',
            'fa-gavel': '🔨',
            'fa-landmark': '🏛️',
            'fa-courthouse': '🏛️',
            'fa-scroll': '📜',
            'fa-file-contract': '📄',
            'fa-file-alt': '📄',
            'fa-file-pdf': '📄',
            'fa-file-word': '📄',
            'fa-file-excel': '📄',

            # === 身份和权限管理 ===
            'fa-user': '👤',
            'fa-users': '👥',
            'fa-user-check': '✅',
            'fa-user-times': '❌',
            'fa-user-plus': '➕',
            'fa-user-minus': '➖',
            'fa-user-cog': '⚙️',
            'fa-id-card': '🪪',
            'fa-passport': '🪪',
            'fa-fingerprint': '👆',

            # === 数据和监控 ===
            'fa-chart-bar': '📊',
            'fa-chart-line': '📈',
            'fa-chart-pie': '📊',
            'fa-chart-area': '📈',
            'fa-table': '📊',
            'fa-database': '🗄',
            'fa-search': '🔍',
            'fa-search-plus': '🔍',
            'fa-search-minus': '🔍',

            # === 攻击和防御 ===
            'fa-swords': '⚔️',
            'fa-crosshairs': '🎯',
            'fa-shield-alt': '🛡',
            'fa-bomb': '💣',
            'fa-hammer': '🔨',
            'fa-wrench': '🔧',
            'fa-tools': '🛠',

            # === 时间和流程 ===
            'fa-clock': '🕐',
            'fa-hourglass': '⏳',
            'fa-hourglass-half': '⏳',
            'fa-calendar': '📅',
            'fa-calendar-alt': '📅',
            'fa-tasks': '☑',
            'fa-list': '📋',
            'fa-clipboard': '📋',
            'fa-clipboard-check': '✅',
            'fa-clipboard-list': '📋',

            # === 系统和设置 ===
            'fa-cog': '⚙',
            'fa-cogs': '⚙️',
            'fa-settings': '⚙️',
            'fa-adjust': '⚙️',
            'fa-sliders-h': '🎚️',
            'fa-toggle-on': '🔛',
            'fa-toggle-off': '🔴',

            # === 文件和数据 ===
            'fa-file': '📄',
            'fa-file-code': '📄',
            'fa-folder': '📁',
            'fa-folder-open': '📂',
            'fa-download': '⬇',
            'fa-upload': '⬆',
            'fa-archive': '📦',
            'fa-file-archive': '📦',

            # === 通信和消息 ===
            'fa-envelope': '✉',
            'fa-envelope-open': '📧',
            'fa-comments': '💬',
            'fa-comment': '💬',
            'fa-comment-dots': '💬',
            'fa-phone': '📞',
            'fa-video': '📹',

            # === 基础图标 ===
            'fa-check': '✓',
            'fa-check-circle': '✓',
            'fa-times': '✗',
            'fa-times-circle': '❌',
            'fa-plus': '+',
            'fa-plus-circle': '⭕',
            'fa-minus': '-',
            'fa-minus-circle': '⭕',
            'fa-arrow-right': '→',
            'fa-arrow-left': '←',
            'fa-arrow-up': '↑',
            'fa-arrow-down': '↓',
            'fa-sync': '🔄',
            'fa-redo': '↻',
            'fa-undo': '↺',
            'fa-play': '▶',
            'fa-pause': '⏸',
            'fa-stop': '⏹',
            'fa-home': '🏠',
            'fa-building': '🏢',

            # === 新增：常用FontAwesome图标 ===
            # 状态和标记
            'fa-info-circle': 'ℹ',
            'fa-question-circle': '❓',
            'fa-asterisk': '*',
            'fa-star': '⭐',
            'fa-heart': '♥',
            'fa-heartbeat': '💓',
            'fa-fire': '🔥',
            'fa-bolt': '⚡',
            'fa-flash': '⚡',
            'fa-magic': '✨',
            'fa-sparkles': '✨',

            # 方向和导航
            'fa-chevron-right': '›',
            'fa-chevron-left': '‹',
            'fa-chevron-up': '⌃',
            'fa-chevron-down': '⌄',
            'fa-angle-right': '›',
            'fa-angle-left': '‹',
            'fa-angle-up': '⌃',
            'fa-angle-down': '⌄',
            'fa-caret-right': '▶',
            'fa-caret-left': '◀',
            'fa-caret-up': '▲',
            'fa-caret-down': '▼',

            # 商务和金融
            'fa-dollar-sign': '$',
            'fa-euro-sign': '€',
            'fa-pound-sign': '£',
            'fa-yen-sign': '¥',
            'fa-coins': '🪙',
            'fa-wallet': '👛',
            'fa-credit-card': '💳',
            'fa-chart-pie': '📊',
            'fa-pie-chart': '📊',
            'fa-chart-simple': '📊',

            # 云和数据
            'fa-cloud': '☁',
            'fa-cloud-arrow-up': '☁️',
            'fa-cloud-arrow-down': '☁️',
            'fa-cloud-download': '☁️',
            'fa-cloud-upload': '☁️',
            'fa-server': '🖥',
            'fa-desktop': '🖥',
            'fa-laptop': '💻',
            'fa-mobile': '📱',
            'fa-tablet': '📱',

            # 编辑和创作
            'fa-edit': '✏️',
            'fa-pen': '🖊️',
            'fa-pencil': '✏️',
            'fa-eraser': '🧹',
            'fa-paint-brush': '🖌️',
            'fa-palette': '🎨',
            'fa-camera': '📷',
            'fa-video': '📹',
            'fa-film': '🎬',
            'fa-music': '🎵',
            'fa-headphones': '🎧',
            'fa-microphone': '🎤',

            # 社交和用户
            'fa-user': '👤',
            'fa-user-circle': '👤',
            'fa-user-group': '👥',
            'fa-users': '👥',
            'fa-user-tie': '👔',
            'fa-user-graduate': '🎓',
            'fa-user-doctor': '👨‍⚕️',
            'fa-user-ninja': '🥷',
            'fa-user-astronaut': '👨‍🚀',

            # 环境和自然
            'fa-tree': '🌳',
            'fa-leaf': '🍃',
            'fa-seedling': '🌱',
            'fa-sun': '☀️',
            'fa-moon': '🌙',
            'fa-star': '⭐',
            'fa-snowflake': '❄️',
            'fa-fire': '🔥',
            'fa-water': '💧',
            'fa-droplet': '💧',

            # 交通和移动
            'fa-car': '🚗',
            'fa-plane': '✈️',
            'fa-ship': '🚢',
            'fa-train': '🚂',
            'fa-bicycle': '🚴',
            'fa-motorcycle': '🏍️',
            'fa-rocket': '🚀',
            'fa-satellite': '🛰️',
            'fa-helicopter': '🚁',

            # 食物和饮料
            'fa-utensils': '🍴',
            'fa-coffee': '☕',
            'fa-glass': '🥤',
            'fa-wine-glass': '🍷',
            'fa-beer': '🍺',
            'fa-pizza-slice': '🍕',
            'fa-hamburger': '🍔',
            'fa-ice-cream': '🍦',

            # 其他新增
            'fa-cloud-showers-heavy': '🌧️',
            'fa-gift': '🎁',
            'fa-tag': '🏷️',
            'fa-tags': '🏷️',
            'fa-certificate': '🎓',
            'fa-award': '🏆',
            'fa-trophy': '🏆',
            'fa-medal': '🏅',
            'fa-ribbon': '🎀',
            'fa-flag': '🚩',
            'fa-bookmark': '🔖',
            'fa-thumbtack': '📌',
            'fa-pushpin': '📌',
        }

        for cls in icon_classes:
            if cls in icon_map:
                return icon_map[cls]

        # 如果找不到匹配，返回默认图标
        return '●'

    def _get_font_size_pt(self, element, default_px: int = 16) -> int:
        """
        获取元素的字体大小（以pt为单位）

        Args:
            element: HTML元素
            default_px: 默认字体大小（像素）

        Returns:
            字体大小（pt）
        """
        # 1. 首先尝试从已有的style_computer获取
        font_size_pt = self.style_computer.get_font_size_pt(element)
        if font_size_pt:
            return int(font_size_pt)

        # 2. 检查元素的style属性
        style = element.get('style', '')
        if 'font-size' in style:
            import re
            match = re.search(r'font-size:\s*(\d+)px', style)
            if match:
                px_size = int(match.group(1))
                # px转pt的近似公式：1px ≈ 0.75pt
                return int(px_size * 0.75)

        # 3. 检查元素的class属性中的Tailwind CSS字体大小类
        classes = element.get('class', [])
        if isinstance(classes, str):
            classes = classes.split()

        # Tailwind字体大小映射
        tailwind_font_sizes = {
            'text-xs': 12,    # 12px = 9pt
            'text-sm': 14,    # 14px = 10.5pt
            'text-base': 16,  # 16px = 12pt
            'text-lg': 18,    # 18px = 13.5pt
            'text-xl': 20,    # 20px = 15pt
            'text-2xl': 24,   # 24px = 18pt
            'text-3xl': 30,   # 30px = 22.5pt
            'text-4xl': 36,   # 36px = 27pt
            'text-5xl': 48,   # 48px = 36pt
            'text-6xl': 60,   # 60px = 45pt
            'text-7xl': 72,   # 72px = 54pt
            'text-8xl': 96,   # 96px = 72pt
            'text-9xl': 128,  # 128px = 96pt
        }

        for cls in classes:
            if cls in tailwind_font_sizes:
                px_size = tailwind_font_sizes[cls]
                return int(px_size * 0.75)

        # 4. 检查父元素（特别是bullet-point）
        parent = element.parent
        if parent:
            parent_classes = parent.get('class', [])
            if isinstance(parent_classes, str):
                parent_classes = parent_classes.split()
            if 'bullet-point' in parent_classes:
                # bullet-point通常有25px的字体大小
                return int(25 * 0.75)  # 19pt

        # 5. 返回默认值
        return int(default_px * 0.75)

    def _get_element_color(self, element):
        """
        获取元素的颜色，支持Tailwind CSS类

        Args:
            element: BeautifulSoup元素

        Returns:
            RGBColor对象，如果没有找到颜色则返回None
        """
        if not element:
            return None

        # 检查Tailwind CSS颜色类
        classes = element.get('class', [])
        for cls in classes:
            # 优先检查primary-color类
            if cls == 'primary-color':
                return ColorParser.get_primary_color()
            elif cls.startswith('text-') and hasattr(self.css_parser, 'tailwind_colors'):
                color = self.css_parser.tailwind_colors.get(cls)
                if color:
                    return ColorParser.parse_color(color)

        # 检查CSS样式中的颜色
        computed_style = self.style_computer.compute_computed_style(element)
        color_str = computed_style.get('color')
        if color_str:
            return ColorParser.parse_color(color_str)

        return None

    def _should_be_bold(self, element):
        """
        判断元素是否应该加粗

        Args:
            element: HTML元素

        Returns:
            bool: 是否应该加粗
        """
        if not element:
            return False

        # 1. 检查内联样式的font-weight
        style_str = element.get('style', '')
        if style_str:
            import re
            weight_match = re.search(r'font-weight:\s*([^;]+)', style_str)
            if weight_match:
                weight_str = weight_match.group(1).strip()
                # 转换常见的font-weight值
                if weight_str in ['bold', '700', '600', '800', '900']:
                    return True
                elif weight_str in ['normal', '400', '300']:
                    return False

        # 2. 检查类名中的加粗相关类
        classes = element.get('class', [])
        bold_classes = ['font-bold', 'font-semibold', 'font-extrabold', 'font-black']
        for cls in classes:
            if cls in bold_classes:
                return True

        # 3. 根据元素类型判断
        tag_name = element.name.lower()
        if tag_name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
            # 标题默认加粗，除非明确指定了font-weight: normal
            if self.css_parser:
                # 获取CSS中定义的font-weight
                css_style = self.css_parser.get_style(tag_name)
                if css_style and 'font-weight' in css_style:
                    weight = css_style['font-weight']
                    if weight in ['normal', '400', '300']:
                        return False
                    elif weight in ['bold', '600', '700', '800', '900']:
                        return True
            # 默认情况下，h1和h2加粗，h3根据CSS决定
            return tag_name in ['h1', 'h2']

        # 4. 检查strong标签
        if element.name == 'strong':
            return True

        # 5. 检查父元素的加粗设置（如b标签内的文本）
        parent = element.parent
        if parent and parent.name in ['strong', 'b']:
            return True

        return False

    def _determine_layout_direction(self, box) -> str:
        """
        智能判断布局方向：水平或垂直

        Args:
            box: stat-box元素

        Returns:
            'horizontal' 或 'vertical'
        """
        # 检查CSS样式，特别是align-items属性
        # align-items: center 通常表示水平布局
        # align-items: flex-start 或未设置通常表示垂直布局

        # 方法1：检查内联样式
        inline_style = box.get('style', '')
        if 'align-items' in inline_style:
            if 'center' in inline_style:
                logger.info("检测到align-items: center，使用水平布局")
                return 'horizontal'
            elif 'flex-start' in inline_style or 'start' in inline_style:
                logger.info("检测到align-items: flex-start，使用垂直布局")
                return 'vertical'

        # 方法2：从CSS解析器获取样式
        computed_styles = self.css_parser.get_style('.stat-box')
        align_items = computed_styles.get('align-items', '').lower() if computed_styles else ''

        if 'center' in align_items:
            logger.info("从CSS检测到align-items: center，使用水平布局")
            return 'horizontal'
        elif 'flex-start' in align_items or 'start' in align_items:
            logger.info("从CSS检测到align-items: flex-start，使用垂直布局")
            return 'vertical'

        # 方法3：检查具体的HTML结构
        # 如果有text-center类，倾向于垂直布局
        box_classes = box.get('class', [])
        if 'text-center' in box_classes:
            logger.info("检测到text-center类，使用垂直布局")
            return 'vertical'

        # 方法4：检查子元素的对齐方式
        title_elem = box.find('div', class_='stat-title')
        if title_elem:
            title_classes = title_elem.get('class', [])
            if 'text-center' in title_classes:
                logger.info("检测到标题居中，使用垂直布局")
                return 'vertical'

        # 默认策略：根据常见模式判断
        # 如果图标存在且有居中类，很可能是垂直布局
        icon = box.find('i')
        if icon:
            icon_parent_classes = icon.parent.get('class', []) if icon.parent else []
            if 'text-center' in icon_parent_classes:
                logger.info("检测到图标居中，使用垂直布局")
                return 'vertical'

        # 默认使用垂直布局（更常见的模式）
        logger.info("未检测到明确的布局方向，使用默认垂直布局")
        return 'vertical'

    def _determine_text_alignment(self, box) -> int:
        """
        智能判断文字对齐方式

        Args:
            box: stat-box元素

        Returns:
            PPTX对齐常量: PP_PARAGRAPH_ALIGNMENT.LEFT, CENTER, RIGHT
        """
        # 方法1：检查内联样式
        inline_style = box.get('style', '')
        if 'text-align' in inline_style:
            if 'center' in inline_style:
                return PP_PARAGRAPH_ALIGNMENT.CENTER
            elif 'right' in inline_style:
                return PP_PARAGRAPH_ALIGNMENT.RIGHT
            elif 'left' in inline_style:
                return PP_PARAGRAPH_ALIGNMENT.LEFT

        # 方法2：检查CSS类
        box_classes = box.get('class', [])
        if 'text-center' in box_classes:
            logger.info("检测到text-center类，使用居中对齐")
            return PP_PARAGRAPH_ALIGNMENT.CENTER
        elif 'text-right' in box_classes:
            logger.info("检测到text-right类，使用右对齐")
            return PP_PARAGRAPH_ALIGNMENT.RIGHT
        elif 'text-left' in box_classes:
            logger.info("检测到text-left类，使用左对齐")
            return PP_PARAGRAPH_ALIGNMENT.LEFT

        # 方法3：检查子元素的对齐类
        title_elem = box.find('div', class_='stat-title')
        if title_elem:
            title_classes = title_elem.get('class', [])
            if 'text-center' in title_classes:
                logger.info("检测到标题居中类，使用居中对齐")
                return PP_PARAGRAPH_ALIGNMENT.CENTER
            elif 'text-right' in title_classes:
                logger.info("检测到标题右对齐类，使用右对齐")
                return PP_PARAGRAPH_ALIGNMENT.RIGHT
            elif 'text-left' in title_classes:
                logger.info("检测到标题左对齐类，使用左对齐")
                return PP_PARAGRAPH_ALIGNMENT.LEFT

        # 方法4：从CSS解析器获取样式
        computed_styles = self.css_parser.get_style('.stat-box')
        text_align = computed_styles.get('text-align', '').lower() if computed_styles else ''

        if 'center' in text_align:
            return PP_PARAGRAPH_ALIGNMENT.CENTER
        elif 'right' in text_align:
            return PP_PARAGRAPH_ALIGNMENT.RIGHT
        elif 'left' in text_align:
            return PP_PARAGRAPH_ALIGNMENT.LEFT

        # 方法5：根据布局方向智能推断
        layout_direction = self._determine_layout_direction(box)
        if layout_direction == 'horizontal':
            # 水平布局通常左对齐更美观
            logger.info("水平布局，默认使用左对齐")
            return PP_PARAGRAPH_ALIGNMENT.LEFT
        else:
            # 垂直布局通常居中对齐更美观
            logger.info("垂直布局，默认使用居中对齐")
            return PP_PARAGRAPH_ALIGNMENT.CENTER

    def _has_numbered_list_pattern(self, container) -> bool:
        """
        检测容器是否包含数字列表模式

        Args:
            container: 容器元素

        Returns:
            是否包含数字列表
        """
        # 检查子元素是否包含数字
        children = container.find_all(recursive=False)
        for child in children:
            text = child.get_text(strip=True)
            if text and text[0].isdigit():
                return True

            # 检查是否包含number相关的类
            child_classes = child.get('class', [])
            if any('number' in cls or 'num' in cls or 'count' in cls for cls in child_classes):
                return True

        return False

    def _convert_numbered_list_container(self, container, pptx_slide, y_start) -> int:
        """
        转换数字列表容器

        Args:
            container: 容器元素
            pptx_slide: PPTX幻灯片
            y_start: 起始Y坐标

        Returns:
            下一个元素的Y坐标
        """
        logger.info("处理数字列表容器")

        # 初始化文本转换器
        text_converter = TextConverter(pptx_slide, self.css_parser)

        # 检查是否是toc-item结构
        if 'toc-item' in container.get('class', []):
            # 处理单个toc-item
            number_elem = container.find('div', class_='toc-number')
            text_elem = container.find('div', class_='toc-title')

            if number_elem and text_elem:
                numbered_item = {
                    'type': 'toc',
                    'container': container,
                    'number_elem': number_elem,
                    'text_elem': text_elem,
                    'number': number_elem.get_text(strip=True),
                    'text': text_elem.get_text(strip=True)
                }
                return text_converter.convert_numbered_list(numbered_item, 80, y_start)

        # 处理其他数字列表格式
        text = container.get_text(strip=True)
        if text and text[0].isdigit():
            # 尝试分离数字和文本
            import re
            match = re.match(r'^(\d+)[\.\)\s]*\s*(.*)', text)
            if match:
                numbered_item = {
                    'type': 'paragraph_numbered',
                    'container': container,
                    'number_elem': container,
                    'text_elem': container,
                    'number': match.group(1),
                    'text': match.group(2)
                }
                return text_converter.convert_numbered_list(numbered_item, 80, y_start)

        # 降级处理为普通段落
        return self._convert_generic_card(container, pptx_slide, y_start, card_type='numbered_list')

    def convert_to_pptx_shared(self, output_dir: str = "output", output_filename: str = None):
        """
        转换HTML到PPTX（共享模式，用于批量转换）
        不创建新的PPTX文件，只将内容添加到共享的presentation中

        Args:
            output_dir: 输出目录（仅用于日志）
            output_filename: 输出文件名（仅用于日志）

        Returns:
            生成的幻灯片数量
        """
        logger.info(f"转换HTML文件到共享PPTX: {self.html_path}")

        # 获取所有幻灯片
        slides = self.html_parser.get_slides()
        slide_count = 0

        for slide_html in slides:
            logger.info(f"\n处理幻灯片...")
            slide_count += 1

            # 创建空白幻灯片
            pptx_slide = self.pptx_builder.add_blank_slide()

            # 初始化转换器
            text_converter = TextConverter(pptx_slide, self.css_parser)
            table_converter = TableConverter(pptx_slide, self.css_parser)
            shape_converter = ShapeConverter(pptx_slide, self.css_parser)

            # 1. 添加顶部装饰条
            shape_converter.add_top_bar()

            # 2. 添加标题和副标题
            title_info = self.html_parser.get_title_info(slide_html)
            if title_info:
                # content-section的padding-top是20px
                title_end_y = text_converter.convert_title(
                    title_info['text'],
                    title_info['subtitle'],
                    x=80,
                    y=20
                )
            else:
                title_end_y = 100

            # 3. 转换主要内容
            main_content = self.html_parser.get_main_content(slide_html)
            if main_content:
                self._convert_content_container(main_content, pptx_slide, title_end_y + 20, shape_converter)

            # 4. 添加页码
            self._add_page_number(pptx_slide, slide_count)

            logger.info(f"成功处理幻灯片 {slide_count}")

        return slide_count


def main():
    """主函数"""
    if len(sys.argv) < 2:
        print("用法: python main.py <html文件路径> [输出pptx路径]")
        sys.exit(1)

    html_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else "output/output.pptx"

    # 执行转换
    converter = HTML2PPTX(html_path)
    converter.convert(output_path)


if __name__ == "__main__":
    main()
