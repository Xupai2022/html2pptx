"""
PPTX构建器
"""

from pptx import Presentation
from pptx.util import Inches
from pathlib import Path

from src.utils.unit_converter import UnitConverter
from src.utils.logger import setup_logger

logger = setup_logger(__name__)


class PPTXBuilder:
    """PPTX构建器"""

    def __init__(self):
        """初始化PPTX构建器"""
        self.prs = Presentation()

        # 设置幻灯片尺寸为16:9 (1920x1080)
        self.prs.slide_width = UnitConverter.px_to_emu(1920)
        self.prs.slide_height = UnitConverter.px_to_emu(1080)

        logger.info("初始化PPTX,尺寸: 1920x1080")

    def add_blank_slide(self):
        """
        添加空白幻灯片

        Returns:
            幻灯片对象
        """
        blank_slide_layout = self.prs.slide_layouts[6]  # 空白布局
        slide = self.prs.slides.add_slide(blank_slide_layout)
        logger.info(f"添加空白幻灯片,当前总数: {len(self.prs.slides)}")
        return slide

    def save(self, output_path: str):
        """
        保存PPTX文件

        Args:
            output_path: 输出路径
        """
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)

        self.prs.save(str(output_path))
        logger.info(f"PPTX已保存: {output_path}")

    def get_presentation(self):
        """获取Presentation对象"""
        return self.prs
