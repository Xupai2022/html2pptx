"""
时间线转换器
处理timeline时间线结构的转换
"""

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

from src.converters.base_converter import BaseConverter
from src.utils.unit_converter import UnitConverter
from src.utils.color_parser import ColorParser
from src.utils.logger import setup_logger
from src.utils.font_manager import get_font_manager
from src.utils.style_computer import get_style_computer

logger = setup_logger(__name__)


class TimelineConverter(BaseConverter):
    """时间线转换器"""

    def convert_timeline(self, timeline_element, x: int, y: int, width: int = 1760) -> int:
        """
        转换时间线结构

        Args:
            timeline_element: Timeline元素
            x: X坐标(px)
            y: Y坐标(px)
            width: 宽度(px)

        Returns:
            下一个元素的Y坐标
        """
        if not timeline_element:
            return y

        # 查找所有timeline-item
        timeline_items = timeline_element.find_all('div', class_='timeline-item')

        if not timeline_items:
            logger.warning("timeline中没有找到timeline-item")
            return y

        logger.info(f"找到 {len(timeline_items)} 个timeline-item")

        current_y = y

        for idx, item in enumerate(timeline_items):
            # 获取时间线图标（数字）
            icon_elem = item.find('div', class_='timeline-icon')
            icon_text = icon_elem.get_text(strip=True) if icon_elem else str(idx + 1)

            # 获取标题
            title_elem = item.find('div', class_='timeline-title')
            title_text = title_elem.get_text(strip=True) if title_elem else ""

            # 获取内容（timeline-content下的p标签）
            content_elem = item.find('div', class_='timeline-content')
            content_text = ""
            if content_elem:
                p_elem = content_elem.find('p')
                if p_elem:
                    content_text = p_elem.get_text(strip=True)

            # 渲染这个timeline-item
            current_y = self._render_timeline_item(
                icon_text, title_text, content_text,
                x, current_y, width
            )

        return current_y

    def _render_timeline_item(
        self,
        icon_text: str,
        title_text: str,
        content_text: str,
        x: int,
        y: int,
        width: int
    ) -> int:
        """
        渲染单个时间线项目

        Args:
            icon_text: 图标文本（数字）
            title_text: 标题文本
            content_text: 内容文本
            x: X坐标(px)
            y: Y坐标(px)
            width: 宽度(px)

        Returns:
            下一个元素的Y坐标
        """
        left_emu = UnitConverter.px_to_emu(x)
        top_emu = UnitConverter.px_to_emu(y)

        # 1. 绘制圆形图标（数字）
        icon_size = 25  # px
        icon_left = UnitConverter.px_to_emu(x)
        icon_top = UnitConverter.px_to_emu(y)
        icon_size_emu = UnitConverter.px_to_emu(icon_size)

        # 使用椭圆形状创建圆形
        from pptx.enum.shapes import MSO_SHAPE
        icon_shape = self.slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            icon_left, icon_top,
            icon_size_emu, icon_size_emu
        )

        # 设置圆形样式
        icon_shape.fill.solid()
        icon_shape.fill.fore_color.rgb = ColorParser.get_primary_color()
        icon_shape.line.fill.background()  # 无边框

        # 添加数字文本到圆形中
        text_frame = icon_shape.text_frame
        text_frame.text = icon_text
        text_frame.vertical_anchor = 1  # 居中
        text_frame.word_wrap = True

        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(12)
                run.font.color.rgb = ColorParser.parse_color('#FFFFFF')
                run.font.bold = True
                run.font.name = get_font_manager(self.css_parser).get_font('body')

        # 2. 绘制左侧竖线（连接线）
        line_left = UnitConverter.px_to_emu(x + 12)  # 圆心位置
        line_top = UnitConverter.px_to_emu(y + icon_size)
        line_height = UnitConverter.px_to_emu(60)  # 竖线高度

        line_shape = self.slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            line_left, line_top,
            UnitConverter.px_to_emu(2), line_height
        )
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = ColorParser.get_primary_color()
        line_shape.line.fill.background()

        # 3. 添加标题文本框
        text_x = x + 40  # 圆形右侧
        text_width = width - 40

        # 获取样式计算器
        style_computer = get_style_computer(self.css_parser)
        font_manager = get_font_manager(self.css_parser)

        # 创建临时timeline-title元素来获取字体大小
        from bs4 import BeautifulSoup
        temp_soup_title = BeautifulSoup('<div class="timeline-title">' + title_text + '</div>', 'html.parser')
        title_elem = temp_soup_title.find('div', class_='timeline-title')

        # 获取timeline-title的字体大小
        title_font_size_pt = style_computer.get_font_size_pt(title_elem)
        # 转换回px用于高度计算
        title_font_size_px = UnitConverter.pt_to_px(title_font_size_pt)
        title_height = int(title_font_size_px * 1.3)  # timeline-title行高约1.3

        title_left = UnitConverter.px_to_emu(text_x)
        title_top = UnitConverter.px_to_emu(y)
        title_width = UnitConverter.px_to_emu(text_width)
        title_height_emu = UnitConverter.px_to_emu(title_height)

        title_box = self.slide.shapes.add_textbox(
            title_left, title_top, title_width, title_height_emu
        )
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_frame.word_wrap = True

        for paragraph in title_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(title_font_size_pt)
                run.font.color.rgb = ColorParser.get_primary_color()
                run.font.bold = True
                run.font.name = font_manager.get_font('body')

        # 4. 添加内容文本框
        # 创建临时p元素来获取字体大小
        temp_soup_content = BeautifulSoup('<p>' + content_text + '</p>', 'html.parser')
        content_elem = temp_soup_content.p

        # 获取内容的字体大小
        content_font_size_pt = style_computer.get_font_size_pt(content_elem)
        # 转换回px用于高度计算
        content_font_size_px = UnitConverter.pt_to_px(content_font_size_pt)
        content_height = int(content_font_size_px * 1.5)  # 内容行高1.5

        content_left = UnitConverter.px_to_emu(text_x)
        content_top = UnitConverter.px_to_emu(y + title_height + 3)  # 标题下方3px间距
        content_width = UnitConverter.px_to_emu(text_width)
        content_height_emu = UnitConverter.px_to_emu(content_height)

        content_box = self.slide.shapes.add_textbox(
            content_left, content_top, content_width, content_height_emu
        )
        content_frame = content_box.text_frame
        content_frame.text = content_text
        content_frame.word_wrap = True

        for paragraph in content_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(content_font_size_pt)
                run.font.color.rgb = ColorParser.parse_color('#333333')
                run.font.name = font_manager.get_font('body')

        # 返回下一个item的Y坐标
        return y + 85  # 每个item占用约85px高度

    def convert(self, element, **kwargs):
        """转换时间线元素"""
        x = kwargs.get('x', 80)
        y = kwargs.get('y', 0)
        width = kwargs.get('width', 1760)
        return self.convert_timeline(element, x, y, width)
