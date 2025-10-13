"""
图表转换器
处理Chart.js图表的转换,支持截图插入
"""

from pathlib import Path
from typing import Optional

from src.converters.base_converter import BaseConverter
from src.utils.unit_converter import UnitConverter
from src.utils.chart_capture import ChartCapture
from src.utils.logger import setup_logger
from src.utils.font_manager import get_font_manager

logger = setup_logger(__name__)


class ChartConverter(BaseConverter):
    """图表转换器"""

    def __init__(self, slide, css_parser, html_path: str = None):
        """
        初始化图表转换器

        Args:
            slide: PPTX幻灯片对象
            css_parser: CSS解析器
            html_path: HTML文件路径(用于截图)
        """
        super().__init__(slide, css_parser)
        self.html_path = html_path
        self.chart_capturer = ChartCapture()

    def convert_chart(
        self,
        canvas_element,
        x: int,
        y: int,
        width: int = 1760,
        height: int = 200,
        use_screenshot: bool = True
    ) -> bool:
        """
        转换图表

        Args:
            canvas_element: Canvas元素
            x: X坐标(px)
            y: Y坐标(px)
            width: 宽度(px)
            height: 高度(px)
            use_screenshot: 是否使用截图(False则显示占位文本)

        Returns:
            是否成功转换
        """
        # 检查Playwright是否可用
        if use_screenshot and not ChartCapture.is_available():
            logger.warning("Playwright不可用,使用占位文本")
            use_screenshot = False

        # 获取canvas ID
        canvas_id = canvas_element.get('id')
        if not canvas_id:
            logger.warning("Canvas元素缺少ID,使用默认选择器")
            canvas_selector = "canvas"
        else:
            canvas_selector = f"#{canvas_id}"

        # 尝试截图
        if use_screenshot and self.html_path:
            screenshot_path = self._capture_chart_screenshot(canvas_selector)

            if screenshot_path:
                return self._insert_chart_image(screenshot_path, x, y, width, height)

        # 截图失败,显示占位文本
        return self._insert_placeholder(x, y, width, height)

    def _capture_chart_screenshot(self, canvas_selector: str) -> Optional[str]:
        """
        截取图表截图

        Args:
            canvas_selector: Canvas选择器

        Returns:
            截图路径,失败返回None
        """
        if not self.html_path:
            logger.error("HTML路径未设置,无法截图")
            return None

        try:
            logger.info(f"开始截取图表: {canvas_selector}")
            screenshot_path = self.chart_capturer.capture_chart(
                self.html_path,
                canvas_selector,
                wait_time=2000
            )

            if screenshot_path:
                logger.info(f"图表截图成功: {screenshot_path}")
                return screenshot_path
            else:
                logger.error("图表截图失败")
                return None

        except Exception as e:
            logger.error(f"截图异常: {e}")
            return None

    def _insert_chart_image(
        self,
        image_path: str,
        x: int,
        y: int,
        width: int,
        height: int
    ) -> bool:
        """
        插入图表图片到PPTX

        Args:
            image_path: 图片路径
            x, y: 坐标(px)
            width, height: 尺寸(px)

        Returns:
            是否成功
        """
        try:
            left = UnitConverter.px_to_emu(x)
            top = UnitConverter.px_to_emu(y)
            pic_width = UnitConverter.px_to_emu(width)
            pic_height = UnitConverter.px_to_emu(height)

            # 插入图片
            self.slide.shapes.add_picture(
                image_path,
                left,
                top,
                width=pic_width,
                height=pic_height
            )

            logger.info(f"图表图片已插入: {image_path}")
            return True

        except Exception as e:
            logger.error(f"插入图片失败: {e}")
            return False

    def _insert_placeholder(
        self,
        x: int,
        y: int,
        width: int,
        height: int
    ) -> bool:
        """
        插入占位文本

        Args:
            x, y: 坐标(px)
            width, height: 尺寸(px)

        Returns:
            是否成功
        """
        try:
            from pptx.util import Pt
            from pptx.enum.text import MSO_ANCHOR
            from src.utils.color_parser import ColorParser

            left = UnitConverter.px_to_emu(x)
            top = UnitConverter.px_to_emu(y)
            box_width = UnitConverter.px_to_emu(width)
            box_height = UnitConverter.px_to_emu(height)

            text_box = self.slide.shapes.add_textbox(left, top, box_width, box_height)
            text_frame = text_box.text_frame
            text_frame.text = "[图表占位 - Chart.js图表]\n(需要安装Playwright: pip install playwright)"
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            text_frame.word_wrap = True

            for paragraph in text_frame.paragraphs:
                paragraph.alignment = 1  # 居中
                for run in paragraph.runs:
                    run.font.size = Pt(20)
                    run.font.color.rgb = ColorParser.parse_color('#999')
                    run.font.name = get_font_manager(self.css_parser).get_font('body')

            logger.info("插入图表占位文本")
            return True

        except Exception as e:
            logger.error(f"插入占位文本失败: {e}")
            return False

    def convert(self, element, **kwargs):
        """转换图表元素"""
        return self.convert_chart(element, **kwargs)
