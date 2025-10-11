"""
表格转换器
"""

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from src.converters.base_converter import BaseConverter
from src.utils.unit_converter import UnitConverter
from src.utils.color_parser import ColorParser
from src.utils.logger import setup_logger

logger = setup_logger(__name__)


class TableConverter(BaseConverter):
    """表格转换器"""

    def convert(self, table_element, x: int, y: int, **kwargs):
        """
        转换表格

        Args:
            table_element: 表格HTML元素
            x: X坐标(px)
            y: Y坐标(px)
        """
        # 提取表格数据
        rows = table_element.find_all('tr')
        if not rows:
            return

        # 计算列数
        first_row = rows[0]
        cols = len(first_row.find_all(['th', 'td']))

        # 创建PPTX表格
        left = UnitConverter.px_to_emu(x)
        top = UnitConverter.px_to_emu(y)
        width = UnitConverter.px_to_emu(1760)

        # 计算行高和列宽
        row_height = UnitConverter.px_to_emu(50)
        col_width = width // cols

        pptx_table = self.slide.shapes.add_table(
            len(rows), cols, left, top, width, row_height * len(rows)
        ).table

        # 填充表格
        for row_idx, row in enumerate(rows):
            cells = row.find_all(['th', 'td'])

            for col_idx, cell in enumerate(cells):
                pptx_cell = pptx_table.cell(row_idx, col_idx)
                text = cell.get_text(strip=True)
                pptx_cell.text = text

                # 设置文本样式
                text_frame = pptx_cell.text_frame
                text_frame.word_wrap = True
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

                for paragraph in text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.LEFT

                    for run in paragraph.runs:
                        run.font.size = Pt(24)
                        run.font.name = 'Microsoft YaHei'

                        # 表头样式
                        if cell.name == 'th':
                            run.font.bold = True
                            run.font.color.rgb = ColorParser.WHITE
                            pptx_cell.fill.solid()
                            pptx_cell.fill.fore_color.rgb = ColorParser.get_primary_color()
                        else:
                            # 第一列加粗
                            if col_idx == 0:
                                run.font.bold = True
                                run.font.color.rgb = ColorParser.get_primary_color()

        logger.info(f"添加表格: {len(rows)}行 x {cols}列")
