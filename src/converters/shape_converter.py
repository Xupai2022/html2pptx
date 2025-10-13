"""
形状转换器
处理装饰条、进度条等形状元素
"""

from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE
from src.converters.base_converter import BaseConverter
from src.utils.unit_converter import UnitConverter
from src.utils.color_parser import ColorParser
from src.utils.logger import setup_logger

logger = setup_logger(__name__)


class ShapeConverter(BaseConverter):
    """形状转换器"""

    def add_top_bar(self):
        """添加顶部装饰条"""
        left = 0
        top = 0
        width = UnitConverter.px_to_emu(1920)
        height = UnitConverter.px_to_emu(10)

        shape = self.slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = ColorParser.get_primary_color()
        shape.line.fill.background()

        logger.info("添加顶部装饰条")

    def add_progress_bar(self, label_text: str, percentage: float, x: int, y: int, width: int = 1600):
        """
        添加进度条

        Args:
            label_text: 标签文本
            percentage: 百分比(0-1之间的浮点数)
            x: X坐标(px)
            y: Y坐标(px)
            width: 宽度(px)
        """
        # 添加标签
        label_left = UnitConverter.px_to_emu(x)
        label_top = UnitConverter.px_to_emu(y)
        label_width = UnitConverter.px_to_emu(width - 100)
        label_height = UnitConverter.px_to_emu(25)

        label_box = self.slide.shapes.add_textbox(
            label_left, label_top, label_width, label_height
        )
        label_frame = label_box.text_frame
        label_frame.text = label_text
        label_frame.word_wrap = False
        label_frame.margin_top = 0
        label_frame.margin_bottom = 0

        for paragraph in label_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(16)
                run.font.name = 'Microsoft YaHei'

        # 添加百分比文本
        percent_text = f"{percentage * 100:.1f}%"
        percent_left = UnitConverter.px_to_emu(x + width - 80)
        percent_box = self.slide.shapes.add_textbox(
            percent_left, label_top, UnitConverter.px_to_emu(80), label_height
        )
        percent_frame = percent_box.text_frame
        percent_frame.text = percent_text

        for paragraph in percent_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(16)
                run.font.name = 'Microsoft YaHei'

        # 添加进度条背景
        bar_top = UnitConverter.px_to_emu(y + 30)
        bar_left = UnitConverter.px_to_emu(x)
        bar_width = UnitConverter.px_to_emu(width)
        bar_height = UnitConverter.px_to_emu(16)

        bg_shape = self.slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            bar_left, bar_top, bar_width, bar_height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = ColorParser.parse_color('#e2e8f0')
        bg_shape.line.fill.background()

        # 添加进度条填充
        fill_width = UnitConverter.px_to_emu(int(width * percentage))
        if fill_width > 0:
            fill_shape = self.slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                bar_left, bar_top, fill_width, bar_height
            )
            fill_shape.fill.solid()
            fill_shape.fill.fore_color.rgb = ColorParser.get_primary_color()
            fill_shape.line.fill.background()

        logger.info(f"添加进度条: {label_text} - {percent_text}")

    def add_page_number(self, page_num: str):
        """
        添加页码

        Args:
            page_num: 页码文本
        """
        left = UnitConverter.px_to_emu(1920 - 100)
        top = UnitConverter.px_to_emu(1030)
        width = UnitConverter.px_to_emu(50)
        height = UnitConverter.px_to_emu(30)

        page_box = self.slide.shapes.add_textbox(left, top, width, height)
        page_frame = page_box.text_frame
        page_frame.text = page_num

        for paragraph in page_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(14)
                run.font.color.rgb = ColorParser.parse_color('#666')
                run.font.name = 'Microsoft YaHei'

        logger.info(f"添加页码: {page_num}")

    def add_stat_box_background(self, x: int, y: int, width: int, height: int):
        """
        添加统计卡片背景

        Args:
            x, y: 坐标(px)
            width, height: 尺寸(px)
        """
        left = UnitConverter.px_to_emu(x)
        top = UnitConverter.px_to_emu(y)
        w = UnitConverter.px_to_emu(width)
        h = UnitConverter.px_to_emu(height)

        shape = self.slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, w, h
        )

        # 从CSS获取stat-box背景颜色
        bg_color_str = self.css_parser.get_background_color('.stat-box')
        if bg_color_str:
            bg_rgb, alpha = ColorParser.parse_rgba(bg_color_str)
            if bg_rgb:
                # 如果有透明度，与白色混合
                if alpha < 1.0:
                    bg_rgb = ColorParser.blend_with_white(bg_rgb, alpha)
                shape.fill.solid()
                shape.fill.fore_color.rgb = bg_rgb
        else:
            # 降级：使用默认颜色 rgba(10, 66, 117, 0.06)
            bg_color = ColorParser.blend_with_white(ColorParser.get_primary_color(), 0.06)
            shape.fill.solid()
            shape.fill.fore_color.rgb = bg_color

        shape.line.fill.background()
        shape.shadow.inherit = False  # 无阴影

    def add_border_left(self, x: int, y: int, height: int, width: int = 4):
        """
        添加左边框

        Args:
            x, y: 坐标(px)
            height: 高度(px)
            width: 边框宽度(px)
        """
        left = UnitConverter.px_to_emu(x)
        top = UnitConverter.px_to_emu(y)
        w = UnitConverter.px_to_emu(width)
        h = UnitConverter.px_to_emu(height)

        shape = self.slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, w, h
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = ColorParser.get_primary_color()
        shape.line.fill.background()

    def convert(self, element, **kwargs):
        """转换形状元素"""
        pass
