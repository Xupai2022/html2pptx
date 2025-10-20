"""
文本转换器
处理H1, H2, P等文本元素
"""

from pptx.util import Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from src.converters.base_converter import BaseConverter
from src.mapper.style_mapper import StyleMapper
from src.utils.unit_converter import UnitConverter
from src.utils.color_parser import ColorParser
from src.utils.logger import setup_logger
from src.utils.font_manager import get_font_manager
from src.utils.style_computer import get_style_computer

logger = setup_logger(__name__)


class TextConverter(BaseConverter):
    """文本元素转换器"""

    def convert_title(self, title_text: str, subtitle_text: str = None, x: int = 80, y: int = 20,
                      is_cover: bool = False, title_classes: list = None, h1_element=None) -> int:
        """
        转换标题和副标题

        Tailwind CSS v2.2.19 规则 (1rem = 16px):
        - mt-10 = 10 × 0.25rem = 2.5rem = 40px
        - mt-32 = 32 × 0.25rem = 8rem = 128px (封面页用)
        - mt-2 = 2 × 0.25rem = 0.5rem = 8px
        - mb-4 = 4 × 0.25rem = 1rem = 16px
        - mb-16 = 16 × 0.25rem = 4rem = 64px (封面页用)
        - h-1 = 1 × 0.25rem = 0.25rem = 4px
        - line-height: 1.5 (Tailwind默认)

        文本高度计算:
        - h1 (font-size: 48px): 48px × 1.5 = 72px
        - h1 cover-title (font-size: 56px): 56px × 1.5 = 84px
        - h2 (font-size: 36px): 36px × 1.5 = 54px

        普通页面布局计算 (从content-section padding-top开始):
        初始y = 20px (content-section padding-top)
        + mt-10: 40px → y = 60px
        + h1: 72px → y = 132px
        + mt-2: 8px → y = 140px
        + h2: 54px → y = 194px
        + line: 4px → y = 198px (装饰线紧接h2)
        + mb-4: 16px → y = 214px (装饰线的下边距)
        标题区域结束: y = 214px

        封面页布局计算 (居中对齐):
        初始y = 20px (content-section padding-top)
        + mt-32: 128px → y = 148px
        + cover-title h1: 84px → y = 232px
        + cover-title h2: 84px → y = 316px
        + mb-16: 64px → y = 380px
        标题区域结束: y = 380px

        Args:
            title_text: 标题文本
            subtitle_text: 副标题文本
            x: X坐标(px)
            y: Y坐标(px) - content-section的padding-top起始位置，默认20px
            is_cover: 是否为封面页
            title_classes: 标题的CSS类列表
            h1_element: H1元素对象（用于获取样式）

        Returns:
            标题区域结束后的Y坐标(px)
        """
        # 初始y值应该是content-section的padding-top (20px)
        current_y = y

        # 获取样式计算器和字体
        style_computer = get_style_computer(self.css_parser)
        font_manager = get_font_manager(self.css_parser)

        # 根据是否为封面页调整边距和对齐方式
        if is_cover:
            # 封面页：mt-32 = 8rem = 128px
            current_y += 128
            # 居中对齐：幻灯片宽度1920px，减去左右padding各80px，内容区1760px
            # 标题框宽度为1760px，居中显示在幻灯片上
            left = UnitConverter.px_to_emu(80)  # 左边距
            width = UnitConverter.px_to_emu(1760)  # 内容区宽度
        else:
            # 普通页面：mt-10 = 2.5rem = 40px
            current_y += 40
            left = UnitConverter.px_to_emu(x)
            width = UnitConverter.px_to_emu(1760)

        # 使用传入的h1_element或创建临时元素
        if h1_element is None:
            from bs4 import BeautifulSoup
            temp_soup = BeautifulSoup('<h1>' + title_text + '</h1>', 'html.parser')
            h1_element = temp_soup.h1

        # 获取h1的字体大小 (现在get_font_size_pt返回pt值)
        h1_font_size_pt = style_computer.get_font_size_pt(h1_element)
        # 转换回px用于高度计算
        h1_font_size_px = UnitConverter.pt_to_px(h1_font_size_pt)
        h1_height = int(h1_font_size_px * 1.5)  # 行高1.5

        logger.debug(f"H1标题字体大小: {h1_font_size_px}px → {h1_font_size_pt}pt, 高度: {h1_height}px")

        # 获取h1的颜色样式
        h1_style = style_computer.compute_computed_style(h1_element)
        h1_inline_style = self._extract_inline_style(h1_element)
        h1_color_str = h1_style.get('color') or h1_inline_style.get('color')

        # 添加h1标题
        top = UnitConverter.px_to_emu(current_y)
        height = UnitConverter.px_to_emu(h1_height)

        title_box = self.slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_frame.word_wrap = True
        title_frame.margin_top = 0
        title_frame.margin_bottom = 0
        title_frame.margin_left = 0
        title_frame.margin_right = 0

        # 设置字体
        font_name = font_manager.get_font('h1')

        for paragraph in title_frame.paragraphs:
            # 封面页标题居中对齐
            if is_cover:
                paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            for run in paragraph.runs:
                # 使用转换后的pt值设置字体大小
                run.font.size = Pt(h1_font_size_pt)
                run.font.bold = True
                run.font.name = font_name

                # 应用颜色：优先使用HTML中定义的颜色，否则根据页面类型使用默认颜色
                if h1_color_str:
                    # HTML中定义了颜色
                    color = ColorParser.parse_color(h1_color_str)
                    if color:
                        run.font.color.rgb = color
                        logger.debug(f"应用H1自定义颜色: {h1_color_str}")
                elif is_cover:
                    # 封面页使用主题色
                    run.font.color.rgb = ColorParser.get_primary_color()
                    logger.debug("封面页H1使用主题色")
                else:
                    # 普通页面也使用主题色（保持与HTML一致）
                    run.font.color.rgb = ColorParser.get_primary_color()
                    logger.debug("普通页面H1使用主题色")

        logger.info(f"添加标题: {title_text} ({'封面页' if is_cover else '普通页面'})")

        current_y += h1_height  # 72px 或 84px

        # 添加副标题
        if subtitle_text:
            # mt-2: 0.5rem = 8px
            current_y += 8

            # 创建临时h2元素来获取字体大小
            from bs4 import BeautifulSoup
            temp_soup_h2 = BeautifulSoup('<h2>' + subtitle_text + '</h2>', 'html.parser')
            h2_element = temp_soup_h2.h2

            # 获取h2的字体大小 (现在get_font_size_pt返回pt值)
            h2_font_size_pt = style_computer.get_font_size_pt(h2_element)
            # 转换回px用于高度计算
            h2_font_size_px = UnitConverter.pt_to_px(h2_font_size_pt)
            h2_height = int(h2_font_size_px * 1.5)  # 行高1.5

            logger.debug(f"H2副标题字体大小: {h2_font_size_px}px → {h2_font_size_pt}pt, 高度: {h2_height}px")

            # 获取h2的颜色样式
            h2_style = style_computer.compute_computed_style(h2_element)
            h2_inline_style = self._extract_inline_style(h2_element)
            h2_color_str = h2_style.get('color') or h2_inline_style.get('color')

            subtitle_top = UnitConverter.px_to_emu(current_y)
            subtitle_box = self.slide.shapes.add_textbox(
                left, subtitle_top, width, UnitConverter.px_to_emu(h2_height)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle_text
            subtitle_frame.word_wrap = True
            subtitle_frame.margin_top = 0
            subtitle_frame.margin_bottom = 0
            subtitle_frame.margin_left = 0
            subtitle_frame.margin_right = 0

            font_name_h2 = font_manager.get_font('h2')

            for paragraph in subtitle_frame.paragraphs:
                # 封面页副标题也居中对齐
                if is_cover:
                    paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                for run in paragraph.runs:
                    # 使用转换后的pt值设置字体大小
                    run.font.size = Pt(h2_font_size_pt)
                    run.font.bold = True
                    run.font.name = font_name_h2

                    # 应用颜色：优先使用HTML中定义的颜色，否则使用主题色
                    if h2_color_str:
                        # HTML中定义了颜色
                        color = ColorParser.parse_color(h2_color_str)
                        if color:
                            run.font.color.rgb = color
                            logger.debug(f"应用H2自定义颜色: {h2_color_str}")
                    else:
                        # 使用主题色
                        run.font.color.rgb = ColorParser.get_primary_color()
                        logger.debug("H2使用主题色")

            logger.info(f"添加副标题: {subtitle_text}")

            current_y += h2_height  # 54px

            # 只有非封面页才添加装饰线
            if not is_cover:
                # 添加装饰线 (w-20 h-1) - 紧接h2，无间距
                line_top = UnitConverter.px_to_emu(current_y)
                line_left = UnitConverter.px_to_emu(x)
                line_width = UnitConverter.px_to_emu(80)  # w-20 = 5rem = 80px
                # h-1: 0.25rem = 4px
                line_height = UnitConverter.px_to_emu(4)

                line_shape = self.slide.shapes.add_shape(
                    1,  # Rectangle
                    line_left, line_top, line_width, line_height
                )
                line_shape.fill.solid()
                line_shape.fill.fore_color.rgb = ColorParser.get_primary_color()
                line_shape.line.fill.background()

                current_y += 4  # 装饰线高度

                # mb-4: 1rem = 16px (装饰线的下边距)
                current_y += 16
            else:
                # 封面页：智能计算装饰线宽度
                # 获取标题和副标题的最大宽度作为参考

                # 计算标题文本的估算宽度（中文字符数 * 字体大小系数）
                title_chars = len(title_text)
                # 对于56px的字体，中文字符宽度约为字体大小的0.8倍
                title_width_px = int(title_chars * h1_font_size_px * 0.8)

                # 如果有副标题，也计算其宽度
                subtitle_width_px = 0
                if subtitle_text:
                    subtitle_chars = len(subtitle_text)
                    # 副标题也是56px（cover-title），使用相同的系数
                    subtitle_width_px = int(subtitle_chars * h2_font_size_px * 0.8)

                # 取标题和副标题的最大宽度
                max_text_width = max(title_width_px, subtitle_width_px)

                # 装饰线宽度策略：
                # - 装饰线宽度约为文本宽度的60-80%，但不能太长或太短
                # - 最小宽度：w-24 (96px)
                # - 最大宽度：w-64 (256px)
                # - 基础宽度：文本宽度的70%
                base_line_width = int(max_text_width * 0.7)

                # 限制在合理范围内
                line_width = max(96, min(256, base_line_width))

                # 确保不超过内容区域宽度的一半
                max_line_width = 1760 // 2
                line_width = min(line_width, max_line_width)

                # 添加装饰线
                line_top = UnitConverter.px_to_emu(current_y)
                # 计算居中位置：
                # 标题框从80px开始，宽度1760px，文本居中对齐
                # 文本中心在 80 + 1760/2 = 960px（幻灯片正中心）
                # 装饰线应该以960px为中心对齐
                # 装饰线左边界 = 960 - line_width/2
                line_left = UnitConverter.px_to_emu(960 - line_width // 2)
                line_width_emu = UnitConverter.px_to_emu(line_width)
                # h-1: 0.25rem = 4px
                line_height = UnitConverter.px_to_emu(4)

                line_shape = self.slide.shapes.add_shape(
                    1,  # Rectangle
                    line_left, line_top, line_width_emu, line_height
                )
                line_shape.fill.solid()
                line_shape.fill.fore_color.rgb = ColorParser.get_primary_color()
                line_shape.line.fill.background()

                logger.info(f"添加封面页装饰线: 宽度={line_width}px, 标题宽度={max_text_width}px")

                current_y += 4  # 装饰线高度

                # mb-16: 4rem = 64px (封面页装饰线的下边距)
                current_y += 64

        # 标题区域结束位置
        # 普通页面: y(20) + mt-10(40) + h1(72) + mt-2(8) + h2(54) + line(4) + mb-4(16) = 214px
        # 封面页: y(20) + mt-32(128) + h1(84) + mt-2(8) + h2(84) + line(4) + mb-16(64) = 392px
        logger.info(f"标题区域结束位置: y={current_y}px ({'封面页' if is_cover else '普通页面'})")
        return current_y

    def convert_paragraph(self, p_element, x: int, y: int, width: int = 1760):
        """
        转换段落

        Args:
            p_element: 段落HTML元素
            x: X坐标(px)
            y: Y坐标(px)
            width: 宽度(px)
        """
        text = p_element.get_text(strip=True)
        if not text:
            return

        # 获取样式计算器和字体
        style_computer = get_style_computer(self.css_parser)
        font_manager = get_font_manager(self.css_parser)

        # 获取段落的字体大小 (现在get_font_size_pt返回pt值)
        p_font_size_pt = style_computer.get_font_size_pt(p_element)
        # 转换回px用于高度计算
        p_font_size_px = UnitConverter.pt_to_px(p_font_size_pt)
        p_height = int(p_font_size_px * 1.5)  # 行高1.5

        logger.debug(f"段落字体大小: {p_font_size_px}px → {p_font_size_pt}pt, 高度: {p_height}px")

        left = UnitConverter.px_to_emu(x)
        top = UnitConverter.px_to_emu(y)
        w = UnitConverter.px_to_emu(width)
        h = UnitConverter.px_to_emu(p_height)

        text_box = self.slide.shapes.add_textbox(left, top, w, h)
        text_frame = text_box.text_frame
        text_frame.text = text
        text_frame.word_wrap = True
        text_frame.margin_top = 0
        text_frame.margin_bottom = 0

        # 应用样式
        p_style = style_computer.compute_computed_style(p_element)
        inline_style = self._extract_inline_style(p_element)
        font_name = font_manager.get_font('p', inline_style)

        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                # 使用转换后的pt值设置字体大小
                run.font.size = Pt(p_font_size_pt)
                run.font.name = font_name

                # 颜色
                color_str = p_style.get('color') or inline_style.get('color')
                if color_str:
                    color = ColorParser.parse_color(color_str)
                    if color:
                        run.font.color.rgb = color

                # 字体粗细
                font_weight = p_style.get('font-weight')
                if font_weight:
                    run.font.bold = StyleMapper.parse_font_weight(font_weight)

    def _extract_inline_style(self, element) -> dict:
        """提取内联样式"""
        style_dict = {}
        style_str = element.get('style', '')

        if not style_str:
            # 检查class属性
            classes = element.get('class', [])
            if 'primary-color' in classes:
                style_dict['color'] = 'rgb(10, 66, 117)'
            # 对于h1和h2标签，应用默认的主题色
            elif element.name in ['h1', 'h2']:
                style_dict['color'] = 'rgb(10, 66, 117)'

        return style_dict

    def convert_numbered_list(self, numbered_item: dict, x: int, y: int, width: int = 1760) -> int:
        """
        转换数字列表项

        Args:
            numbered_item: 数字列表项信息
            x: X坐标(px)
            y: Y坐标(px)
            width: 宽度(px)

        Returns:
            下一项的Y坐标(px)
        """
        style_computer = get_style_computer(self.css_parser)
        font_manager = get_font_manager(self.css_parser)

        # 获取基础字体大小（使用p标签作为参考）
        from bs4 import BeautifulSoup
        temp_soup = BeautifulSoup('<p>Temp</p>', 'html.parser')
        p_element = temp_soup.p
        p_font_size_pt = style_computer.get_font_size_pt(p_element)
        p_font_size_px = UnitConverter.pt_to_px(p_font_size_pt)
        line_height = int(p_font_size_px * 1.6)  # 使用1.6行高，与HTML一致

        # 获取数字样式
        number_style = style_computer.compute_computed_style(numbered_item['number_elem'])
        number_inline = self._extract_inline_style(numbered_item['number_elem'])
        number_font_name = font_manager.get_font('p', number_inline)

        # 获取文本样式
        text_style = style_computer.compute_computed_style(numbered_item['text_elem'])
        text_inline = self._extract_inline_style(numbered_item['text_elem'])
        text_font_name = font_manager.get_font('p', text_inline)

        # 根据类型调整布局
        if numbered_item['type'] == 'toc':
            # TOC格式：数字和文本水平排列
            number_width = 60  # 数字区域宽度
            text_width = width - number_width - 20  # 文本区域宽度，留20px间距

            # 添加数字
            number_left = UnitConverter.px_to_emu(x)
            number_top = UnitConverter.px_to_emu(y)
            number_w = UnitConverter.px_to_emu(number_width)
            number_h = UnitConverter.px_to_emu(line_height)

            number_box = self.slide.shapes.add_textbox(number_left, number_top, number_w, number_h)
            number_frame = number_box.text_frame
            number_frame.text = numbered_item['number']
            number_frame.margin_top = 0
            number_frame.margin_bottom = 0
            number_frame.margin_left = 0

            # 设置数字样式
            for paragraph in number_frame.paragraphs:
                paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT  # 右对齐
                for run in paragraph.runs:
                    run.font.size = Pt(p_font_size_pt)
                    run.font.bold = True
                    run.font.name = number_font_name
                    # 应用数字颜色（通常是主题色）
                    run.font.color.rgb = ColorParser.get_primary_color()

            # 添加文本
            text_left = UnitConverter.px_to_emu(x + number_width + 20)
            text_top = UnitConverter.px_to_emu(y)
            text_w = UnitConverter.px_to_emu(text_width)
            text_h = UnitConverter.px_to_emu(line_height)

            text_box = self.slide.shapes.add_textbox(text_left, text_top, text_w, text_h)
            text_frame = text_box.text_frame
            text_frame.text = numbered_item['text']
            text_frame.margin_top = 0
            text_frame.margin_bottom = 0
            text_frame.margin_left = 0

            # 设置文本样式
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(p_font_size_pt)
                    run.font.name = text_font_name

                    # 应用文本颜色
                    color_str = text_style.get('color') or text_inline.get('color')
                    if color_str:
                        color = ColorParser.parse_color(color_str)
                        if color:
                            run.font.color.rgb = color

            logger.info(f"添加目录项: {numbered_item['number']} - {numbered_item['text']}")

        else:
            # 其他格式：数字和文本在同一个文本框中
            text_left = UnitConverter.px_to_emu(x)
            text_top = UnitConverter.px_to_emu(y)
            text_w = UnitConverter.px_to_emu(width)
            text_h = UnitConverter.px_to_emu(line_height)

            text_box = self.slide.shapes.add_textbox(text_left, text_top, text_w, text_h)
            text_frame = text_box.text_frame
            text_frame.margin_top = 0
            text_frame.margin_bottom = 0
            text_frame.margin_left = 0

            # 清除默认段落
            text_frame.clear()

            # 创建段落
            p = text_frame.paragraphs[0]

            # 添加数字部分
            number_run = p.add_run()
            number_run.text = numbered_item['number']
            number_run.font.size = Pt(p_font_size_pt)
            number_run.font.bold = True
            number_run.font.name = number_font_name
            number_run.font.color.rgb = ColorParser.get_primary_color()

            # 添加分隔符
            sep_run = p.add_run()
            sep_run.text = ". " if numbered_item['type'] in ['ordered_list', 'paragraph_numbered'] else " "
            sep_run.font.size = Pt(p_font_size_pt)
            sep_run.font.name = number_font_name

            # 添加文本部分
            text_run = p.add_run()
            text_run.text = numbered_item['text']
            text_run.font.size = Pt(p_font_size_pt)
            text_run.font.name = text_font_name

            # 应用文本颜色
            color_str = text_style.get('color') or text_inline.get('color')
            if color_str:
                color = ColorParser.parse_color(color_str)
                if color:
                    text_run.font.color.rgb = color

            logger.info(f"添加数字列表项: {numbered_item['number']} - {numbered_item['text']}")

        # 返回下一行的Y坐标（添加项目间距）
        item_spacing = 18 if numbered_item['type'] == 'toc' else 10
        return y + line_height + item_spacing

    def convert(self, element, **kwargs):
        """转换文本元素"""
        # 注意：这个方法在main.py中用于单个元素的转换
        # 但是convert_title会创建新的文本框，可能导致重复文字
        # 为了避免重复，这里应该使用convert_paragraph方法处理所有文本元素

        tag_name = element.name
        if tag_name in ['h1', 'h2', 'h3', 'p']:
            # 统一使用段落转换方法，传入实际元素
            x = kwargs.get('x', 80)
            y = kwargs.get('y', 0)
            width = kwargs.get('width', 1760)
            self.convert_paragraph(element, x, y, width)
