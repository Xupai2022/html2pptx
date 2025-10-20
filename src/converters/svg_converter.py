"""
SVG转换器
将SVG图表转换为PPTX内容，支持截图和内容提取两种方式
"""

from typing import Optional, Dict, List, Tuple
from pathlib import Path
import re

try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Pt

from src.converters.base_converter import BaseConverter
from src.utils.chart_capture import ChartCapture
from src.utils.unit_converter import UnitConverter
from src.utils.logger import setup_logger

logger = setup_logger(__name__)


class SvgConverter(BaseConverter):
    """SVG图表转换器"""

    def __init__(self, slide, css_parser, html_path: str = None, use_stable_chart_capture: bool = False):
        """
        初始化SVG转换器

        Args:
            slide: PPTX幻灯片对象
            css_parser: CSS解析器
            html_path: HTML文件路径（用于截图）
            use_stable_chart_capture: 是否使用稳定的截图版本
        """
        super().__init__(slide, css_parser)
        self.html_path = html_path
        self.generated_png_files = []  # 记录生成的PNG文件路径

        if use_stable_chart_capture:
            from src.utils.chart_capture_working import get_working_capturer
            self.chart_capturer = get_working_capturer()
            logger.info("使用工作的图表截图版本")
        else:
            self.chart_capturer = ChartCapture()
            logger.info("使用标准的图表截图版本")

    def convert(self, element, **kwargs):
        """
        转换SVG元素（实现BaseConverter的抽象方法）

        Args:
            element: SVG元素
            **kwargs: 其他参数

        Returns:
            转换结果
        """
        # 获取参数
        x = kwargs.get('x', 80)
        y = kwargs.get('y', 100)
        width = kwargs.get('width', 1760)
        chart_index = kwargs.get('chart_index', 0)

        # 获取容器
        container = kwargs.get('container', element.parent)

        # 调用SVG转换方法
        return self.convert_svg(element, container, x, y, width, chart_index)

    def convert_svg(
        self,
        svg_element,
        container,
        x: int,
        y: int,
        width: int,
        chart_index: int = 0
    ) -> int:
        """
        转换SVG图表

        Args:
            svg_element: SVG元素
            container: 容器元素
            x: X坐标
            y: Y坐标
            width: 宽度
            chart_index: 图表索引（用于多个图表时的截图）

        Returns:
            实际高度
        """
        logger.info(f"开始转换SVG图表 {chart_index}")

        # 获取SVG的原始尺寸
        svg_width, svg_height = self._get_svg_dimensions(svg_element)

        # 计算目标尺寸，保持宽高比
        target_width = width
        target_height = int(target_width * svg_height / svg_width) if svg_width > 0 else 250

        logger.info(f"SVG原始尺寸: {svg_width}x{svg_height}, 目标尺寸: {target_width}x{target_height}")

        # 尝试截图（优先方案）
        actual_height = target_height
        screenshot_success = False

        if self.html_path:
            screenshot_path = self._capture_svg_screenshot(svg_element, chart_index, container)
            if screenshot_path:
                actual_height = self._insert_svg_screenshot(
                    self.slide, screenshot_path, x, y, target_width, target_height
                )
                screenshot_success = True
                logger.info(f"SVG图表 {chart_index} 截图成功，尺寸: {target_width}x{actual_height}")

        # 不使用降级渲染，截图失败则跳过
        if not screenshot_success:
            logger.error(f"SVG图表 {chart_index} 截图失败，跳过")
            return 0

        return actual_height

    def _get_svg_dimensions(self, svg_element) -> Tuple[int, int]:
        """
        获取SVG的尺寸

        Args:
            svg_element: SVG元素

        Returns:
            (宽度, 高度)
        """
        # 优先从width和height属性获取
        width_str = svg_element.get('width', '')
        height_str = svg_element.get('height', '')

        # 如果有明确的width和height属性，使用它们
        if width_str and height_str:
            width = self._parse_dimension(width_str)
            height = self._parse_dimension(height_str)
            if width > 0 and height > 0:
                logger.info(f"使用SVG属性尺寸: {width}x{height}")
                return width, height

        # 如果没有width/height或解析失败，尝试从viewBox获取
        viewbox = svg_element.get('viewBox')
        if viewbox:
            try:
                values = viewbox.split()
                if len(values) >= 4:
                    # viewBox格式: min-x min-y width height
                    vb_width = float(values[2])
                    vb_height = float(values[3])

                    # 如果viewBox的宽高比合理，直接使用
                    if vb_width > 0 and vb_height > 0:
                        # 对于没有明确尺寸的SVG，使用合理的默认尺寸
                        # 但保持viewBox的宽高比
                        if vb_width > vb_height:
                            # 横向SVG，最大宽度400
                            width = 400
                            height = int(400 * vb_height / vb_width)
                        elif vb_height > vb_width:
                            # 纵向SVG，最大高度300
                            height = 300
                            width = int(300 * vb_width / vb_height)
                        else:
                            # 正方形
                            size = min(400, int(vb_width))
                            width = height = size

                        logger.info(f"从viewBox推导尺寸: {width}x{height} (viewBox: {vb_width}x{vb_height})")
                        return width, height
            except (ValueError, IndexError) as e:
                logger.warning(f"解析viewBox失败: {e}")

        # 最后使用默认值
        logger.warning("无法获取SVG尺寸，使用默认值400x250")
        return 400, 250

    def _parse_dimension(self, dimension_str: str) -> int:
        """
        解析尺寸字符串（处理px、em等单位）

        Args:
            dimension_str: 尺寸字符串

        Returns:
            像素值
        """
        if not dimension_str:
            return 0

        # 移除空格并转为小写
        dimension_str = dimension_str.strip().lower()

        # 提取数字部分
        match = re.search(r'([\d.]+)', dimension_str)
        if match:
            value = float(match.group(1))

            # 处理单位
            if dimension_str.endswith('px'):
                return int(value)
            elif dimension_str.endswith('em'):
                return int(value * 16)  # 假设1em=16px
            elif dimension_str.endswith('rem'):
                return int(value * 16)  # 假设1rem=16px
            elif dimension_str.endswith('%'):
                # 百分比需要基于父容器，这里返回0让调用方处理
                logger.warning(f"SVG尺寸不支持百分比单位: {dimension_str}")
                return 0
            elif dimension_str.endswith('pt'):
                return int(value * 96 / 72)  # 1pt = 96/72px
            elif dimension_str.endswith('pc'):
                return int(value * 16)  # 1pc = 16px
            elif dimension_str.endswith('in'):
                return int(value * 96)  # 1in = 96px
            elif dimension_str.endswith('cm'):
                return int(value * 96 / 2.54)  # 1cm = 96/2.54px
            elif dimension_str.endswith('mm'):
                return int(value * 96 / 25.4)  # 1mm = 96/25.4px
            else:
                # 无单位，假设为像素
                return int(value)

        return 0

    def _generate_svg_signature(self, svg_element) -> str:
        """
        生成SVG的唯一签名，用于缓存和识别

        Args:
            svg_element: SVG元素

        Returns:
            SVG签名字符串
        """
        import hashlib

        # 收集SVG的关键特征
        features = []

        # 1. viewBox
        viewbox = svg_element.get('viewBox', '')
        if viewbox:
            features.append(f"vb:{viewbox}")

        # 2. width和height
        width = svg_element.get('width', '')
        height = svg_element.get('height', '')
        if width:
            features.append(f"w:{width}")
        if height:
            features.append(f"h:{height}")

        # 3. 子元素统计
        child_counts = {}
        for child in svg_element.children:
            if hasattr(child, 'name') and child.name:
                child_counts[child.name] = child_counts.get(child.name, 0) + 1

        # 按元素名排序并添加到特征
        for elem_name in sorted(child_counts.keys()):
            features.append(f"{elem_name}:{child_counts[elem_name]}")

        # 4. 文本内容（如果有）
        texts = svg_element.find_all('text')
        if texts:
            text_content = ' '.join([t.get_text(strip=True)[:50] for t in texts])
            if text_content:
                features.append(f"text:{text_content}")

        # 5. 特殊属性（如class, id等）
        for attr in ['class', 'id']:
            val = svg_element.get(attr, '')
            if val:
                features.append(f"{attr}:{val}")

        # 生成签名
        signature_str = '|'.join(features)
        signature_hash = hashlib.md5(signature_str.encode()).hexdigest()[:12]

        logger.debug(f"SVG特征: {signature_str}")
        logger.debug(f"SVG签名哈希: {signature_hash}")

        return signature_hash

    def _is_same_svg(self, svg1, svg2) -> bool:
        """
        比较两个SVG元素是否相同（基于内容特征）

        Args:
            svg1: 第一个SVG元素
            svg2: 第二个SVG元素

        Returns:
            是否为同一个SVG
        """
        # 比较viewBox
        vb1 = svg1.get('viewBox', '')
        vb2 = svg2.get('viewBox', '')
        if vb1 != vb2:
            return False

        # 比较width和height
        w1 = svg1.get('width', '')
        w2 = svg2.get('width', '')
        h1 = svg1.get('height', '')
        h2 = svg2.get('height', '')
        if w1 != w2 or h1 != h2:
            return False

        # 比较主要子元素的数量和类型
        children1 = []
        for child in svg1.children:
            if hasattr(child, 'name') and child.name:
                children1.append(child.name)

        children2 = []
        for child in svg2.children:
            if hasattr(child, 'name') and child.name:
                children2.append(child.name)

        # 简单的元素数量比较
        from collections import Counter
        return Counter(children1) == Counter(children2)

    def _capture_svg_screenshot(self, svg_element, chart_index: int, container) -> Optional[str]:
        """
        截取SVG截图（禁用缓存，确保获取当前HTML的实际图片）

        Args:
            svg_element: SVG元素
            chart_index: 图表索引（在容器中的索引）
            container: 容器元素

        Returns:
            截图路径，失败返回None
        """
        try:
            # 计算SVG在整个HTML中的实际索引
            # 找到所有在它之前的SVG元素
            from bs4 import BeautifulSoup
            import os

            # 读取HTML文件
            with open(self.html_path, 'r', encoding='utf-8') as f:
                html_content = f.read()

            soup = BeautifulSoup(html_content, 'html.parser')
            all_svgs = soup.find_all('svg')

            # 找到当前SVG在整个HTML中的索引
            actual_svg_index = None
            # 通过比较SVG的特征来找到正确的索引
            current_signature = self._generate_svg_signature(svg_element)

            for i, svg in enumerate(all_svgs):
                soup_signature = self._generate_svg_signature(svg)
                if soup_signature == current_signature:
                    actual_svg_index = i
                    break

            # 如果找不到匹配的，使用传入的索引
            if actual_svg_index is None:
                logger.warning(f"未找到匹配的SVG，使用传入索引: {chart_index}")
                actual_svg_index = chart_index

            logger.info(f"开始SVG截图，容器索引: {chart_index}, 实际HTML索引: {actual_svg_index}")

            # 生成唯一的输出路径（使用当前HTML文件名+时间戳）
            import time
            html_basename = os.path.basename(self.html_path)
            unique_id = f"{html_basename}_{actual_svg_index}_{int(time.time() * 1000)}"

            # 第1层：快速截图（等待0.5秒）
            logger.debug("尝试快速截图（500ms）")
            png_path = f"svg_screenshot_{unique_id}.png"
            self.generated_png_files.append(png_path)  # 记录文件路径
            result = self.chart_capturer.capture_svg_by_index(
                self.html_path,
                actual_svg_index,  # 使用实际索引
                output_path=png_path,  # 禁用缓存
                wait_time=500
            )
            if result:
                logger.info(f"快速截图成功: {result}")
                return result

            # 第2层：标准截图（等待1.5秒）
            logger.debug("尝试标准截图（1500ms）")
            png_path = f"svg_screenshot_{unique_id}_2.png"
            self.generated_png_files.append(png_path)  # 记录文件路径
            result = self.chart_capturer.capture_svg_by_index(
                self.html_path,
                actual_svg_index,  # 使用实际索引
                output_path=png_path,  # 禁用缓存
                wait_time=1500
            )
            if result:
                logger.info(f"标准截图成功: {result}")
                return result

            # 第3层：慢速截图（等待3秒）
            logger.debug("尝试慢速截图（3000ms）")
            png_path = f"svg_screenshot_{unique_id}_3.png"
            self.generated_png_files.append(png_path)  # 记录文件路径
            result = self.chart_capturer.capture_svg_by_index(
                self.html_path,
                actual_svg_index,  # 使用实际索引
                output_path=png_path,  # 禁用缓存
                wait_time=3000
            )
            if result:
                logger.info(f"慢速截图成功: {result}")
                return result

            # 所有截图方案都失败
            logger.error(f"SVG截图失败（索引: {chart_index}）")
            return None

        except Exception as e:
            logger.error(f"SVG截图异常: {e}")
            return None

    def _insert_svg_screenshot(
        self,
        pptx_slide,
        screenshot_path: str,
        x: int,
        y: int,
        width: int,
        height: int
    ) -> int:
        """
        插入SVG截图到PPTX，使用原始尺寸

        Args:
            pptx_slide: PPTX幻灯片
            screenshot_path: 截图路径
            x: X坐标
            y: Y坐标
            width: 期望宽度（忽略，使用截图原始尺寸）
            height: 期望高度（忽略，使用截图原始尺寸）

        Returns:
            实际高度
        """
        try:
            # 使用PIL库获取截图的实际尺寸
            if PIL_AVAILABLE:
                with Image.open(screenshot_path) as img:
                    actual_width, actual_height = img.size

                logger.info(f"截图实际尺寸: {actual_width}x{actual_height}px")
                logger.info(f"期望插入尺寸: {width}x{height}px（将被忽略）")

                # 使用截图的原始尺寸，不进行缩放
                final_width = actual_width
                final_height = actual_height

                logger.info(f"实际插入尺寸: {final_width}x{final_height}px（使用原始尺寸）")
            else:
                logger.warning("PIL库不可用，使用期望尺寸")
                final_width = width
                final_height = height

            # 添加图片，使用原始尺寸
            pic = pptx_slide.shapes.add_picture(
                screenshot_path,
                UnitConverter.px_to_emu(x),
                UnitConverter.px_to_emu(y),
                UnitConverter.px_to_emu(final_width),
                UnitConverter.px_to_emu(final_height)
            )

            return final_height

        except Exception as e:
            logger.error(f"插入SVG截图失败: {e}")
            # 不使用降级渲染，直接返回
            return 0

    def convert_multiple_svgs(
        self,
        container,
        x: int,
        y: int,
        total_width: int,
        gap: int = 24
    ) -> int:
        """
        转换容器中的多个SVG图表（水平布局）

        Args:
            container: 容器元素
            x: 起始X坐标
            y: Y坐标
            total_width: 总宽度
            gap: 图表间距

        Returns:
            实际高度
        """
        svg_elements = container.find_all('svg')
        num_svgs = len(svg_elements)

        if num_svgs == 0:
            logger.warning("容器中未找到SVG元素")
            return 0

        logger.info(f"找到 {num_svgs} 个SVG元素")

        # 计算每个SVG的宽度
        chart_width = (total_width - (num_svgs - 1) * gap) // num_svgs

        # 计算起始X位置（水平居中）
        total_charts_width = num_svgs * chart_width + (num_svgs - 1) * gap
        start_x = x + (total_width - total_charts_width) // 2

        current_y = y
        max_height = 0

        # 处理每个SVG图表
        for i, svg_elem in enumerate(svg_elements):
            chart_x = start_x + i * (chart_width + gap)

            # 转换SVG
            chart_height = self.convert_svg(
                svg_elem,
                container,
                chart_x,
                current_y,
                chart_width,
                i
            )

            # 更新最大高度
            max_height = max(max_height, chart_height)

        return max_height

    def cleanup_temp_files(self):
        """
        清理生成的临时PNG文件
        """
        import os
        for png_path in self.generated_png_files:
            try:
                if os.path.exists(png_path):
                    os.remove(png_path)
                    logger.info(f"已删除临时文件: {png_path}")
            except Exception as e:
                logger.warning(f"删除临时文件失败 {png_path}: {e}")
        self.generated_png_files.clear()

    def __del__(self):
        """
        析构函数，自动清理临时文件
        """
        self.cleanup_temp_files()


# 导入必要的RGBColor和其他常量
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_PARAGRAPH_ALIGNMENT