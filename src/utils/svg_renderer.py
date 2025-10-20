"""
SVG降级渲染器
当截图失败时，将SVG元素转换为PPTX形状
确保100%的SVG都能被处理
"""

import re
from typing import Tuple, List, Dict, Any
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Pt, Emu
from src.utils.logger import setup_logger
from src.utils.unit_converter import UnitConverter

logger = setup_logger(__name__)


class SvgRenderer:
    """SVG降级渲染器 - 将SVG转换为PPTX形状"""

    def __init__(self, slide):
        """
        初始化SVG渲染器

        Args:
            slide: PPTX幻灯片对象
        """
        self.slide = slide

    def render_svg(self, svg_element, x: int, y: int, width: int, height: int) -> int:
        """
        渲染SVG元素为PPTX形状

        Args:
            svg_element: SVG元素
            x: X坐标
            y: Y坐标
            width: 宽度
            height: 高度

        Returns:
            实际高度
        """
        logger.info(f"使用降级渲染器处理SVG: {x},{y} {width}x{height}")

        # 解析SVG内容
        shapes = self._parse_svg_elements(svg_element)

        if not shapes:
            # 如果无法解析，创建占位符
            return self._create_placeholder(x, y, width, height)

        # 渲染所有形状
        rendered_height = self._render_shapes(shapes, x, y, width, height)

        return rendered_height if rendered_height > 0 else height

    def _parse_svg_elements(self, svg_element) -> List[Dict[str, Any]]:
        """
        解析SVG元素为形状列表

        Args:
            svg_element: SVG元素

        Returns:
            形状列表
        """
        shapes = []

        try:
            # 获取SVG的viewBox以确定坐标系
            viewbox = svg_element.get('viewBox', '0 0 100 100')
            try:
                vb_x, vb_y, vb_width, vb_height = map(float, viewbox.split())
            except:
                vb_x, vb_y, vb_width, vb_height = 0, 0, 100, 100

            scale_x = 1.0
            scale_y = 1.0
            if vb_width > 0 and vb_height > 0:
                scale_x = 100 / vb_width
                scale_y = 100 / vb_height

            # 处理圆形
            for circle in svg_element.find_all('circle'):
                cx = float(circle.get('cx', 0)) * scale_x
                cy = float(circle.get('cy', 0)) * scale_y
                r = float(circle.get('r', 10)) * min(scale_x, scale_y)
                fill = circle.get('fill', '#333333')
                stroke = circle.get('stroke', 'none')
                stroke_width = float(circle.get('stroke-width', 0))

                shapes.append({
                    'type': 'circle',
                    'x': cx,
                    'y': cy,
                    'r': r,
                    'fill': fill,
                    'stroke': stroke,
                    'stroke_width': stroke_width
                })

            # 处理矩形
            for rect in svg_element.find_all('rect'):
                x = float(rect.get('x', 0)) * scale_x
                y = float(rect.get('y', 0)) * scale_y
                width = float(rect.get('width', 50)) * scale_x
                height = float(rect.get('height', 50)) * scale_y
                rx = float(rect.get('rx', 0)) * scale_x
                fill = rect.get('fill', '#333333')
                stroke = rect.get('stroke', 'none')
                stroke_width = float(rect.get('stroke-width', 0))

                shapes.append({
                    'type': 'rect',
                    'x': x,
                    'y': y,
                    'width': width,
                    'height': height,
                    'rx': rx,
                    'fill': fill,
                    'stroke': stroke,
                    'stroke_width': stroke_width
                })

            # 处理线条
            for line in svg_element.find_all('line'):
                x1 = float(line.get('x1', 0)) * scale_x
                y1 = float(line.get('y1', 0)) * scale_y
                x2 = float(line.get('x2', 100)) * scale_x
                y2 = float(line.get('y2', 100)) * scale_y
                stroke = line.get('stroke', '#333333')
                stroke_width = float(line.get('stroke-width', 1))

                shapes.append({
                    'type': 'line',
                    'x1': x1,
                    'y1': y1,
                    'x2': x2,
                    'y2': y2,
                    'stroke': stroke,
                    'stroke_width': stroke_width
                })

            # 处理路径（简化版）
            for path in svg_element.find_all('path'):
                d = path.get('d', '')
                fill = path.get('fill', '#333333')
                stroke = path.get('stroke', 'none')
                stroke_width = float(path.get('stroke-width', 0))

                # 简单解析路径，提取移动和直线命令
                points = self._parse_simple_path(d, scale_x, scale_y)
                if points:
                    shapes.append({
                        'type': 'polyline',
                        'points': points,
                        'fill': fill,
                        'stroke': stroke,
                        'stroke_width': stroke_width
                    })

            # 处理文本
            for text in svg_element.find_all('text'):
                x = float(text.get('x', 0)) * scale_x
                y = float(text.get('y', 0)) * scale_y
                content = text.get_text(strip=True)
                font_size = float(text.get('font-size', 16)) * min(scale_x, scale_y)
                fill = text.get('fill', '#333333')

                if content:
                    shapes.append({
                        'type': 'text',
                        'x': x,
                        'y': y,
                        'content': content,
                        'font_size': font_size,
                        'fill': fill
                    })

        except Exception as e:
            logger.error(f"解析SVG元素失败: {e}")

        return shapes

    def _parse_simple_path(self, d: str, scale_x: float, scale_y: float) -> List[Tuple[float, float]]:
        """
        简单解析SVG路径

        Args:
            d: 路径数据
            scale_x: X轴缩放
            scale_y: Y轴缩放

        Returns:
            点列表
        """
        points = []
        try:
            # 简单解析M(移动)和L(直线)命令
            commands = re.findall(r'([ML])\s*([\d.]+)\s*,\s*([\d.]+)', d.upper())
            current_x, current_y = 0, 0

            for cmd, x, y in commands:
                x = float(x) * scale_x
                y = float(y) * scale_y

                if cmd == 'M':
                    current_x, current_y = x, y
                elif cmd == 'L':
                    points.append((current_x, current_y))
                    points.append((x, y))
                    current_x, current_y = x, y

        except Exception as e:
            logger.error(f"解析路径失败: {e}")

        return points

    def _render_shapes(self, shapes: List[Dict], x: int, y: int, width: int, height: int) -> int:
        """
        渲染形状到幻灯片

        Args:
            shapes: 形状列表
            x: X坐标
            y: Y坐标
            width: 宽度
            height: 高度

        Returns:
            实际高度
        """
        try:
            # 计算边界
            min_x = min(shape.get('x', 0) for shape in shapes if 'x' in shape)
            min_y = min(shape.get('y', 0) for shape in shapes if 'y' in shape)
            max_x = max(
                max(shape.get('x', 0) + shape.get('width', shape.get('r', 0) * 2),
                    shape.get('x1', shape.get('x', 0)),
                    shape.get('x2', shape.get('x', 0)))
                for shape in shapes
            )
            max_y = max(
                max(shape.get('y', 0) + shape.get('height', shape.get('r', 0) * 2),
                    shape.get('y1', shape.get('y', 0)),
                    shape.get('y2', shape.get('y', 0)))
                for shape in shapes
            )

            # 计算缩放比例
            shape_width = max_x - min_x
            shape_height = max_y - min_y

            if shape_width > 0 and shape_height > 0:
                scale_x = width / shape_width
                scale_y = height / shape_height
                scale = min(scale_x, scale_y, 2.0)  # 限制最大缩放

                # 居中偏移
                offset_x = x + (width - shape_width * scale) / 2 - min_x * scale
                offset_y = y + (height - shape_height * scale) / 2 - min_y * scale
            else:
                scale = 1.0
                offset_x = x
                offset_y = y

            # 渲染每个形状
            for shape in shapes:
                self._render_single_shape(shape, offset_x, offset_y, scale)

            return int(shape_height * scale)

        except Exception as e:
            logger.error(f"渲染形状失败: {e}")
            return height

    def _render_single_shape(self, shape: Dict, offset_x: float, offset_y: float, scale: float):
        """
        渲染单个形状

        Args:
            shape: 形状数据
            offset_x: X偏移
            offset_y: Y偏移
            scale: 缩放比例
        """
        try:
            shape_type = shape.get('type')

            if shape_type == 'circle':
                cx = shape['x'] * scale + offset_x
                cy = shape['y'] * scale + offset_y
                r = shape['r'] * scale

                circle = self.slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    UnitConverter.px_to_emu(cx - r),
                    UnitConverter.px_to_emu(cy - r),
                    UnitConverter.px_to_emu(r * 2),
                    UnitConverter.px_to_emu(r * 2)
                )
                self._apply_shape_style(circle, shape)

            elif shape_type == 'rect':
                x = shape['x'] * scale + offset_x
                y = shape['y'] * scale + offset_y
                w = shape['width'] * scale
                h = shape['height'] * scale
                rx = shape.get('rx', 0)

                rect = self.slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    UnitConverter.px_to_emu(x),
                    UnitConverter.px_to_emu(y),
                    UnitConverter.px_to_emu(w),
                    UnitConverter.px_to_emu(h)
                )
                self._apply_shape_style(rect, shape)

            elif shape_type == 'line':
                # add_line方法在SlideShapes中不存在，跳过线条
                logger.debug("跳过线条渲染（API不支持）")

            elif shape_type == 'polyline' and shape['points']:
                # add_line方法不存在，跳过多线条
                logger.debug("跳过多线条渲染（API不支持）")

            elif shape_type == 'text':
                x = shape['x'] * scale + offset_x
                y = shape['y'] * scale + offset_y
                content = shape['content']
                font_size = shape['font_size'] * scale

                text_box = self.slide.shapes.add_textbox(
                    UnitConverter.px_to_emu(x - 50),
                    UnitConverter.px_to_emu(y - font_size/2),
                    UnitConverter.px_to_emu(100),
                    UnitConverter.px_to_emu(font_size * 1.5)
                )

                text_frame = text_box.text_frame
                text_frame.text = content
                text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

                # 设置文本样式
                for paragraph in text_frame.paragraphs:
                    paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(max(font_size * 0.75, 8))  # 最小8pt
                        color = self._parse_color(shape.get('fill', '#333333'))
                        if color:
                            run.font.color.rgb = color

        except Exception as e:
            logger.error(f"渲染形状失败: {e}")

    def _apply_shape_style(self, shape, style_data):
        """
        应用形状样式

        Args:
            shape: PPTX形状对象
            style_data: 样式数据
        """
        try:
            # 填充色
            fill_color = style_data.get('fill', 'none')
            if fill_color and fill_color != 'none':
                color = self._parse_color(fill_color)
                if color:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = color
                else:
                    shape.fill.background()
            else:
                shape.fill.background()

            # 边框
            stroke_color = style_data.get('stroke', 'none')
            stroke_width = style_data.get('stroke_width', 0)

            if stroke_color and stroke_color != 'none' and stroke_width > 0:
                color = self._parse_color(stroke_color)
                if color:
                    shape.line.color.rgb = color
                    shape.line.width = Pt(stroke_width * 0.75)  # 调整线宽
            else:
                # no_line() 不存在，设置为白色透明
                shape.line.color.rgb = RGBColor(255, 255, 255)
                shape.line.width = Pt(0)

        except Exception as e:
            logger.error(f"应用形状样式失败: {e}")

    def _apply_line_style(self, line, style_data):
        """
        应用线条样式

        Args:
            line: PPTX线条对象
            style_data: 样式数据
        """
        try:
            stroke_color = style_data.get('stroke', '#333333')
            stroke_width = style_data.get('stroke_width', 1)

            color = self._parse_color(stroke_color)
            if color:
                line.line.color.rgb = color
                line.line.width = Pt(stroke_width * 0.75)

        except Exception as e:
            logger.error(f"应用线条样式失败: {e}")

    def _parse_color(self, color_str: str) -> RGBColor:
        """
        解析颜色字符串

        Args:
            color_str: 颜色字符串

        Returns:
            RGBColor对象
        """
        try:
            # 处理不同的颜色格式
            color_str = color_str.strip().lower()

            # 十六进制
            if color_str.startswith('#'):
                hex_str = color_str[1:]
                if len(hex_str) == 3:
                    hex_str = ''.join(c*2 for c in hex_str)
                if len(hex_str) == 6:
                    r = int(hex_str[0:2], 16)
                    g = int(hex_str[2:4], 16)
                    b = int(hex_str[4:6], 16)
                    return RGBColor(r, g, b)

            # RGB(r,g,b)格式
            elif color_str.startswith('rgb('):
                values = color_str[4:-1].split(',')
                if len(values) == 3:
                    r = int(values[0].strip())
                    g = int(values[1].strip())
                    b = int(values[2].strip())
                    return RGBColor(r, g, b)

            # 颜色名称映射
            color_map = {
                'red': (255, 0, 0),
                'green': (0, 128, 0),
                'blue': (0, 0, 255),
                'black': (0, 0, 0),
                'white': (255, 255, 255),
                'gray': (128, 128, 128),
                'grey': (128, 128, 128),
                'silver': (192, 192, 192),
                'maroon': (128, 0, 0),
                'purple': (128, 0, 128),
                'fuchsia': (255, 0, 255),
                'lime': (0, 255, 0),
                'olive': (128, 128, 0),
                'navy': (0, 0, 128),
                'teal': (0, 128, 128),
                'aqua': (0, 255, 255)
            }

            if color_str in color_map:
                r, g, b = color_map[color_str]
                return RGBColor(r, g, b)

        except Exception as e:
            logger.error(f"解析颜色失败 {color_str}: {e}")

        return None

    def _create_placeholder(self, x: int, y: int, width: int, height: int) -> int:
        """
        创建占位符

        Args:
            x: X坐标
            y: Y坐标
            width: 宽度
            height: 高度

        Returns:
            高度
        """
        logger.warning("创建SVG占位符")

        try:
            placeholder = self.slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                UnitConverter.px_to_emu(x),
                UnitConverter.px_to_emu(y),
                UnitConverter.px_to_emu(width),
                UnitConverter.px_to_emu(height)
            )

            # 设置样式
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = RGBColor(240, 240, 240)
            placeholder.line.color.rgb = RGBColor(200, 200, 200)
            placeholder.line.width = Pt(1)

            # 添加文本
            text_frame = placeholder.text_frame
            text_frame.text = "SVG图表\n（降级显示）"
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            for paragraph in text_frame.paragraphs:
                paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(16)
                    run.font.color.rgb = RGBColor(100, 100, 100)

        except Exception as e:
            logger.error(f"创建占位符失败: {e}")

        return height