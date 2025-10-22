"""
验证所有问题是否修复的脚本
"""
from pptx import Presentation
from pptx.util import Pt, Inches
import os

def verify_slide_004(pptx_path):
    """验证slide_004: 第一个和第二个容器高度/间距"""
    print("\n=== 验证 slide_004 ===")
    prs = Presentation(pptx_path)
    slide = prs.slides[0]
    
    # 查找stat-card背景（grid中的）
    stat_card_bg = None
    for shape in slide.shapes:
        if hasattr(shape, 'fill') and shape.fill.type == 1:
            rgb = shape.fill.fore_color.rgb
            # stat-card背景色: rgba(10, 66, 117, 0.08)
            if 230 < rgb[0] < 245 and 235 < rgb[1] < 250 and 240 < rgb[2] < 255:
                top_px = int(shape.top / 9525)
                height_px = int(shape.height / 9525)
                if stat_card_bg is None or top_px < stat_card_bg['top']:
                    stat_card_bg = {
                        'top': top_px,
                        'height': height_px,
                        'bottom': top_px + height_px
                    }
    
    # 查找第一个data-card背景
    data_card_bg = None
    for shape in slide.shapes:
        if hasattr(shape, 'fill') and shape.fill.type == 1:
            rgb = shape.fill.fore_color.rgb
            # data-card背景色: rgba(10, 66, 117, 0.03)
            if 240 < rgb[0] < 255 and 245 < rgb[1] < 255 and 248 < rgb[2] < 255:
                top_px = int(shape.top / 9525)
                height_px = int(shape.height / 9525)
                if data_card_bg is None or top_px < data_card_bg['top']:
                    data_card_bg = {
                        'top': top_px,
                        'height': height_px,
                        'bottom': top_px + height_px
                    }
    
    if stat_card_bg and data_card_bg:
        print(f"  找到第一个stat-card (grid): top={stat_card_bg['top']}px, height={stat_card_bg['height']}px, bottom={stat_card_bg['bottom']}px")
        print(f"  找到第一个data-card: top={data_card_bg['top']}px, height={data_card_bg['height']}px")
        
        spacing = data_card_bg['top'] - stat_card_bg['bottom']
        print(f"  grid到data-card的间距: {spacing}px")
    
    if stat_card_bg and data_card_bg:
        # 间距应该约等于24px（mb-6）
        if 20 <= spacing <= 30:
            print(f"  ✅ 间距正确 (期望24px，实际{spacing}px)")
            return True
        else:
            print(f"  ❌ 间距不正确 (期望24px，实际{spacing}px)")
            return False
    else:
        print(f"  ❌ 未找到关键容器")
        return False

def verify_slide_005(pptx_path):
    """验证slide_005: stat-card背景高度100%填充"""
    print("\n=== 验证 slide_005 ===")
    prs = Presentation(pptx_path)
    slide = prs.slides[0]
    
    # 查找stat-card的背景矩形
    stat_card_backgrounds = []
    stat_card_texts = []
    
    for shape in slide.shapes:
        if hasattr(shape, "top") and hasattr(shape, "height"):
            # 使用EMU转换，更精确：1px = 9525 EMU
            top_px = int(shape.top / 9525)
            height_px = int(shape.height / 9525)
            
            # 矩形背景，必须是stat-card背景色
            if hasattr(shape, 'fill') and shape.fill.type == 1:  # Solid fill
                rgb = shape.fill.fore_color.rgb
                # stat-card背景色: rgba(10, 66, 117, 0.08) blended with white ≈ (235, 243, 248)
                if 230 < rgb[0] < 245 and 235 < rgb[1] < 250 and 240 < rgb[2] < 255:
                    if 200 < top_px < 250:  # 大概在标题后的位置
                        stat_card_backgrounds.append({
                            'top': top_px,
                            'height': height_px,
                            'shape': shape
                        })
            
            # 文本
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if "100%" in text or "防火墙防护率" in text or "EDR" in text:
                    stat_card_texts.append({
                        'text': text[:30],
                        'top': top_px,
                        'height': height_px
                    })
    
    print(f"  找到{len(stat_card_backgrounds)}个stat-card背景")
    print(f"  找到{len(stat_card_texts)}个stat-card文本")
    
    if len(stat_card_backgrounds) >= 2:
        # 检查两个stat-card的高度是否相同（等高）
        bg1 = stat_card_backgrounds[0]
        bg2 = stat_card_backgrounds[1]
        
        print(f"  stat-card #1: 高度={bg1['height']}px")
        print(f"  stat-card #2: 高度={bg2['height']}px")
        
        if abs(bg1['height'] - bg2['height']) < 5:
            print(f"  ✅ 两个stat-card高度一致 (100%填充)")
            return True
        else:
            print(f"  ❌ 两个stat-card高度不一致")
            return False
    else:
        print(f"  ❌ 未找到足够的stat-card背景")
        return False

def verify_slide_006(pptx_path):
    """验证slide_006: h3标题不重复"""
    print("\n=== 验证 slide_006 ===")
    prs = Presentation(pptx_path)
    slide = prs.slides[0]
    
    # 统计所有文本
    all_texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text = shape.text.strip()
            if text:
                all_texts.append(text)
    
    # 检查"弱口令数量"是否重复
    weak_password_count = sum(1 for text in all_texts if "弱口令数量" in text)
    
    print(f"  '弱口令数量'出现次数: {weak_password_count}")
    
    if weak_password_count == 1:
        print(f"  ✅ h3标题没有重复")
        return True
    else:
        print(f"  ❌ h3标题重复了 (出现{weak_password_count}次)")
        # 显示所有包含"弱口令"的文本
        for i, text in enumerate(all_texts):
            if "弱口令" in text:
                print(f"    #{i}: {text[:50]}")
        return False

def verify_slide_008(pptx_path):
    """验证slide_008: 第二个容器包含"41.0分钟" """
    print("\n=== 验证 slide_008 ===")
    prs = Presentation(pptx_path)
    slide = prs.slides[0]
    
    # 统计所有文本
    all_texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text = shape.text.strip()
            if text:
                all_texts.append(text)
    
    # 检查"41.0分钟"是否存在
    has_41_minutes = any("41.0" in text and "分钟" in text for text in all_texts)
    has_avg_response = any("平均响应时间" in text for text in all_texts)
    
    print(f"  找到'平均响应时间': {has_avg_response}")
    print(f"  找到'41.0分钟': {has_41_minutes}")
    
    if has_41_minutes:
        print(f"  ✅ 第二个data-card包含'41.0分钟'")
        return True
    else:
        print(f"  ❌ 第二个data-card缺少'41.0分钟'")
        # 显示所有包含"响应"或"分钟"的文本
        for text in all_texts:
            if "响应" in text or "分钟" in text:
                print(f"    - {text[:60]}")
        return False

# 主验证流程
print("=" * 60)
print("开始验证所有修复")
print("=" * 60)

results = {}

# 验证slide_004
if os.path.exists("output/slide_004.pptx"):
    results['slide_004'] = verify_slide_004("output/slide_004.pptx")
else:
    print("\n=== slide_004.pptx 不存在 ===")
    results['slide_004'] = False

# 验证slide_005
if os.path.exists("output/slide_005.pptx"):
    results['slide_005'] = verify_slide_005("output/slide_005.pptx")
else:
    print("\n=== slide_005.pptx 不存在 ===")
    results['slide_005'] = False

# 验证slide_006
if os.path.exists("output/slide_006.pptx"):
    results['slide_006'] = verify_slide_006("output/slide_006.pptx")
else:
    print("\n=== slide_006.pptx 不存在 ===")
    results['slide_006'] = False

# 验证slide_008
if os.path.exists("output/slide_008.pptx"):
    results['slide_008'] = verify_slide_008("output/slide_008.pptx")
else:
    print("\n=== slide_008.pptx 不存在 ===")
    results['slide_008'] = False

# 总结
print("\n" + "=" * 60)
print("验证结果总结")
print("=" * 60)
for slide, result in results.items():
    status = "✅ 通过" if result else "❌ 失败"
    print(f"{slide}: {status}")

all_passed = all(results.values())
if all_passed:
    print("\n🎉 所有问题都已修复！")
else:
    print("\n⚠️ 还有问题需要修复")
    failed = [k for k, v in results.items() if not v]
    print(f"失败的测试: {', '.join(failed)}")

