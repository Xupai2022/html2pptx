"""
HTML批量转PPTX工具（完整样式保留版）
使用修改后的HTML2PPTX类实现完美的样式保留
"""

import os
import sys
import glob
from pathlib import Path
from datetime import datetime
import shutil
import signal
import time
from concurrent.futures import ThreadPoolExecutor, TimeoutError as FutureTimeoutError

# 设置控制台编码为UTF-8（Windows）
if sys.platform == 'win32':
    import locale
    os.system('chcp 65001 >nul')
    sys.stdout.reconfigure(encoding='utf-8')
    sys.stderr.reconfigure(encoding='utf-8')

# 添加src目录到路径
sys.path.insert(0, str(Path(__file__).parent / 'src'))

from src.main import HTML2PPTX
from pptx import Presentation
import logging

# 设置日志
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)


class BatchConverter:
    """批量转换器 - 完整样式保留版"""

    def __init__(self, timeout_seconds: int = 30, svg_timeout_seconds: int = 5):
        """
        初始化批量转换器

        Args:
            timeout_seconds: 每个HTML文件的处理超时时间（秒）
            svg_timeout_seconds: 每个SVG的处理超时时间（秒）
        """
        self.output_dir = Path("output")
        self.timeout_seconds = timeout_seconds
        self.svg_timeout_seconds = svg_timeout_seconds
        self.processed_count = 0
        self.failed_count = 0
        self.start_time = None

    def find_slide_files(self):
        """
        查找input目录下所有以slide开头的HTML文件

        Returns:
            排序后的HTML文件路径列表
        """
        # 检查input目录是否存在
        input_dir = Path("input")
        if not input_dir.exists():
            logger.error("input目录不存在！")
            return []

        # 查找input目录下所有以slide开头的HTML文件
        pattern = "input/slide*.html"
        html_files = sorted(glob.glob(pattern))

        logger.info(f"找到 {len(html_files)} 个slide文件")
        return html_files

    def create_output_directory(self):
        """创建输出目录"""
        self.output_dir.mkdir(exist_ok=True)
        logger.info(f"输出目录: {self.output_dir}")

    def convert_all_to_single_pptx(self):
        """
        将所有HTML文件转换为单个PPTX文件
        使用共享的Presentation对象确保样式一致性
        带有超时控制
        """
        try:
            # 查找HTML文件
            html_files = self.find_slide_files()

            if not html_files:
                logger.warning("未找到任何slide开头的HTML文件")
                return

            # 创建输出目录
            self.create_output_directory()

            # 初始化计数器
            self.start_time = time.time()
            self.processed_count = 0
            self.failed_count = 0

            logger.info(f"\n开始批量转换 {len(html_files)} 个HTML文件...")
            logger.info(f"超时设置: 文件处理 {self.timeout_seconds}秒, SVG处理 {self.svg_timeout_seconds}秒")

            # 创建一个共享的Presentation对象
            # 这确保了所有幻灯片使用相同的样式和格式
            shared_presentation = Presentation()

            # 设置幻灯片尺寸
            from src.utils.unit_converter import UnitConverter
            shared_presentation.slide_width = UnitConverter.px_to_emu(1920)
            shared_presentation.slide_height = UnitConverter.px_to_emu(1080)

            # 移除默认的标题幻灯片（如果有）
            if len(shared_presentation.slides) > 0:
                xml_slides = shared_presentation.slides._sldIdLst
                xml_slides.remove(xml_slides[0])

            total_slides_processed = 0

            # 处理每个HTML文件
            for i, html_file in enumerate(html_files):
                logger.info(f"\n处理第 {i+1}/{len(html_files)} 个文件: {os.path.basename(html_file)}")
                logger.info(f"进度: 成功 {self.processed_count}, 失败 {self.failed_count}")

                # 使用线程池执行器来处理超时
                with ThreadPoolExecutor(max_workers=1) as executor:
                    try:
                        # 提交任务
                        future = executor.submit(
                            self._process_single_html,
                            html_file,
                            shared_presentation
                        )

                        # 等待任务完成，带超时
                        result = future.result(timeout=self.timeout_seconds)

                        if result:
                            total_slides_processed += 1
                            self.processed_count += 1
                            logger.info(f"  [✓] 成功处理 ({time.time() - self.start_time:.1f}s)")
                        else:
                            self.failed_count += 1
                            logger.error(f"  [✗] 处理失败")

                    except FutureTimeoutError:
                        self.failed_count += 1
                        logger.error(f"  [✗] 处理超时（>{self.timeout_seconds}秒），跳过")
                        # 取消任务
                        future.cancel()

                    except Exception as e:
                        self.failed_count += 1
                        logger.error(f"  [✗] 处理异常: {str(e)}")
                        import traceback
                        traceback.print_exc()

            # 保存最终文件
            elapsed_time = time.time() - self.start_time
            if total_slides_processed > 0:
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                output_file = self.output_dir / f"merged_presentation_{timestamp}.pptx"
                shared_presentation.save(output_file)

                # 清理PNG截图文件
                self._cleanup_svg_screenshots()

                logger.info(f"\n" + "="*60)
                logger.info(f"[SUCCESS] 批量转换完成！")
                logger.info(f"输出文件: {output_file}")
                logger.info(f"总文件数: {len(html_files)}")
                logger.info(f"成功处理: {self.processed_count}")
                logger.info(f"处理失败: {self.failed_count}")
                logger.info(f"总耗时: {elapsed_time:.1f}秒")
                logger.info(f"平均每文件: {elapsed_time/max(1, self.processed_count):.1f}秒")
                logger.info(f"最终PPTX幻灯片数: {len(shared_presentation.slides)}")
                logger.info("="*60)
            else:
                logger.error("\n没有成功处理任何HTML文件")

        except KeyboardInterrupt:
            logger.error("\n用户中断批量转换")
            raise
        except Exception as e:
            logger.error(f"\n批量转换失败: {str(e)}")
            import traceback
            traceback.print_exc()
            raise

    def _process_single_html(self, html_file, shared_presentation) -> bool:
        """
        处理单个HTML文件（在独立线程中运行）

        Args:
            html_file: HTML文件路径
            shared_presentation: 共享的Presentation对象

        Returns:
            是否成功
        """
        try:
            # 清理全局缓存，避免不同HTML文件间的样式污染
            from src.utils.style_computer import _style_computer_instance
            from src.utils.font_manager import _font_manager_instance
            from src.utils.font_size_extractor import _font_size_extractor_instance

            if _style_computer_instance is not None:
                _style_computer_instance.clear_cache()
            if _font_manager_instance is not None:
                _font_manager_instance._cached_fonts.clear()
            if _font_size_extractor_instance is not None:
                _font_size_extractor_instance.clear_cache()

            # 设置SVG超时
            if hasattr(self, 'svg_timeout_seconds'):
                # 设置全局SVG超时（通过修改chart_capture的默认值）
                import src.utils.chart_capture
                original_wait_time = getattr(src.utils.chart_capture.ChartCapture, 'default_wait_time', 1000)
                src.utils.chart_capture.ChartCapture.default_wait_time = self.svg_timeout_seconds * 1000

            # 创建转换器，传入共享的Presentation对象
            converter = HTML2PPTX(
                html_path=html_file,
                existing_presentation=shared_presentation
            )

            # 使用特殊的批量转换方法
            self._convert_html_to_shared_presentation(converter, html_file)

            return True

        except Exception as e:
            logger.error(f"处理HTML失败: {e}")
            return False

    def _convert_html_to_shared_presentation(self, converter, html_file):
        """
        将HTML内容转换到共享的Presentation对象
        这是核心方法，确保样式完全保留
        """
        # 获取HTML解析器
        html_parser = converter.html_parser

        # 获取所有幻灯片
        slides = html_parser.get_slides()

        for slide_html in slides:
            logger.info(f"  处理幻灯片...")

            # 使用共享的pptx_builder添加幻灯片
            # 这里会使用完全相同的转换逻辑
            pptx_slide = converter.pptx_builder.add_blank_slide()

            # 初始化转换器
            from src.converters.text_converter import TextConverter
            from src.converters.table_converter import TableConverter
            from src.converters.shape_converter import ShapeConverter
            from src.utils.font_manager import get_font_manager
            from src.utils.style_computer import get_style_computer

            text_converter = TextConverter(pptx_slide, converter.css_parser)
            table_converter = TableConverter(pptx_slide, converter.css_parser)
            shape_converter = ShapeConverter(pptx_slide, converter.css_parser)

            # 1. 添加顶部装饰条
            shape_converter.add_top_bar()

            # 2. 添加标题和副标题
            title_info = html_parser.get_title_info(slide_html)
            if title_info:
                # content-section的padding-top是20px
                title_end_y = text_converter.convert_title(
                    title_info['text'],
                    title_info['subtitle'],
                    x=80,
                    y=20,
                    is_cover=title_info.get('is_cover', False),
                    title_classes=title_info.get('classes', []),
                    h1_element=title_info.get('h1_element')
                )
                # space-y-10的第一个子元素紧接标题区域（无上间距）
                y_offset = title_end_y
            else:
                # 没有标题时使用默认位置（content-section padding-top）
                y_offset = 20

            # 3. 处理内容区域
            # 统一处理所有容器：优先查找space-y-10容器，如果没有则处理content-section的直接子元素
            space_y_container = slide_html.find('div', class_='space-y-10')
            content_section = slide_html.find('div', class_='content-section')

            if space_y_container:
                # 按顺序遍历直接子元素
                is_first_container = True
                for container in space_y_container.find_all(recursive=False):
                    if not container.name:
                        continue

                    container_classes = container.get('class', [])

                    # space-y-10: 第一个元素无上间距，后续元素有40px间距
                    if not is_first_container:
                        y_offset += 40  # space-y-10间距
                    is_first_container = False

                    # 根据class路由到对应的处理方法
                    y_offset = converter._process_container(container, pptx_slide, y_offset, shape_converter)
            else:
                # 处理content-section的直接子元素
                logger.info("  未找到space-y-10容器，处理content-section的直接子元素")

                # 跳过标题区域
                containers = []
                skip_first_mb = True
                for child in content_section.children:
                    if hasattr(child, 'get') and child.get('class'):
                        classes = child.get('class', [])
                        has_title = child.find('h1') or child.find('h2')

                        is_title_container = False
                        if skip_first_mb and any(cls in ['mb-6', 'mb-4', 'mb-8'] for cls in classes) and has_title:
                            has_content = any(cls in classes for cls in ['grid', 'stat-card', 'data-card', 'risk-card', 'flex'])
                            if not has_content:
                                is_title_container = True
                                skip_first_mb = False

                        if is_title_container:
                            continue

                        if child.name:
                            containers.append(child)

                # 处理所有容器
                for container in containers:
                    if container.name:
                        if containers.index(container) > 0:
                            y_offset += 40
                        y_offset = converter._process_container(container, pptx_slide, y_offset, shape_converter)

            # 4. 添加页码
            page_num = html_parser.get_page_number(slide_html)
            if page_num:
                shape_converter.add_page_number(page_num)

    def _cleanup_svg_screenshots(self):
        """清理PNG截图文件"""
        try:
            # 查找所有PNG截图文件
            png_pattern = "svg_screenshot_*.png"
            png_files = glob.glob(png_pattern)

            if png_files:
                logger.info(f"\n清理 {len(png_files)} 个PNG截图文件...")
                deleted_count = 0
                failed_count = 0

                for png_file in png_files:
                    try:
                        os.remove(png_file)
                        deleted_count += 1
                    except Exception as e:
                        logger.warning(f"  删除失败: {png_file} - {e}")
                        failed_count += 1

                logger.info(f"  成功删除: {deleted_count} 个文件")
                if failed_count > 0:
                    logger.warning(f"  删除失败: {failed_count} 个文件")
            else:
                logger.debug("  未找到PNG截图文件")

        except Exception as e:
            logger.error(f"清理PNG截图文件失败: {e}")


def main():
    """主函数"""
    import argparse

    # 解析命令行参数
    parser = argparse.ArgumentParser(
        description='HTML批量转PPTX工具（完整样式保留版）',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
示例:
  python batch_convert.py                    # 使用默认超时（30秒）
  python batch_convert.py -t 60              # 设置文件处理超时60秒
  python batch_convert.py -t 60 -s 10        # 文件超时60秒，SVG超时10秒
  python batch_convert.py --fast             # 快速模式（15秒/3秒）
        """
    )
    parser.add_argument(
        '-t', '--timeout',
        type=int,
        default=30,
        help='每个HTML文件的处理超时时间（秒，默认30）'
    )
    parser.add_argument(
        '-s', '--svg-timeout',
        type=int,
        default=5,
        help='每个SVG的处理超时时间（秒，默认5）'
    )
    parser.add_argument(
        '--fast',
        action='store_true',
        help='快速模式（等同于 -t 15 -s 3）'
    )

    args = parser.parse_args()

    # 快速模式
    if args.fast:
        args.timeout = 15
        args.svg_timeout = 3

    print("=" * 60)
    print("HTML批量转PPTX工具（完整样式保留版）")
    print("=" * 60)
    print(f"超时设置: 文件 {args.timeout}秒, SVG {args.svg_timeout}秒")
    print("=" * 60)

    # 检查input目录
    input_dir = Path("input")
    if not input_dir.exists():
        print("\n错误：input目录不存在！")
        print("请创建input目录并放入slide开头的HTML文件")
        return

    # 查找HTML文件
    html_files = sorted(glob.glob("input/slide*.html"))
    if not html_files:
        print("\n未找到任何slide开头的HTML文件！")
        print(f"请检查input目录中的文件")
        return

    print(f"\n找到 {len(html_files)} 个HTML文件:")
    for i, f in enumerate(html_files[:10]):  # 最多显示10个
        print(f"  - {os.path.basename(f)}")
    if len(html_files) > 10:
        print(f"  ... 还有 {len(html_files) - 10} 个文件")

    # 创建转换器并执行转换
    converter = BatchConverter(
        timeout_seconds=args.timeout,
        svg_timeout_seconds=args.svg_timeout
    )

    try:
        converter.convert_all_to_single_pptx()
    except KeyboardInterrupt:
        print("\n\n用户中断转换")
        sys.exit(1)
    except Exception as e:
        print(f"\n错误：{str(e)}")
        sys.exit(1)

    print("\n" + "=" * 60)
    print("转换结束")
    print("=" * 60)


if __name__ == "__main__":
    main()