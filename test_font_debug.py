"""调试字体设置"""
from src.parser.html_parser import HTMLParser
from src.parser.css_parser import CSSParser
from src.utils.font_manager import get_font_manager

html_path = 'test_custom_font.html'
html_parser = HTMLParser(html_path)
css_parser = CSSParser(html_parser.soup)

# 打印CSS规则
print("=== CSS规则 ===")
for selector, props in css_parser.style_rules.items():
    if 'font-family' in props:
        print(f"{selector}: {props['font-family']}")

# 测试字体管理器
font_manager = get_font_manager(css_parser)

print("\n=== 字体获取测试 ===")
print(f"body字体: {font_manager.get_font('body')}")
print(f"h1字体: {font_manager.get_font('h1')}")
print(f"h2字体: {font_manager.get_font('h2')}")
print(f"p字体: {font_manager.get_font('p')}")

# 检查实际设置
print("\n=== 实际转换测试 ===")
from src.renderer.pptx_builder import PPTXBuilder
from src.converters.text_converter import TextConverter
from pptx.util import Pt

pptx_builder = PPTXBuilder()
pptx_slide = pptx_builder.add_blank_slide()
text_converter = TextConverter(pptx_slide, css_parser)

# 转换标题
title_end_y = text_converter.convert_title("测试标题", "测试副标题", x=80, y=20)

# 检查文本框中的字体
print("\n=== PPTX中的实际字体 ===")
for shape in pptx_slide.shapes:
    if hasattr(shape, 'text_frame'):
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                print(f"文本: {run.text[:20]}... → 字体: {run.font.name}")

