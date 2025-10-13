"""检查生成的PPTX文件中的字体"""
from pptx import Presentation

pptx_path = 'output/test_custom_font.pptx'
prs = Presentation(pptx_path)

print("=== 检查PPTX中的所有字体 ===\n")

for slide_idx, slide in enumerate(prs.slides, 1):
    print(f"幻灯片 {slide_idx}:")
    for shape_idx, shape in enumerate(slide.shapes, 1):
        if hasattr(shape, 'text_frame') and shape.text:
            print(f"  形状 {shape_idx}: {shape.text[:30]}...")
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip():
                        font_name = run.font.name if run.font.name else "(未设置)"
                        print(f"    文本: '{run.text[:20]}' → 字体: {font_name}")
    print()

